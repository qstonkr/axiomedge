#!/usr/bin/env python3
"""Confluence Full Content Crawler

위키의 전체 본문, 첨부파일 내용까지 수집하는 확장 크롤러.

기능:
1. 전체 본문 저장 (200자 preview → 전체)
2. HTML 원본 저장 (테이블 구조 보존)
3. 첨부파일 다운로드 및 내용 추출
   - PDF: PyMuPDF로 텍스트/테이블 추출
   - Excel: openpyxl로 시트 데이터 추출
   - Word: python-docx로 텍스트/테이블 추출
   - PPT: python-pptx로 슬라이드 텍스트 추출
   - 이미지: 메타데이터 + Base64 (OCR 대기)
4. @멘션, 담당자 테이블 파싱

Usage:
    cd <project-root>
    uv run python scripts/confluence_full_crawler.py --page-id 318803690 --sample 10
    uv run python scripts/confluence_full_crawler.py --page-id 318803690 --full
    uv run python scripts/confluence_full_crawler.py --spaces "318803690,6685934" --full

Created: 2026-02-05
Updated: 2026-02-05 - 첨부파일 파싱 완전 구현
"""

from __future__ import annotations

import asyncio
import io
import json
import os
import re
import signal
import sys
import time
import warnings
from dataclasses import dataclass, field
from datetime import datetime, timezone
from html.parser import HTMLParser
from pathlib import Path
from typing import Any

import httpx

# Suppress warnings from document parsing libraries
warnings.filterwarnings("ignore", category=UserWarning)

sys.path.insert(0, str(Path(__file__).parent.parent))

from rich.console import Console
from rich.panel import Panel
from rich.progress import Progress, SpinnerColumn, TextColumn
from rich.table import Table

console = Console()

# =============================================================================
# Configuration (Environment Variables)
# =============================================================================

# .env.local 자동 로드
_env_local = Path(__file__).parent.parent / ".env.local"
if _env_local.exists():
    with open(_env_local) as _f:
        for _line in _f:
            _line = _line.strip()
            if _line and not _line.startswith("#") and "=" in _line:
                _key, _, _val = _line.partition("=")
                _key, _val = _key.strip(), _val.strip()
                if _key and _key not in os.environ:
                    os.environ[_key] = _val

CONFLUENCE_BASE_URL = os.getenv("CONFLUENCE_BASE_URL", "https://wiki.gsretail.com")
CONFLUENCE_PAT = os.getenv("CONFLUENCE_PAT")

if not CONFLUENCE_PAT:
    console.print("[bold red]Error:[/] CONFLUENCE_PAT 환경변수가 설정되지 않았습니다.")
    console.print("[dim]export CONFLUENCE_PAT='your-personal-access-token'[/]")
    sys.exit(1)


def _load_knowledge_sources() -> dict:
    """환경변수에서 지식 소스 로드 (KNOWLEDGE_SOURCES_JSON 또는 기본값 사용)"""
    sources_json = os.getenv("KNOWLEDGE_SOURCES_JSON")
    if sources_json:
        try:
            return json.loads(sources_json)
        except json.JSONDecodeError as e:
            console.print(f"[yellow]Warning:[/] KNOWLEDGE_SOURCES_JSON 파싱 실패: {e}")
            console.print("[dim]기본 지식 소스를 사용합니다.[/]")

    # 기본 지식 소스 (환경변수 미설정 시)
    # space_key는 API 응답에서 자동으로 가져옴
    sources: dict[str, dict[str, str]] = {}
    infra_pid = os.getenv("KNOWLEDGE_SOURCE_INFRA_PAGE_ID", "312765276")
    if infra_pid:
        sources["infra"] = {
            "page_id": infra_pid,
            "name": os.getenv("KNOWLEDGE_SOURCE_INFRA_NAME", "infra"),
        }
    hs_pid = os.getenv("KNOWLEDGE_SOURCE_HS_PAGE_ID")
    if hs_pid:
        sources["homeshopping"] = {
            "page_id": hs_pid,
            "name": os.getenv("KNOWLEDGE_SOURCE_HS_NAME", "homeshopping"),
        }
    # AX본부 KB sources (small Confluence spaces)
    faq_pid = os.getenv("KNOWLEDGE_SOURCE_FAQ_PAGE_ID", "388722996")
    if faq_pid:
        sources["faq"] = {
            "page_id": faq_pid,
            "name": os.getenv("KNOWLEDGE_SOURCE_FAQ_NAME", "AX FAQ"),
        }
    system_pid = os.getenv("KNOWLEDGE_SOURCE_SYSTEM_PAGE_ID", "388722906")
    if system_pid:
        sources["system"] = {
            "page_id": system_pid,
            "name": os.getenv("KNOWLEDGE_SOURCE_SYSTEM_NAME", "AX 시스템별상세정보"),
        }
    dictionary_pid = os.getenv("KNOWLEDGE_SOURCE_DICTIONARY_PAGE_ID", "388722934")
    if dictionary_pid:
        sources["dictionary"] = {
            "page_id": dictionary_pid,
            "name": os.getenv("KNOWLEDGE_SOURCE_DICTIONARY_NAME", "AX 용어사전"),
        }
    # AX챗봇 지식관리
    axchat_pid = os.getenv("KNOWLEDGE_SOURCE_AXCHAT_PAGE_ID", "373865276")
    if axchat_pid:
        sources["axchat_know"] = {
            "page_id": axchat_pid,
            "name": os.getenv("KNOWLEDGE_SOURCE_AXCHAT_NAME", "AX챗봇 지식관리"),
        }
    # Large Confluence sources
    hax_pid = os.getenv("KNOWLEDGE_SOURCE_HAX_PAGE_ID", "6685934")
    if hax_pid:
        sources["hax"] = {
            "page_id": hax_pid,
            "name": os.getenv("KNOWLEDGE_SOURCE_HAX_NAME", "홈쇼핑AX부문 업무 위키"),
        }
    itops_pid = os.getenv("KNOWLEDGE_SOURCE_ITOPS_PAGE_ID", "14225186")
    if itops_pid:
        sources["itops"] = {
            "page_id": itops_pid,
            "name": os.getenv("KNOWLEDGE_SOURCE_ITOPS_NAME", "홈쇼핑 IT운영 업무 가이드"),
        }
    if not sources:
        console.print("[bold red]Error:[/] 지식 소스가 설정되지 않았습니다.")
        console.print("[dim].env.local에 KNOWLEDGE_SOURCE_INFRA_PAGE_ID 등을 설정하세요.[/]")
        sys.exit(1)
    return sources


KNOWLEDGE_SOURCES = _load_knowledge_sources()


def _resolve_output_dir() -> Path:
    """Resolve crawl output directory with safe fallback for local execution."""
    candidates: list[tuple[Path, str]] = []
    configured_dir = os.getenv("CONFLUENCE_OUTPUT_DIR", "").strip()
    if configured_dir:
        candidates.append((Path(configured_dir), "CONFLUENCE_OUTPUT_DIR"))
    fallback_dir = Path.home() / ".axiomedge" / "crawl"
    candidates.append((fallback_dir, "home fallback"))

    last_error: Exception | None = None
    for output_dir, label in candidates:
        try:
            output_dir.mkdir(parents=True, exist_ok=True)
            return output_dir
        except OSError as exc:
            last_error = exc
            if label == "CONFLUENCE_OUTPUT_DIR":
                console.print(
                    f"[yellow]⚠️  CONFLUENCE_OUTPUT_DIR '{output_dir}' is not writable; "
                    f"fallback to {fallback_dir}[/yellow]"
                )
            else:
                console.print(
                    f"[yellow]⚠️  Failed to create fallback output dir '{output_dir}': {exc}[/yellow]"
                )

    if last_error is not None:
        raise RuntimeError(
            "Unable to resolve writable CONFLUENCE_OUTPUT_DIR"
        ) from last_error
    return fallback_dir


# 저장 경로 (환경변수로 설정 가능)
OUTPUT_DIR = _resolve_output_dir()
ATTACHMENTS_DIR = OUTPUT_DIR / "attachments"
ATTACHMENTS_DIR.mkdir(parents=True, exist_ok=True)


def _env_bool(name: str, default: bool) -> bool:
    raw = os.getenv(name)
    if raw is None:
        return default
    return raw.strip().lower() in {"1", "true", "yes", "on"}


def _env_int(name: str) -> int | None:
    raw = os.getenv(name)
    if raw is None or raw.strip() == "":
        return None
    try:
        return int(raw)
    except ValueError:
        return None


_DEFAULT_ATTACHMENT_OCR_MODE = "force"
_DEFAULT_OCR_MIN_TEXT_CHARS = 100
_DEFAULT_OCR_MAX_PDF_PAGES = 1_000_000
_DEFAULT_OCR_MAX_PPT_SLIDES = 1_000_000
_DEFAULT_OCR_MAX_IMAGES_PER_ATTACHMENT = 1

_SOURCE_ATTACHMENT_OCR_DEFAULTS: dict[str, dict[str, Any]] = {
    "itops": {
        "attachment_ocr_mode": "auto",
        "ocr_min_text_chars": 100,
        "ocr_max_pdf_pages": 10,
        "ocr_max_ppt_slides": 10,
        "ocr_max_images_per_attachment": 1,
        "slide_render_enabled": False,
        "layout_analysis_enabled": False,
    }
}


@dataclass(frozen=True)
class AttachmentOCRPolicy:
    attachment_ocr_mode: str
    ocr_min_text_chars: int
    ocr_max_pdf_pages: int
    ocr_max_ppt_slides: int
    ocr_max_images_per_attachment: int
    slide_render_enabled: bool
    layout_analysis_enabled: bool


@dataclass
class AttachmentParseResult:
    extracted_text: str
    extracted_tables: list[dict]
    confidence: float
    ocr_mode: str | None = None
    ocr_applied: bool = False
    ocr_skip_reason: str | None = None
    ocr_units_attempted: int = 0
    ocr_units_extracted: int = 0
    ocr_units_deferred: int = 0
    native_text_chars: int = 0
    ocr_text_chars: int = 0


# =============================================================================
# OCR Subprocess Isolation (SIGSEGV Defense)
# =============================================================================
# PaddleOCR C++ inference can SIGSEGV on certain images (signal 11).
# Python try/except CANNOT catch OS signals — the entire process dies.
# Solution: run OCR in a forked subprocess via ProcessPoolExecutor.
# If SIGSEGV occurs, only the child process dies; the parent skips the image.
# Fork inherits the loaded PaddleOCR model via copy-on-write (no reload cost).


def _ocr_worker_fn(image_bytes: bytes) -> tuple:
    """OCR worker function executed in forked subprocess.

    Inherits PaddleOCR model from parent via fork COW — no model reload.
    If SIGSEGV occurs here, only this subprocess dies.
    """
    try:
        if not hasattr(_ocr_worker_fn, "_ocr"):
            try:
                from src.nlp.ocr.paddle_ocr_provider import PaddleOCRProvider

                _ocr_worker_fn._ocr = PaddleOCRProvider()
            except ImportError:
                return (None, 0.0, [])
        result = _ocr_worker_fn._ocr.extract(image_bytes)
        tables = [
            {"headers": t.headers, "rows": t.rows, "source": "ocr"}
            for t in (result.tables or [])
        ]
        return (result.text, result.confidence, tables)
    except Exception:
        return (None, 0.0, [])


# =============================================================================
# Attachment Parsers
# =============================================================================
class AttachmentParser:
    """첨부파일 내용 추출기"""

    _active_source_key = ""
    _ocr_policy = AttachmentOCRPolicy(
        attachment_ocr_mode=_DEFAULT_ATTACHMENT_OCR_MODE,
        ocr_min_text_chars=_DEFAULT_OCR_MIN_TEXT_CHARS,
        ocr_max_pdf_pages=_DEFAULT_OCR_MAX_PDF_PAGES,
        ocr_max_ppt_slides=_DEFAULT_OCR_MAX_PPT_SLIDES,
        ocr_max_images_per_attachment=_DEFAULT_OCR_MAX_IMAGES_PER_ATTACHMENT,
        slide_render_enabled=_env_bool("KNOWLEDGE_SLIDE_RENDER_ENABLED", True),
        layout_analysis_enabled=_env_bool("KNOWLEDGE_LAYOUT_ANALYSIS_ENABLED", True),
    )

    @classmethod
    def configure_run(cls, source_key: str, overrides: dict[str, Any] | None = None) -> AttachmentOCRPolicy:
        """Resolve run-local OCR policy once per crawl source."""
        overrides = overrides or {}
        source_defaults = _SOURCE_ATTACHMENT_OCR_DEFAULTS.get(source_key, {})
        legacy_slide_render = _env_bool("KNOWLEDGE_SLIDE_RENDER_ENABLED", True)
        legacy_layout_analysis = _env_bool("KNOWLEDGE_LAYOUT_ANALYSIS_ENABLED", True)

        def _resolve_mode() -> str:
            raw_mode = (
                overrides.get("attachment_ocr_mode")
                or os.getenv("KNOWLEDGE_CRAWL_ATTACHMENT_OCR_MODE")
                or source_defaults.get("attachment_ocr_mode")
                or _DEFAULT_ATTACHMENT_OCR_MODE
            )
            value = str(raw_mode).strip().lower()
            return value if value in {"auto", "off", "force"} else _DEFAULT_ATTACHMENT_OCR_MODE

        def _resolve_int_value(
            override_key: str,
            env_key: str,
            source_key_name: str,
            legacy_default: int,
        ) -> int:
            if override_key in overrides and overrides[override_key] is not None:
                return max(0, int(overrides[override_key]))
            env_value = _env_int(env_key)
            if env_value is not None:
                return max(0, env_value)
            if source_key_name in source_defaults:
                return max(0, int(source_defaults[source_key_name]))
            return legacy_default

        def _resolve_bool_value(
            override_key: str,
            env_key: str,
            source_key_name: str,
            legacy_default: bool,
        ) -> bool:
            if override_key in overrides and overrides[override_key] is not None:
                return bool(overrides[override_key])
            if os.getenv(env_key) is not None:
                return _env_bool(env_key, legacy_default)
            if source_key_name in source_defaults:
                return bool(source_defaults[source_key_name])
            return legacy_default

        cls._active_source_key = source_key
        cls._ocr_policy = AttachmentOCRPolicy(
            attachment_ocr_mode=_resolve_mode(),
            ocr_min_text_chars=_resolve_int_value(
                "ocr_min_text_chars",
                "KNOWLEDGE_CRAWL_OCR_MIN_TEXT_CHARS",
                "ocr_min_text_chars",
                _DEFAULT_OCR_MIN_TEXT_CHARS,
            ),
            ocr_max_pdf_pages=_resolve_int_value(
                "ocr_max_pdf_pages",
                "KNOWLEDGE_CRAWL_OCR_MAX_PDF_PAGES",
                "ocr_max_pdf_pages",
                _DEFAULT_OCR_MAX_PDF_PAGES,
            ),
            ocr_max_ppt_slides=_resolve_int_value(
                "ocr_max_ppt_slides",
                "KNOWLEDGE_CRAWL_OCR_MAX_PPT_SLIDES",
                "ocr_max_ppt_slides",
                _DEFAULT_OCR_MAX_PPT_SLIDES,
            ),
            ocr_max_images_per_attachment=_resolve_int_value(
                "ocr_max_images_per_attachment",
                "KNOWLEDGE_CRAWL_OCR_MAX_IMAGES_PER_ATTACHMENT",
                "ocr_max_images_per_attachment",
                _DEFAULT_OCR_MAX_IMAGES_PER_ATTACHMENT,
            ),
            slide_render_enabled=_resolve_bool_value(
                "slide_render_enabled",
                "KNOWLEDGE_CRAWL_SLIDE_RENDER_ENABLED",
                "slide_render_enabled",
                legacy_slide_render,
            ),
            layout_analysis_enabled=_resolve_bool_value(
                "layout_analysis_enabled",
                "KNOWLEDGE_CRAWL_LAYOUT_ANALYSIS_ENABLED",
                "layout_analysis_enabled",
                legacy_layout_analysis,
            ),
        )
        return cls._ocr_policy

    @classmethod
    def current_policy(cls) -> AttachmentOCRPolicy:
        return cls._ocr_policy

    @staticmethod
    def _emit_status(heartbeat_fn, message: str) -> None:
        if heartbeat_fn:
            heartbeat_fn(message)
        else:
            print(f"[status] {message}")

    @staticmethod
    def _text_chars(value: str | None) -> int:
        return len(value.strip()) if value and value.strip() else 0

    @classmethod
    def parse_pdf(cls, file_path: Path, heartbeat_fn=None) -> AttachmentParseResult:
        """PDF에서 텍스트와 테이블 추출

        Strategy:
        1. PyMuPDF 텍스트 레이어 추출 시도
        2. 텍스트가 빈 페이지 → 이미지 렌더링 → PaddleOCR fallback
           (이미지 기반 PDF — PPT를 PDF로 내보낸 파일 등)

        Returns:
            (extracted_text, tables, confidence)
        """
        try:
            policy = cls.current_policy()
            import fitz  # PyMuPDF

            doc = fitz.open(file_path)
            text_parts = []
            tables = []
            native_text_chars = 0
            ocr_text_chars = 0
            ocr_units_attempted = 0
            ocr_units_extracted = 0
            ocr_units_deferred = 0
            total_pages = len(doc)
            textless_pages = 0

            cls._emit_status(heartbeat_fn, f"native_extract pdf pages={total_pages}")

            for page_num, page in enumerate(doc, 1):
                # 1차: 텍스트 레이어 추출
                page_text = page.get_text("text")
                if page_text.strip():
                    text_parts.append(f"[Page {page_num}]\n{page_text}")
                    native_text_chars += cls._text_chars(page_text)
                else:
                    textless_pages += 1
                    if policy.attachment_ocr_mode == "off":
                        continue
                    if ocr_units_attempted >= policy.ocr_max_pdf_pages:
                        ocr_units_deferred += 1
                        continue
                    # 2차: 이미지 렌더링 → PaddleOCR fallback
                    try:
                        ocr_units_attempted += 1
                        cls._emit_status(
                            heartbeat_fn,
                            f"ocr_processing pdf page={page_num}/{total_pages}",
                        )
                        zoom = 2.0  # 144 DPI (72 * 2)
                        mat = fitz.Matrix(zoom, zoom)
                        pix = page.get_pixmap(matrix=mat)
                        png_bytes = pix.tobytes("png")
                        del pix  # Free native memory immediately

                        try:
                            from src.core.feature_flags import FeatureFlags
                            _pdf_postprocess = FeatureFlags.is_knowledge_ocr_postprocess_enabled()
                        except ImportError:
                            _pdf_postprocess = os.getenv("KNOWLEDGE_OCR_POSTPROCESS_ENABLED", "true").lower() == "true"

                        ocr_text = cls._ocr_slide_image(
                            png_bytes, page_num,
                            preprocess=True,
                            layout_analysis=policy.layout_analysis_enabled,
                            postprocess=_pdf_postprocess,
                        )
                        if ocr_text and ocr_text.strip():
                            text_parts.append(f"[Page {page_num}]\n{ocr_text}")
                            ocr_units_extracted += 1
                            ocr_text_chars += cls._text_chars(ocr_text)
                    except Exception as ocr_err:
                        print(f"[PDF OCR] Page {page_num} OCR failed: {ocr_err}")

                if heartbeat_fn and page_num % 5 == 0:
                    heartbeat_fn(f"pdf_ocr: {page_num}/{total_pages}")

                # 테이블 추출 시도
                try:
                    page_tables = page.find_tables()
                    for table in page_tables:
                        table_data = table.extract()
                        if table_data and len(table_data) > 1:
                            headers = table_data[0] if table_data else []
                            rows = table_data[1:] if len(table_data) > 1 else []
                            tables.append({
                                "page": page_num,
                                "headers": headers,
                                "rows": [dict(zip(headers, row)) for row in rows if len(row) == len(headers)],
                            })
                except Exception:
                    pass  # 테이블 추출 실패 시 건너뛰기

            doc.close()

            full_text = "\n\n".join(text_parts)
            if ocr_units_extracted > 0:
                print(f"[PDF OCR] {ocr_units_extracted}/{total_pages} pages used OCR fallback")
            if ocr_units_deferred > 0:
                cls._emit_status(
                    heartbeat_fn,
                    f"ocr_skipped_budget pdf deferred={ocr_units_deferred}",
                )
            confidence = 0.9 if full_text.strip() and ocr_units_extracted == 0 else (0.7 if full_text.strip() else 0.0)

            skip_reason = None
            if textless_pages > 0 and policy.attachment_ocr_mode == "off":
                skip_reason = "disabled"
            elif ocr_units_deferred > 0:
                skip_reason = "budget_exceeded"

            return AttachmentParseResult(
                extracted_text=full_text,
                extracted_tables=tables,
                confidence=confidence,
                ocr_mode=policy.attachment_ocr_mode,
                ocr_applied=ocr_units_extracted > 0,
                ocr_skip_reason=skip_reason,
                ocr_units_attempted=ocr_units_attempted,
                ocr_units_extracted=ocr_units_extracted,
                ocr_units_deferred=ocr_units_deferred,
                native_text_chars=native_text_chars,
                ocr_text_chars=ocr_text_chars,
            )

        except Exception as e:
            return AttachmentParseResult(
                extracted_text=f"[PDF 파싱 오류: {e}]",
                extracted_tables=[],
                confidence=0.0,
                ocr_mode=cls.current_policy().attachment_ocr_mode,
                ocr_skip_reason="parse_error",
            )

    @staticmethod
    def parse_excel(file_path: Path) -> AttachmentParseResult:
        """Excel에서 시트 데이터 추출"""
        try:
            from openpyxl import load_workbook

            wb = load_workbook(file_path, read_only=True, data_only=True)
            text_parts = []
            tables = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]

                # 시트 데이터를 테이블로 변환
                rows_data = []
                for row in sheet.iter_rows(values_only=True):
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    if any(v.strip() for v in row_values):  # 빈 행 제외
                        rows_data.append(row_values)

                if rows_data:
                    headers = rows_data[0] if rows_data else []
                    data_rows = rows_data[1:] if len(rows_data) > 1 else []

                    tables.append({
                        "sheet": sheet_name,
                        "headers": headers,
                        "rows": [dict(zip(headers, row)) for row in data_rows if len(row) == len(headers)],
                        "row_count": len(data_rows),
                    })

                    # 텍스트 표현
                    text_parts.append(f"[Sheet: {sheet_name}]")
                    text_parts.append(" | ".join(headers))
                    for row in data_rows[:10]:  # 최대 10행만 텍스트로
                        text_parts.append(" | ".join(row))
                    if len(data_rows) > 10:
                        text_parts.append(f"... 외 {len(data_rows) - 10}행")

            wb.close()

            full_text = "\n".join(text_parts)
            confidence = 0.95 if tables else 0.0

            return AttachmentParseResult(
                extracted_text=full_text,
                extracted_tables=tables,
                confidence=confidence,
                native_text_chars=AttachmentParser._text_chars(full_text),
            )

        except Exception as e:
            return AttachmentParseResult(
                extracted_text=f"[Excel 파싱 오류: {e}]",
                extracted_tables=[],
                confidence=0.0,
                ocr_skip_reason="parse_error",
            )

    @staticmethod
    def _extract_doc_olefile(file_path: Path) -> str | None:
        """Pure Python .doc 텍스트 추출 (olefile OLE2 스트림 파싱)"""
        try:
            import olefile
        except ImportError:
            return None

        try:
            ole = olefile.OleFileIO(str(file_path))
        except Exception:
            return None

        try:
            if not ole.exists("WordDocument"):
                return None

            raw = ole.openstream("WordDocument").read()
            if len(raw) < 20:
                return None

            # Word .doc: 텍스트는 FIB 이후 영역에 저장됨
            # 인코딩 우선순위: cp949(한국어) → cp1252(서양어) → utf-16-le
            for encoding in ("cp949", "cp1252"):
                try:
                    decoded = raw.decode(encoding, errors="ignore")
                    # 제어 문자 제거 (탭/줄바꿈 제외)
                    cleaned = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", decoded)
                    # 유효 텍스트 비율 검사 (한글/영문/숫자/공백)
                    readable = re.findall(r"[\uAC00-\uD7AF\u3000-\u9FFFa-zA-Z0-9\s,.!?()]+", cleaned)
                    text = " ".join(readable).strip()
                    text = re.sub(r"\s{3,}", "\n\n", text)
                    if len(text) > 50:
                        return text
                except Exception:
                    continue

            return None
        finally:
            ole.close()

    @staticmethod
    def _parse_legacy_doc(file_path: Path) -> AttachmentParseResult:
        """레거시 .doc (OLE2) 파일에서 텍스트 추출 (antiword → catdoc → olefile fallback)"""
        import shutil
        import subprocess

        # 1차: antiword (테이블 구조 보존, Docker 환경)
        antiword_path = shutil.which("antiword")
        if antiword_path:
            try:
                result = subprocess.run(
                    [antiword_path, str(file_path)],
                    capture_output=True, text=True, timeout=30,
                )
                if result.returncode == 0 and result.stdout.strip():
                    text = result.stdout.strip()
                    return AttachmentParseResult(
                        extracted_text=text,
                        extracted_tables=[],
                        confidence=0.7,
                        native_text_chars=AttachmentParser._text_chars(text),
                    )
            except Exception:
                pass

        # 2차: catdoc (antiword 없을 때, Docker 환경)
        catdoc_path = shutil.which("catdoc")
        if catdoc_path:
            try:
                result = subprocess.run(
                    [catdoc_path, "-w", str(file_path)],
                    capture_output=True, text=True, timeout=30,
                )
                if result.returncode == 0 and result.stdout.strip():
                    text = result.stdout.strip()
                    return AttachmentParseResult(
                        extracted_text=text,
                        extracted_tables=[],
                        confidence=0.6,
                        native_text_chars=AttachmentParser._text_chars(text),
                    )
            except Exception:
                pass

        # 3차: olefile pure Python 추출 (로컬 환경, 시스템 도구 없을 때)
        text = AttachmentParser._extract_doc_olefile(file_path)
        if text:
            return AttachmentParseResult(
                extracted_text=text,
                extracted_tables=[],
                confidence=0.5,
                native_text_chars=AttachmentParser._text_chars(text),
            )

        return AttachmentParseResult(
            extracted_text="[.doc 파싱 실패: antiword/catdoc 미설치 및 olefile 추출 실패]",
            extracted_tables=[],
            confidence=0.0,
            ocr_skip_reason="parse_error",
        )

    @staticmethod
    def parse_word(file_path: Path) -> AttachmentParseResult:
        """Word에서 텍스트와 테이블 추출 (.docx: python-docx, .doc: antiword/catdoc)"""
        try:
            if str(file_path).lower().endswith(".doc"):
                return AttachmentParser._parse_legacy_doc(file_path)
            from docx import Document

            doc = Document(file_path)
            text_parts = []
            tables = []

            # 문단 추출
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # 테이블 추출
            for idx, table in enumerate(doc.tables, 1):
                rows_data = []
                for row in table.rows:
                    row_values = [cell.text.strip() for cell in row.cells]
                    rows_data.append(row_values)

                if rows_data:
                    headers = rows_data[0] if rows_data else []
                    data_rows = rows_data[1:] if len(rows_data) > 1 else []

                    tables.append({
                        "table_index": idx,
                        "headers": headers,
                        "rows": [dict(zip(headers, row)) for row in data_rows if len(row) == len(headers)],
                    })

            full_text = "\n\n".join(text_parts)
            confidence = 0.9 if full_text.strip() else 0.0

            return AttachmentParseResult(
                extracted_text=full_text,
                extracted_tables=tables,
                confidence=confidence,
                native_text_chars=AttachmentParser._text_chars(full_text),
            )

        except Exception as e:
            return AttachmentParseResult(
                extracted_text=f"[Word 파싱 오류: {e}]",
                extracted_tables=[],
                confidence=0.0,
                ocr_skip_reason="parse_error",
            )

    @staticmethod
    def _extract_ppt_olefile(file_path: Path) -> str | None:
        """Pure Python .ppt 텍스트 추출 (OLE2 PowerPoint 레코드 파싱)"""
        try:
            import olefile
        except ImportError:
            return None

        import struct

        try:
            ole = olefile.OleFileIO(str(file_path))
        except Exception:
            return None

        try:
            if not ole.exists("PowerPoint Document"):
                return None

            raw = ole.openstream("PowerPoint Document").read()
            text_parts: list[str] = []
            offset = 0

            while offset < len(raw) - 8:
                try:
                    rec_type = struct.unpack_from("<H", raw, offset + 2)[0]
                    rec_len = struct.unpack_from("<I", raw, offset + 4)[0]
                except struct.error:
                    break

                data_start = offset + 8
                data_end = data_start + rec_len

                if rec_len > 0 and data_end <= len(raw):
                    # TextCharsAtom (0x0FA0): UTF-16LE 텍스트
                    if rec_type == 0x0FA0:
                        text = raw[data_start:data_end].decode("utf-16-le", errors="ignore").strip()
                        if text:
                            text_parts.append(text)
                    # TextBytesAtom (0x0FA8): ANSI 텍스트
                    elif rec_type == 0x0FA8:
                        data = raw[data_start:data_end]
                        try:
                            text = data.decode("cp949").strip()
                        except UnicodeDecodeError:
                            text = data.decode("cp1252", errors="ignore").strip()
                        if text:
                            text_parts.append(text)

                offset = data_end if rec_len > 0 else offset + 1

            return "\n\n".join(text_parts) if text_parts else None
        finally:
            ole.close()

    @staticmethod
    def _parse_legacy_ppt(file_path: Path, heartbeat_fn=None) -> AttachmentParseResult:
        """레거시 .ppt (OLE2) 파일에서 텍스트 추출

        Strategy:
        1. LibreOffice headless로 .ppt → .pptx 변환 → parse_ppt 재사용 (OCR 포함)
        2. Fallback: catppt → olefile 텍스트 추출
        """
        import shutil
        import subprocess
        import tempfile

        # 1차: LibreOffice .ppt → .pptx 변환 → parse_ppt (OCR 포함)
        from scripts.slide_renderer import _find_soffice

        soffice = _find_soffice()
        if soffice:
            try:
                with tempfile.TemporaryDirectory(prefix="ppt_convert_") as tmpdir:
                    lo_profile = os.path.join(tmpdir, "lo_profile")
                    os.makedirs(lo_profile, exist_ok=True)

                    result = subprocess.run(
                        [
                            soffice,
                            "--headless",
                            "--norestore",
                            f"-env:UserInstallation=file://{lo_profile}",
                            "--convert-to", "pptx",
                            "--outdir", tmpdir,
                            str(file_path),
                        ],
                        capture_output=True,
                        text=True,
                        timeout=120,
                    )
                    if result.returncode == 0:
                        pptx_files = list(Path(tmpdir).glob("*.pptx"))
                        if pptx_files:
                            print(f"[PPT] Converted {file_path.name} → {pptx_files[0].name}")
                            return AttachmentParser.parse_ppt(
                                pptx_files[0], heartbeat_fn=heartbeat_fn,
                            )
                    print(f"[PPT] LibreOffice conversion failed: {result.stderr[:200]}")
            except subprocess.TimeoutExpired:
                print(f"[PPT] LibreOffice conversion timeout for {file_path.name}")
            except Exception as e:
                print(f"[PPT] LibreOffice conversion error: {e}")

        # 2차: catppt (Docker 환경)
        catppt_path = shutil.which("catppt")
        if catppt_path:
            try:
                result = subprocess.run(
                    [catppt_path, str(file_path)],
                    capture_output=True, text=True, timeout=30,
                )
                if result.returncode == 0 and result.stdout.strip():
                    text = result.stdout.strip()
                    return AttachmentParseResult(
                        extracted_text=text,
                        extracted_tables=[],
                        confidence=0.6,
                        native_text_chars=AttachmentParser._text_chars(text),
                    )
            except Exception:
                pass

        # 3차: olefile pure Python 추출 (로컬 환경)
        text = AttachmentParser._extract_ppt_olefile(file_path)
        if text:
            return AttachmentParseResult(
                extracted_text=text,
                extracted_tables=[],
                confidence=0.5,
                native_text_chars=AttachmentParser._text_chars(text),
            )

        return AttachmentParseResult(
            extracted_text="[.ppt 파싱 실패: LibreOffice/catppt/olefile 모두 실패]",
            extracted_tables=[],
            confidence=0.0,
            ocr_mode=AttachmentParser.current_policy().attachment_ocr_mode,
            ocr_skip_reason="parse_error",
        )

    @staticmethod
    def _iter_shapes(shapes):
        """GroupShape 포함 재귀 shape 탐색."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        for shape in shapes:
            yield shape
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from AttachmentParser._iter_shapes(shape.shapes)

    @classmethod
    def parse_ppt(cls, file_path: Path, heartbeat_fn=None) -> AttachmentParseResult:
        """PPT에서 슬라이드 텍스트 추출 (.pptx: python-pptx, .ppt: catppt)

        Args:
            file_path: PPT/PPTX 파일 경로
            heartbeat_fn: Temporal heartbeat 콜백 (thread-safe, OCR 중 호출)
        """
        try:
            if str(file_path).lower().endswith(".ppt"):
                return cls._parse_legacy_ppt(file_path, heartbeat_fn=heartbeat_fn)
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            policy = cls.current_policy()
            prs = Presentation(file_path)
            text_parts = []
            tables = []
            image_shapes = []  # (slide_num, shape) tuples for deferred processing

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []

                for shape in AttachmentParser._iter_shapes(slide.shapes):
                    # 텍스트 프레임
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)

                    # 테이블
                    if shape.has_table:
                        table = shape.table
                        rows_data = []
                        for row in table.rows:
                            row_values = [cell.text.strip() for cell in row.cells]
                            rows_data.append(row_values)

                        if rows_data:
                            headers = rows_data[0] if rows_data else []
                            data_rows = rows_data[1:] if len(rows_data) > 1 else []
                            tables.append({
                                "slide": slide_num,
                                "headers": headers,
                                "rows": [dict(zip(headers, row)) for row in data_rows if len(row) == len(headers)],
                            })

                    # 이미지 shape 수집 (OCR 처리 대상)
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image_bytes = shape.image.blob
                            if len(image_bytes) > 10_000:  # 10KB 이상만
                                image_shapes.append((slide_num, image_bytes))
                        except Exception:
                            pass

                # 슬라이드 노트 추출
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_texts.append(f"[Notes] {notes_text}")

                if slide_texts:
                    text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))

            native_text = "\n\n".join(text_parts)
            native_text_chars = cls._text_chars(native_text)
            cls._emit_status(
                heartbeat_fn,
                f"native_extract ppt slides={len(prs.slides)} chars={native_text_chars}",
            )

            should_ocr = policy.attachment_ocr_mode == "force" or (
                policy.attachment_ocr_mode == "auto"
                and native_text_chars < policy.ocr_min_text_chars
            )
            if policy.attachment_ocr_mode == "off":
                should_ocr = False

            # ============================================================
            # Phase 1: Slide rendering (LibreOffice PPTX → PNG)
            # ============================================================
            slide_rendered = False
            try:
                from src.core.feature_flags import FeatureFlags
                ocr_preprocess_enabled = FeatureFlags.is_knowledge_ocr_preprocess_enabled()
                ocr_postprocess_enabled = FeatureFlags.is_knowledge_ocr_postprocess_enabled()
            except ImportError:
                ocr_preprocess_enabled = os.getenv("KNOWLEDGE_OCR_PREPROCESS_ENABLED", "true").lower() == "true"
                ocr_postprocess_enabled = os.getenv("KNOWLEDGE_OCR_POSTPROCESS_ENABLED", "true").lower() == "true"

            ocr_units_attempted = 0
            ocr_units_extracted = 0
            ocr_units_deferred = 0
            ocr_text_chars = 0
            extracted_slides: set[int] = set()

            if should_ocr and policy.slide_render_enabled and file_path:
                try:
                    from scripts.slide_renderer import render_slides_as_images

                    rendered_slides = render_slides_as_images(Path(str(file_path)))
                    if rendered_slides:
                        slide_rendered = True
                        print(f"[OCR] Slide rendering: {len(rendered_slides)} slides from {file_path}")
                        for slide_num, png_bytes in rendered_slides:
                            if ocr_units_attempted >= policy.ocr_max_ppt_slides:
                                ocr_units_deferred += 1
                                continue
                            ocr_units_attempted += 1
                            cls._emit_status(
                                heartbeat_fn,
                                f"ocr_processing ppt slide={slide_num}/{len(rendered_slides)}",
                            )
                            ocr_text = cls._ocr_slide_image(
                                png_bytes, slide_num,
                                preprocess=ocr_preprocess_enabled,
                                layout_analysis=policy.layout_analysis_enabled,
                                postprocess=ocr_postprocess_enabled,
                            )
                            if ocr_text:
                                text_parts.append(f"[Slide {slide_num} OCR]\n{ocr_text}")
                                ocr_units_extracted += 1
                                ocr_text_chars += cls._text_chars(ocr_text)
                                extracted_slides.add(slide_num)
                            if heartbeat_fn and slide_num % 5 == 0:
                                heartbeat_fn(f"slide_render_ocr: {slide_num}/{len(rendered_slides)}")
                except Exception as render_err:
                    print(f"[OCR] Slide rendering failed, falling back to shape OCR: {render_err}")
                    slide_rendered = False

            # ============================================================
            # Fallback: Shape-by-shape OCR (original path)
            # ============================================================
            if should_ocr and not slide_rendered:
                timed_out_images = []  # (slide_num, png_bytes) — timeout된 이미지 재시도용
                ocr_processed = 0
                ocr_total = len(image_shapes)
                attempted_slides: set[int] = set()
                for slide_num, image_bytes in image_shapes:
                    if slide_num not in attempted_slides and len(attempted_slides) >= policy.ocr_max_ppt_slides:
                        ocr_units_deferred += 1
                        continue
                    try:
                        ocr = cls._get_ocr_instance()
                        if ocr is None:
                            break
                        from PIL import Image
                        img = Image.open(io.BytesIO(image_bytes))
                        img = cls._resize_image_if_needed(img)
                        if img is None:
                            continue  # 극단적 종횡비/너무 작은 이미지 스킵
                        if slide_num not in attempted_slides:
                            attempted_slides.add(slide_num)
                            ocr_units_attempted += 1
                            cls._emit_status(
                                heartbeat_fn,
                                f"ocr_processing ppt slide={slide_num}/{len(prs.slides)}",
                            )

                        # Keep original color image for layout analysis (PP-Structure
                        # uses color cues for table/figure boundary detection)
                        img_original = img.copy()
                        if img_original.mode != "RGB":
                            img_original = img_original.convert("RGB")

                        # Phase 2: Image preprocessing (for standard OCR path)
                        if ocr_preprocess_enabled:
                            try:
                                from scripts.ocr_preprocessor import preprocess_for_ocr
                                img = preprocess_for_ocr(img, mode="auto")
                            except Exception as preproc_err:
                                print(f"[OCR] Preprocess failed, using original: {preproc_err}")
                                if img.mode != "RGB":
                                    img = img.convert("RGB")
                        else:
                            if img.mode != "RGB":
                                img = img.convert("RGB")

                        img_buffer = io.BytesIO()
                        img.save(img_buffer, format="PNG")
                        png_bytes = img_buffer.getvalue()

                        # Phase 2: Layout analysis (PP-Structure) for shape images
                        # Uses ORIGINAL color image (not preprocessed) for better region detection
                        ocr_text = None
                        ocr_conf = 0.0
                        if policy.layout_analysis_enabled:
                            try:
                                from scripts.ocr_preprocessor import analyze_layout_and_ocr
                                regions = analyze_layout_and_ocr(img_original)
                                if regions:
                                    ocr_text = "\n".join(r["content"] for r in regions if r.get("content"))
                                    ocr_conf = 0.7  # Layout-based OCR default confidence
                            except Exception:
                                pass  # Fall through to standard OCR

                        if not ocr_text:
                            # subprocess 격리로 OCR 실행 (SIGSEGV가 메인 프로세스를 죽이지 않도록)
                            with AttachmentParser._ocr_lock:
                                ocr_text, ocr_conf, _ = AttachmentParser._ocr_extract_safe(
                                    png_bytes, f"slide_{slide_num}"
                                )

                        # Phase 3: Post-processing
                        if ocr_text and ocr_postprocess_enabled:
                            try:
                                from scripts.ocr_postprocessor import postprocess_ocr_text
                                ocr_text, ocr_conf = postprocess_ocr_text(ocr_text, ocr_conf)
                            except Exception as postproc_err:
                                print(f"[OCR] Postprocess failed: {postproc_err}")

                        if ocr_text and ocr_conf > 0.3:
                            text_parts.append(f"[Slide {slide_num} Image OCR]\n{ocr_text}")
                            if slide_num not in extracted_slides:
                                extracted_slides.add(slide_num)
                                ocr_units_extracted += 1
                            ocr_text_chars += cls._text_chars(ocr_text)
                        elif ocr_text is None:
                            # timeout 또는 SIGSEGV — 재시도 대상
                            timed_out_images.append((slide_num, png_bytes))
                        # Temporal heartbeat (thread-safe) — 10장마다 전송
                        ocr_processed += 1
                        if heartbeat_fn and ocr_processed % 10 == 0:
                            heartbeat_fn(f"ocr: {ocr_processed}/{ocr_total} images, slide_{slide_num}")
                    except Exception as ocr_err:
                        print(f"[OCR Warning] Slide {slide_num} image: {ocr_err}")

                # timeout된 이미지 순차 재시도 (동시 경합 없이 단독 실행)
                if timed_out_images:
                    print(f"[OCR Retry] {len(timed_out_images)} timed-out images, retrying sequentially...")
                    if heartbeat_fn:
                        heartbeat_fn(f"ocr_retry: {len(timed_out_images)} images to retry")
                    for slide_num, png_bytes in timed_out_images:
                        if slide_num not in attempted_slides and len(attempted_slides) >= policy.ocr_max_ppt_slides:
                            ocr_units_deferred += 1
                            continue
                        try:
                            if slide_num not in attempted_slides:
                                attempted_slides.add(slide_num)
                                ocr_units_attempted += 1
                            with AttachmentParser._ocr_lock:
                                ocr_text, ocr_conf, _ = AttachmentParser._ocr_extract_safe(
                                    png_bytes, f"retry_slide_{slide_num}"
                                )
                            # Post-processing on retry
                            if ocr_text and ocr_postprocess_enabled:
                                try:
                                    from scripts.ocr_postprocessor import postprocess_ocr_text
                                    ocr_text, ocr_conf = postprocess_ocr_text(ocr_text, ocr_conf)
                                except Exception:
                                    pass
                            if ocr_text and ocr_conf > 0.3:
                                text_parts.append(f"[Slide {slide_num} Image OCR]\n{ocr_text}")
                                if slide_num not in extracted_slides:
                                    extracted_slides.add(slide_num)
                                    ocr_units_extracted += 1
                                ocr_text_chars += cls._text_chars(ocr_text)
                                print(f"[OCR Retry] slide_{slide_num}: OK ({len(ocr_text)} chars)")
                            else:
                                print(f"[OCR Retry] slide_{slide_num}: still failed")
                        except Exception as retry_err:
                            print(f"[OCR Retry] slide_{slide_num}: error - {retry_err}")

            full_text = "\n\n".join(text_parts)

            # ============================================================
            # Final fallback: PPTX → PDF → image OCR
            # If both text extraction and slide rendering produced nothing
            # (e.g., PPTX with embedded images that python-pptx can't read,
            # or LibreOffice rendering failed), convert PPTX to PDF and use
            # the PDF image OCR pipeline as last resort.
            # ============================================================
            if should_ocr and (not full_text.strip() or len(full_text.strip()) < 50):
                print(f"[PPT] Empty result after all extraction ({len(full_text)} chars), trying PDF fallback")
                try:
                    import subprocess
                    import tempfile

                    from scripts.slide_renderer import _find_soffice

                    soffice = _find_soffice()
                    if soffice:
                        with tempfile.TemporaryDirectory(prefix="pptx_pdf_fallback_") as tmpdir:
                            lo_profile = os.path.join(tmpdir, "lo_profile")
                            os.makedirs(lo_profile, exist_ok=True)

                            result = subprocess.run(
                                [
                                    soffice, "--headless", "--norestore",
                                    f"-env:UserInstallation=file://{lo_profile}",
                                    "--convert-to", "pdf",
                                    "--outdir", tmpdir,
                                    str(file_path),
                                ],
                                capture_output=True, text=True, timeout=120,
                            )
                            if result.returncode == 0:
                                pdf_files = list(Path(tmpdir).glob("*.pdf"))
                                if pdf_files:
                                    print(f"[PPT] PDF fallback: converted {file_path.name} → {pdf_files[0].name}")
                                    pdf_result = cls.parse_pdf(
                                        pdf_files[0], heartbeat_fn=heartbeat_fn,
                                    )
                                    if pdf_result.extracted_text.strip() and len(pdf_result.extracted_text.strip()) > len(full_text.strip()):
                                        full_text = pdf_result.extracted_text
                                        tables = pdf_result.extracted_tables or tables
                                        ocr_text_chars = max(ocr_text_chars, pdf_result.ocr_text_chars)
                                        print(f"[PPT] PDF fallback produced {len(full_text)} chars")
                except Exception as fb_err:
                    print(f"[PPT] PDF fallback failed: {fb_err}")

            if ocr_units_deferred > 0:
                cls._emit_status(
                    heartbeat_fn,
                    f"ocr_skipped_budget ppt deferred={ocr_units_deferred}",
                )
            confidence = 0.85 if full_text.strip() else 0.0

            skip_reason = None
            if policy.attachment_ocr_mode == "off":
                skip_reason = "disabled"
            elif policy.attachment_ocr_mode == "auto" and not should_ocr:
                skip_reason = "native_text_sufficient"
            elif ocr_units_deferred > 0:
                skip_reason = "budget_exceeded"

            return AttachmentParseResult(
                extracted_text=full_text,
                extracted_tables=tables,
                confidence=confidence,
                ocr_mode=policy.attachment_ocr_mode,
                ocr_applied=ocr_units_extracted > 0,
                ocr_skip_reason=skip_reason,
                ocr_units_attempted=ocr_units_attempted,
                ocr_units_extracted=ocr_units_extracted,
                ocr_units_deferred=ocr_units_deferred,
                native_text_chars=native_text_chars,
                ocr_text_chars=ocr_text_chars,
            )

        except Exception as e:
            return AttachmentParseResult(
                extracted_text=f"[PPT 파싱 오류: {e}]",
                extracted_tables=[],
                confidence=0.0,
                ocr_mode=cls.current_policy().attachment_ocr_mode,
                ocr_skip_reason="parse_error",
            )

    # 싱글톤 OCR 인스턴스 (메모리 효율화)
    _ocr_instance = None
    _ocr_type = None  # "paddle" only (amd64 Crawler Pod)
    _ocr_lock = __import__("threading").Lock()  # PaddleOCR는 thread-safe 보장 없음

    @classmethod
    def _get_ocr_instance(cls):
        """싱글톤 PaddleOCR 인스턴스 반환 (amd64 only, no fallback)."""
        if cls._ocr_instance is None:
            try:
                from src.nlp.ocr.paddle_ocr_provider import PaddleOCRProvider
                cls._ocr_instance = PaddleOCRProvider()
                cls._ocr_type = "paddle"
                print("[OCR] PaddleOCR singleton created")
            except ImportError:
                print("[OCR] PaddleOCR not available (requires amd64)")
                return None
        return cls._ocr_instance

    @classmethod
    def cleanup_ocr(cls):
        """OCR 인스턴스 정리 및 메모리 해제."""
        import gc
        cls._ocr_instance = None
        cls._ocr_type = None
        # Shutdown subprocess pool if active
        if cls._ocr_process_pool is not None:
            try:
                cls._ocr_process_pool.shutdown(wait=False)
            except Exception:
                pass
            cls._ocr_process_pool = None
        gc.collect()
        print("[OCR] 메모리 정리 완료")

    # --- Subprocess-isolated OCR (SIGSEGV defense) ---
    _ocr_process_pool = None
    _ocr_pool_lock = __import__("threading").Lock()

    @classmethod
    def _ocr_extract_safe(
        cls, image_bytes: bytes, file_name: str = "", timeout: int = 1800
    ) -> tuple[str | None, float, list]:
        """Execute OCR in a forked subprocess to survive PaddleOCR SIGSEGV.

        PaddleOCR's C++ inference engine can crash with SIGSEGV on certain
        images. Python try/except cannot catch OS signals (signal 11).
        By running OCR in a forked subprocess via ProcessPoolExecutor:
        - SIGSEGV kills only the child process
        - Parent detects BrokenProcessPool, recreates pool, skips image
        - Fork inherits loaded model via COW (no ~30s model reload)
        """
        from concurrent.futures import ProcessPoolExecutor
        from concurrent.futures.process import BrokenProcessPool

        import multiprocessing as mp

        with cls._ocr_pool_lock:
            if cls._ocr_process_pool is None:
                ctx = mp.get_context("fork")
                cls._ocr_process_pool = ProcessPoolExecutor(
                    max_workers=1, mp_context=ctx
                )
            pool = cls._ocr_process_pool

        try:
            future = pool.submit(_ocr_worker_fn, image_bytes)
            return future.result(timeout=timeout)
        except BrokenProcessPool:
            print(
                f"[OCR SIGSEGV] Worker crashed on {file_name} "
                "— restarting pool, skipping image"
            )
            with cls._ocr_pool_lock:
                cls._ocr_process_pool = None
            return None, 0.0, []
        except TimeoutError:
            print(f"[OCR Timeout] {file_name} exceeded {timeout}s — skipped")
            return None, 0.0, []
        except Exception as e:
            print(f"[OCR Error] {file_name}: {e} — skipped")
            return None, 0.0, []

    # PaddleOCR PP-OCRv5 det 모델이 극단적 종횡비에서 SIGSEGV 발생
    _OCR_MIN_DIMENSION = 32  # 최소 폭/높이 (px)
    _OCR_MAX_ASPECT_RATIO = 8.0  # 최대 종횡비

    @staticmethod
    def _resize_image_if_needed(img, max_size: int = 2048):
        """큰 이미지 리사이즈 (메모리 최적화).

        Args:
            img: PIL Image 객체
            max_size: 최대 폭/높이 (기본: 2048px)

        Returns:
            리사이즈된 이미지 (또는 원본).
            극단적 종횡비(>8:1)나 너무 작은(<32px) 이미지는 None 반환.
        """
        width, height = img.size

        # 원본이 이미 너무 작으면 OCR 무의미
        if width < AttachmentParser._OCR_MIN_DIMENSION or height < AttachmentParser._OCR_MIN_DIMENSION:
            print(f"[OCR] 이미지 스킵 (너무 작음): {width}x{height}")
            return None

        if width > max_size or height > max_size:
            ratio = min(max_size / width, max_size / height)
            new_size = (int(width * ratio), int(height * ratio))

            # 리사이즈 후 최소 크기 검증
            if new_size[0] < AttachmentParser._OCR_MIN_DIMENSION or new_size[1] < AttachmentParser._OCR_MIN_DIMENSION:
                print(f"[OCR] 이미지 스킵 (리사이즈 후 너무 작음): {width}x{height} → {new_size[0]}x{new_size[1]}")
                return None

            img = img.resize(new_size, resample=3)  # Pillow LANCZOS = 3
            print(f"[OCR] 이미지 리사이즈: {width}x{height} → {new_size[0]}x{new_size[1]}")

        # 종횡비 검증 (리사이즈 후 기준) — 극단적이면 패딩으로 보정
        w, h = img.size
        aspect = max(w, h) / max(min(w, h), 1)
        if aspect > AttachmentParser._OCR_MAX_ASPECT_RATIO:
            # 짧은 쪽에 흰색 패딩 추가하여 8:1 이내로 보정
            target_short = max(w, h) // int(AttachmentParser._OCR_MAX_ASPECT_RATIO)
            from PIL import Image as _PilImage
            if w > h:
                padded = _PilImage.new("RGB", (w, target_short), (255, 255, 255))
                padded.paste(img, (0, (target_short - h) // 2))
            else:
                padded = _PilImage.new("RGB", (target_short, h), (255, 255, 255))
                padded.paste(img, ((target_short - w) // 2, 0))
            print(f"[OCR] 이미지 패딩 (종횡비 {aspect:.1f}:1 → 8:1): {w}x{h} → {padded.size[0]}x{padded.size[1]}")
            return padded

        return img

    @classmethod
    def _ocr_slide_image(
        cls,
        png_bytes: bytes,
        slide_num: int,
        preprocess: bool = True,
        layout_analysis: bool = True,
        postprocess: bool = True,
    ) -> str | None:
        """OCR a rendered slide image with preprocessing and layout analysis.

        Args:
            png_bytes: PNG image bytes of the rendered slide.
            slide_num: Slide number (for logging).
            preprocess: Apply OpenCV preprocessing.
            layout_analysis: Use PP-Structure layout analysis.
            postprocess: Apply OCR text post-processing.

        Returns:
            Extracted text or None if OCR fails.
        """
        from PIL import Image

        try:
            img = Image.open(io.BytesIO(png_bytes))
            if img.mode != "RGB":
                img = img.convert("RGB")

            # Keep original color image for layout analysis (PP-Structure
            # uses color cues for region boundary detection)
            img_original = img.copy()

            # Phase 2a: Image preprocessing (for standard OCR fallback path)
            if preprocess:
                try:
                    from scripts.ocr_preprocessor import preprocess_for_ocr
                    img = preprocess_for_ocr(img, mode="slide")
                except Exception as e:
                    print(f"[OCR] Slide {slide_num} preprocess failed: {e}")

            # Phase 2b: Layout analysis (PP-Structure) — uses ORIGINAL color image
            ocr_text = None
            if layout_analysis:
                try:
                    from scripts.ocr_preprocessor import analyze_layout_and_ocr
                    regions = analyze_layout_and_ocr(img_original)
                    if regions:
                        ocr_text = "\n".join(r["content"] for r in regions if r.get("content"))
                        print(f"[OCR] Slide {slide_num}: layout analysis found {len(regions)} regions")
                except Exception as e:
                    print(f"[OCR] Slide {slide_num} layout analysis failed: {e}")

            # Fallback: standard OCR on preprocessed image
            if not ocr_text:
                img_buffer = io.BytesIO()
                img.save(img_buffer, format="PNG")
                with cls._ocr_lock:
                    ocr_text, ocr_conf, _ = cls._ocr_extract_safe(
                        img_buffer.getvalue(), f"rendered_slide_{slide_num}"
                    )
                if not ocr_text or ocr_conf <= 0.3:
                    return None

            # Phase 3: Post-processing
            if ocr_text and postprocess:
                try:
                    from scripts.ocr_postprocessor import postprocess_ocr_text
                    ocr_text, _ = postprocess_ocr_text(ocr_text)
                except Exception as e:
                    print(f"[OCR] Slide {slide_num} postprocess failed: {e}")

            # Phase 4: OCR noise filter — remove lines with repeated characters
            # (e.g., "폐폐폐폐폐" from low-quality PDF OCR)
            if ocr_text:
                clean_lines = []
                for line in ocr_text.split("\n"):
                    stripped = line.strip()
                    if not stripped:
                        continue
                    # Detect repeated single-char patterns (e.g., "폐폐폐폐폐")
                    if len(stripped) >= 3:
                        unique_chars = set(stripped.replace(" ", ""))
                        if len(unique_chars) <= 2 and len(stripped) >= 5:
                            continue  # Skip noise line (1-2 unique chars repeated)
                    clean_lines.append(line)
                ocr_text = "\n".join(clean_lines)

            return ocr_text if ocr_text and ocr_text.strip() else None

        except Exception as e:
            print(f"[OCR] Slide {slide_num} OCR error: {e}")
            return None

    @classmethod
    def _parse_image_sync(
        cls,
        file_path: Path,
        content: bytes,
        use_ocr: bool = True,
    ) -> AttachmentParseResult:
        """이미지 OCR 및 메타데이터 추출 (동기 내부 구현).

        PaddleOCR(한국어 특화)을 사용하여 텍스트와 표를 추출합니다.
        싱글톤 패턴으로 OCR 인스턴스를 재사용하여 메모리를 최적화합니다.
        """
        try:
            policy = cls.current_policy()
            from PIL import Image
            import gc

            img = Image.open(io.BytesIO(content))
            width, height = img.size
            format_type = img.format or "unknown"

            # 이미지 메타데이터
            metadata_text = f"[Image: {width}x{height}, {format_type}, {len(content):,} bytes]"
            tables = []
            confidence = 0.5  # 기본값

            if policy.attachment_ocr_mode == "off" or not use_ocr:
                return AttachmentParseResult(
                    extracted_text=metadata_text,
                    extracted_tables=tables,
                    confidence=confidence,
                    ocr_mode=policy.attachment_ocr_mode,
                    ocr_skip_reason="disabled",
                )

            if policy.ocr_max_images_per_attachment <= 0:
                cls._emit_status(None, "ocr_skipped_budget image deferred=1")
                return AttachmentParseResult(
                    extracted_text=metadata_text,
                    extracted_tables=tables,
                    confidence=confidence,
                    ocr_mode=policy.attachment_ocr_mode,
                    ocr_skip_reason="budget_exceeded",
                    ocr_units_deferred=1,
                )

            # OCR 수행 조건: 10MB 미만 && use_ocr 플래그
            if len(content) < 10_000_000:
                try:
                    cls._emit_status(None, f"ocr_processing image file={file_path.name}")
                    # 싱글톤 OCR 인스턴스 사용 (메모리 최적화)
                    ocr = cls._get_ocr_instance()
                    if ocr is None:
                        return AttachmentParseResult(
                            extracted_text=metadata_text,
                            extracted_tables=tables,
                            confidence=confidence,
                            ocr_mode=policy.attachment_ocr_mode,
                            ocr_skip_reason="ocr_unavailable",
                        )

                    # 큰 이미지 리사이즈 (메모리 최적화 + 극단적 종횡비 방어)
                    img = cls._resize_image_if_needed(img)
                    if img is None:
                        return AttachmentParseResult(
                            extracted_text=metadata_text,
                            extracted_tables=tables,
                            confidence=confidence,
                            ocr_mode=policy.attachment_ocr_mode,
                            ocr_skip_reason="guard_rejected",
                        )

                    # RGB 변환 및 바이트 변환
                    if img.mode != "RGB":
                        img = img.convert("RGB")

                    # 리사이즈된 이미지를 바이트로 변환
                    img_buffer = io.BytesIO()
                    img.save(img_buffer, format="PNG")
                    resized_content = img_buffer.getvalue()

                    # OCR 실행 (subprocess isolation: PaddleOCR SIGSEGV 방어)
                    # SIGSEGV는 try/except로 잡을 수 없음 → 별도 프로세스에서 실행
                    with cls._ocr_lock:
                        ocr_text, ocr_conf, ocr_tables = cls._ocr_extract_safe(
                            resized_content, file_path.name
                        )

                    if ocr_text and ocr_conf > 0.3:
                        # OCR 성공
                        extracted_text = ocr_text
                        confidence = ocr_conf
                        tables = ocr_tables

                        # 메타데이터 + OCR 텍스트
                        full_text = f"{metadata_text}\n\n{extracted_text}"

                        # 메모리 정리 (이미지별)
                        del img, img_buffer, resized_content
                        gc.collect()

                        return AttachmentParseResult(
                            extracted_text=full_text,
                            extracted_tables=tables,
                            confidence=confidence,
                            ocr_mode=policy.attachment_ocr_mode,
                            ocr_applied=True,
                            ocr_units_attempted=1,
                            ocr_units_extracted=1,
                            ocr_text_chars=cls._text_chars(extracted_text),
                        )

                except Exception as ocr_error:
                    # OCR 실패 시 로그만 남기고 메타데이터 반환
                    print(f"[OCR Warning] {file_path.name}: {ocr_error}")

            # OCR 실패 또는 미수행 시 메타데이터만 반환
            return AttachmentParseResult(
                extracted_text=metadata_text,
                extracted_tables=tables,
                confidence=confidence,
                ocr_mode=policy.attachment_ocr_mode,
                ocr_skip_reason="ocr_failed" if len(content) < 10_000_000 else "image_too_large",
                ocr_units_attempted=1 if len(content) < 10_000_000 else 0,
            )

        except Exception as e:
            return AttachmentParseResult(
                extracted_text=f"[이미지 파싱 오류: {e}]",
                extracted_tables=[],
                confidence=0.0,
                ocr_mode=cls.current_policy().attachment_ocr_mode,
                ocr_skip_reason="parse_error",
            )

    @classmethod
    def parse_image(
        cls,
        file_path: Path,
        content: bytes,
        use_ocr: bool = True,
    ) -> AttachmentParseResult:
        """이미지 OCR 및 메타데이터 추출 (동기 호출용 래퍼)."""
        return cls._parse_image_sync(file_path, content, use_ocr)

    @classmethod
    async def parse_image_async(
        cls,
        file_path: Path,
        content: bytes,
        use_ocr: bool = True,
    ) -> AttachmentParseResult:
        """이미지 OCR 비동기 처리.

        CPU-bound OCR을 asyncio.to_thread로 오프로드하여
        다른 첨부파일 다운로드와 병렬 처리합니다.
        OMP_NUM_THREADS=2 설정으로 스레드당 2코어만 사용하여
        동시 여러 OCR 처리가 가능합니다.
        """
        return await asyncio.to_thread(cls._parse_image_sync, file_path, content, use_ocr)


# =============================================================================
# Data Models
# =============================================================================
@dataclass
class ExtractedTable:
    """추출된 테이블"""
    headers: list[str]
    rows: list[dict[str, str]]
    section: str | None = None
    table_type: str | None = None  # owner, system, schedule, status


@dataclass
class ExtractedMention:
    """추출된 @멘션"""
    user_id: str | None
    display_name: str | None
    context: str  # 멘션 주변 텍스트
    email: str | None = None  # 이메일 주소


@dataclass
class ExtractedEmail:
    """추출된 이메일 링크"""
    email: str
    display_name: str | None
    context: str  # 주변 텍스트


@dataclass
class ExtractedMacro:
    """추출된 Confluence 매크로"""
    macro_type: str  # expand, panel, note, info, warning, status, etc.
    title: str | None
    content: str
    parameters: dict = field(default_factory=dict)


@dataclass
class ExtractedComment:
    """추출된 댓글"""
    comment_id: str
    author: str
    author_email: str | None
    content: str
    created_at: str
    parent_id: str | None = None  # 대댓글인 경우


@dataclass
class ExtractedLabel:
    """추출된 라벨(태그)"""
    name: str
    prefix: str | None = None  # global, my, etc.


@dataclass
class ExtractedLink:
    """추출된 링크"""
    link_type: str  # "internal" or "external"
    target_page_id: str | None = None  # 내부 링크인 경우 페이지 ID
    target_url: str | None = None  # 외부 링크인 경우 URL
    anchor_text: str | None = None  # 링크 텍스트
    context: str = ""  # 주변 텍스트


@dataclass
class ExtractedRestriction:
    """추출된 접근 제한 정보"""
    operation: str  # "read" or "update"
    restriction_type: str  # "user" or "group"
    name: str  # 사용자명 또는 그룹명
    account_id: str | None = None  # 사용자인 경우 account ID


@dataclass
class AttachmentContent:
    """첨부파일 내용"""
    id: str
    filename: str
    media_type: str
    file_size: int
    download_path: str | None = None
    download_url: str | None = None
    extracted_text: str | None = None
    extracted_tables: list[dict] = field(default_factory=list)
    ocr_confidence: float | None = None
    parse_error: str | None = None
    has_visual_content: bool = False
    visual_analysis_version: str | None = None
    ocr_mode: str | None = None
    ocr_applied: bool | None = None
    ocr_skip_reason: str | None = None
    ocr_units_attempted: int = 0
    ocr_units_extracted: int = 0
    ocr_units_deferred: int = 0
    native_text_chars: int = 0
    ocr_text_chars: int = 0


@dataclass
class FullPageContent:
    """전체 페이지 내용"""
    # 필수 필드 (기본값 없음)
    page_id: str
    title: str
    # 본문 (3중 저장)
    content_text: str           # 1. 전체 plain text (검색용)
    content_html: str           # 2. HTML 원본 (폴백/재파싱용)
    content_preview: str        # 미리보기 (UI용)
    # 구조
    tables: list[ExtractedTable]
    mentions: list[ExtractedMention]
    sections: list[dict]        # {"level": 1, "title": "...", "content": "..."}
    # 메타
    creator: str
    last_modifier: str
    version: int
    url: str
    created_at: str
    updated_at: str
    # 선택적 필드 (기본값 있음 - 반드시 필수 필드 뒤에 위치)
    content_ir: dict | None = None  # 3. Structured IR (RAG/임베딩용)
    code_blocks: list[dict] = field(default_factory=list)  # 코드 블록
    creator_name: str | None = None
    creator_team: str | None = None
    creator_email: str | None = None  # NEW: 작성자 이메일
    attachments: list[AttachmentContent] = field(default_factory=list)
    # NEW: 추가 메타데이터
    labels: list[ExtractedLabel] = field(default_factory=list)  # 라벨/태그
    comments: list[ExtractedComment] = field(default_factory=list)  # 댓글
    emails: list[ExtractedEmail] = field(default_factory=list)  # 이메일 링크
    macros: list[ExtractedMacro] = field(default_factory=list)  # 매크로
    space_key: str | None = None  # 스페이스 키
    ancestors: list[dict] = field(default_factory=list)  # 상위 페이지 계층
    # NEW: 링크 및 권한
    internal_links: list[ExtractedLink] = field(default_factory=list)  # 내부 문서 링크
    external_links: list[ExtractedLink] = field(default_factory=list)  # 외부 URL 링크
    restrictions: list[ExtractedRestriction] = field(default_factory=list)  # 접근 제한
    version_history: list[dict] = field(default_factory=list)  # 버전 이력 메타데이터


@dataclass
class CrawlSpaceResult:
    """단일 소스 크롤링 결과."""

    pages: list[FullPageContent]
    page_dicts: list[dict]
    interrupted: bool = False
    jsonl_path: str = ""  # 스트리밍 모드에서 사용할 JSONL 경로
    source_key: str = ""  # checkpoint 정리용


# =============================================================================
# HTML Parsers
# =============================================================================
class TableExtractor(HTMLParser):
    """HTML에서 테이블 추출"""

    def __init__(self):
        super().__init__()
        self.tables: list[ExtractedTable] = []
        self.current_table: dict | None = None
        self.current_row: list[str] = []
        self.current_cell: str = ""
        self.in_table = False
        self.in_header = False
        self.in_row = False
        self.in_cell = False

    def handle_starttag(self, tag, attrs):
        if tag == "table":
            self.in_table = True
            self.current_table = {"headers": [], "rows": []}
        elif tag == "thead":
            self.in_header = True
        elif tag == "tr":
            self.in_row = True
            self.current_row = []
        elif tag in ("td", "th"):
            self.in_cell = True
            self.current_cell = ""

    def handle_endtag(self, tag):
        if tag == "table" and self.current_table:
            if self.current_table["headers"] or self.current_table["rows"]:
                # 첫 번째 행이 헤더일 수 있음
                if not self.current_table["headers"] and self.current_table["rows"]:
                    self.current_table["headers"] = self.current_table["rows"].pop(0) if self.current_table["rows"] else []

                # 테이블 타입 추론
                table_type = self._infer_table_type(self.current_table["headers"])

                self.tables.append(ExtractedTable(
                    headers=self.current_table["headers"],
                    rows=[
                        dict(zip(self.current_table["headers"], row))
                        for row in self.current_table["rows"]
                        if len(row) == len(self.current_table["headers"])
                    ],
                    table_type=table_type,
                ))
            self.in_table = False
            self.current_table = None
        elif tag == "thead":
            self.in_header = False
        elif tag == "tr" and self.current_table:
            if self.in_header:
                self.current_table["headers"] = self.current_row
            else:
                self.current_table["rows"].append(self.current_row)
            self.in_row = False
        elif tag in ("td", "th"):
            self.current_row.append(self.current_cell.strip())
            self.in_cell = False

    def handle_data(self, data):
        if self.in_cell:
            self.current_cell += data

    def _infer_table_type(self, headers: list[str]) -> str | None:
        """테이블 타입 추론"""
        headers_str = " ".join(headers).lower()

        if any(kw in headers_str for kw in ["담당자", "담당", "pm", "tl"]):
            return "owner_table"
        elif any(kw in headers_str for kw in ["시스템", "서비스", "api", "url"]):
            return "system_table"
        elif any(kw in headers_str for kw in ["일정", "마감", "주기"]):
            return "schedule_table"
        elif any(kw in headers_str for kw in ["상태", "진행"]):
            return "status_table"
        return None


class MentionExtractor(HTMLParser):
    """HTML에서 @멘션 추출 (이메일 정보 포함)"""

    def __init__(self):
        super().__init__()
        self.mentions: list[ExtractedMention] = []
        self.current_text = ""
        self.in_link = False
        self.current_user_id: str | None = None

    def handle_starttag(self, tag, attrs):
        attrs_dict = dict(attrs)
        # Confluence 멘션 패턴
        if tag == "ri:user":
            user_id = attrs_dict.get("ri:account-id") or attrs_dict.get("ri:userkey")
            self.current_user_id = user_id
            self.mentions.append(ExtractedMention(
                user_id=user_id,
                display_name=None,
                context="",
            ))
        elif tag == "ac:link":
            self.in_link = True

    def handle_endtag(self, tag):
        if tag == "ac:link":
            self.in_link = False
            self.current_user_id = None

    def handle_data(self, data):
        self.current_text = data
        # @ 패턴 찾기
        for match in re.finditer(r"@([가-힣]+(?:\s[가-힣]+)?)", data):
            self.mentions.append(ExtractedMention(
                user_id=None,
                display_name=match.group(1),
                context=data[:100],
            ))
        # ac:link 내부의 텍스트가 사용자 이름일 수 있음
        if self.in_link and self.mentions and self.mentions[-1].display_name is None:
            self.mentions[-1].display_name = data.strip()


class EmailExtractor(HTMLParser):
    """HTML에서 mailto 이메일 링크 추출"""

    def __init__(self):
        super().__init__()
        self.emails: list[ExtractedEmail] = []
        self.current_email: str | None = None
        self.in_mailto_link = False
        self.link_text = ""
        self.context_buffer = ""

    def handle_starttag(self, tag, attrs):
        attrs_dict = dict(attrs)
        if tag == "a":
            href = attrs_dict.get("href", "")
            if href.startswith("mailto:"):
                self.current_email = href.replace("mailto:", "").split("?")[0]  # ?subject= 등 제거
                self.in_mailto_link = True
                self.link_text = ""

    def handle_endtag(self, tag):
        if tag == "a" and self.in_mailto_link and self.current_email:
            self.emails.append(ExtractedEmail(
                email=self.current_email,
                display_name=self.link_text.strip() if self.link_text.strip() else None,
                context=self.context_buffer[-100:] if self.context_buffer else "",
            ))
            self.in_mailto_link = False
            self.current_email = None

    def handle_data(self, data):
        self.context_buffer += data
        if len(self.context_buffer) > 200:
            self.context_buffer = self.context_buffer[-200:]
        if self.in_mailto_link:
            self.link_text += data


class MacroExtractor(HTMLParser):
    """HTML에서 Confluence 매크로 추출 (expand, panel, note, info, warning, status 등)"""

    # 추출 대상 매크로 타입
    TARGET_MACROS = {"expand", "panel", "note", "info", "warning", "tip", "status", "toc", "children", "excerpt"}

    def __init__(self):
        super().__init__()
        self.macros: list[ExtractedMacro] = []
        self.macro_stack: list[dict] = []  # 중첩 매크로 처리용
        self.current_param_name: str | None = None
        self.in_body = False
        self.body_content = ""

    def handle_starttag(self, tag, attrs):
        attrs_dict = dict(attrs)

        if tag == "ac:structured-macro":
            macro_name = attrs_dict.get("ac:name", "")
            if macro_name in self.TARGET_MACROS:
                self.macro_stack.append({
                    "type": macro_name,
                    "title": None,
                    "content": "",
                    "parameters": {},
                })

        elif tag == "ac:parameter" and self.macro_stack:
            self.current_param_name = attrs_dict.get("ac:name")

        elif tag == "ac:rich-text-body" and self.macro_stack:
            self.in_body = True
            self.body_content = ""

        elif tag == "ac:plain-text-body" and self.macro_stack:
            self.in_body = True
            self.body_content = ""

    def handle_endtag(self, tag):
        if tag == "ac:structured-macro" and self.macro_stack:
            macro_data = self.macro_stack.pop()
            self.macros.append(ExtractedMacro(
                macro_type=macro_data["type"],
                title=macro_data.get("title") or macro_data["parameters"].get("title"),
                content=macro_data["content"],
                parameters=macro_data["parameters"],
            ))

        elif tag == "ac:parameter":
            self.current_param_name = None

        elif tag in ("ac:rich-text-body", "ac:plain-text-body"):
            if self.macro_stack:
                self.macro_stack[-1]["content"] = self.body_content.strip()
            self.in_body = False

    def handle_data(self, data):
        if self.current_param_name and self.macro_stack:
            self.macro_stack[-1]["parameters"][self.current_param_name] = data.strip()
            # title 파라미터는 별도 저장
            if self.current_param_name == "title":
                self.macro_stack[-1]["title"] = data.strip()

        if self.in_body:
            self.body_content += data


class LinkExtractor(HTMLParser):
    """HTML에서 내부/외부 링크 추출"""

    # 무시할 URL 패턴 (스타일, 스크립트 등)
    IGNORE_PATTERNS = {"javascript:", "#", "data:", "blob:"}

    def __init__(self, base_url: str = ""):
        super().__init__()
        self.base_url = base_url
        self.internal_links: list[ExtractedLink] = []
        self.external_links: list[ExtractedLink] = []
        self.current_link: dict | None = None
        self.in_link = False
        self.link_text = ""
        self.context_buffer = ""

    def handle_starttag(self, tag, attrs):
        attrs_dict = dict(attrs)

        # Confluence 내부 링크: ac:link + ri:page
        if tag == "ac:link":
            self.current_link = {"type": "internal", "page_id": None, "anchor": None}
            self.in_link = True
            self.link_text = ""

        elif tag == "ri:page" and self.current_link:
            # 내부 페이지 링크
            content_id = attrs_dict.get("ri:content-id")
            if content_id:
                self.current_link["page_id"] = content_id

        elif tag == "ri:attachment" and self.current_link:
            # 첨부파일 링크는 건너뜀
            self.current_link = None
            self.in_link = False

        # 일반 a 태그 링크
        elif tag == "a":
            href = attrs_dict.get("href", "")

            # 무시할 패턴 체크
            if any(href.startswith(p) for p in self.IGNORE_PATTERNS):
                return

            # mailto는 EmailExtractor에서 처리
            if href.startswith("mailto:"):
                return

            self.in_link = True
            self.link_text = ""

            # 내부 링크 판별
            if "/pages/viewpage.action" in href or "/display/" in href:
                # Confluence 내부 링크
                page_id = None
                if "pageId=" in href:
                    try:
                        page_id = href.split("pageId=")[1].split("&")[0]
                    except (IndexError, ValueError):
                        pass
                self.current_link = {"type": "internal", "page_id": page_id, "url": href}
            elif href.startswith("http://") or href.startswith("https://"):
                # 외부 링크
                self.current_link = {"type": "external", "url": href}
            elif href.startswith("/"):
                # 상대 경로 (내부)
                self.current_link = {"type": "internal", "url": self.base_url + href, "page_id": None}

    def handle_endtag(self, tag):
        if tag == "ac:link" and self.current_link and self.current_link.get("page_id"):
            self.internal_links.append(ExtractedLink(
                link_type="internal",
                target_page_id=self.current_link.get("page_id"),
                anchor_text=self.link_text.strip() if self.link_text.strip() else None,
                context=self.context_buffer[-100:] if self.context_buffer else "",
            ))
            self.current_link = None
            self.in_link = False

        elif tag == "a" and self.current_link:
            link_type = self.current_link.get("type", "external")

            if link_type == "internal":
                self.internal_links.append(ExtractedLink(
                    link_type="internal",
                    target_page_id=self.current_link.get("page_id"),
                    target_url=self.current_link.get("url"),
                    anchor_text=self.link_text.strip() if self.link_text.strip() else None,
                    context=self.context_buffer[-100:] if self.context_buffer else "",
                ))
            else:
                self.external_links.append(ExtractedLink(
                    link_type="external",
                    target_url=self.current_link.get("url"),
                    anchor_text=self.link_text.strip() if self.link_text.strip() else None,
                    context=self.context_buffer[-100:] if self.context_buffer else "",
                ))

            self.current_link = None
            self.in_link = False

    def handle_data(self, data):
        self.context_buffer += data
        if len(self.context_buffer) > 200:
            self.context_buffer = self.context_buffer[-200:]

        if self.in_link:
            self.link_text += data


class SectionExtractor(HTMLParser):
    """HTML에서 섹션 구조 추출"""

    def __init__(self):
        super().__init__()
        self.sections: list[dict] = []
        self.current_heading: dict | None = None
        self.in_heading = False
        self.heading_text = ""

    def handle_starttag(self, tag, attrs):
        if tag in ("h1", "h2", "h3", "h4"):
            self.in_heading = True
            self.current_heading = {"level": int(tag[1]), "title": ""}
            self.heading_text = ""

    def handle_endtag(self, tag):
        if tag in ("h1", "h2", "h3", "h4") and self.current_heading:
            self.current_heading["title"] = self.heading_text.strip()
            self.sections.append(self.current_heading)
            self.in_heading = False
            self.current_heading = None

    def handle_data(self, data):
        if self.in_heading:
            self.heading_text += data


class PlainTextExtractor(HTMLParser):
    """HTML에서 plain text 추출"""

    def __init__(self):
        super().__init__()
        self.text_parts = []
        self.in_script = False
        self.in_style = False

    def handle_starttag(self, tag, attrs):
        if tag == "script":
            self.in_script = True
        elif tag == "style":
            self.in_style = True
        elif tag in ("br", "p", "div", "tr", "li", "h1", "h2", "h3", "h4"):
            self.text_parts.append("\n")

    def handle_endtag(self, tag):
        if tag == "script":
            self.in_script = False
        elif tag == "style":
            self.in_style = False

    def handle_data(self, data):
        if not self.in_script and not self.in_style:
            self.text_parts.append(data)

    def get_text(self) -> str:
        text = "".join(self.text_parts).strip()
        text = re.sub(r"\n{3,}", "\n\n", text)
        text = re.sub(r" {2,}", " ", text)
        return text


class CodeBlockExtractor(HTMLParser):
    """HTML에서 코드 블록 추출"""

    def __init__(self):
        super().__init__()
        self.code_blocks: list[dict] = []
        self.current_block: dict | None = None
        self.in_code = False
        self.code_content = ""

    def handle_starttag(self, tag, attrs):
        attrs_dict = dict(attrs)
        # Confluence 코드 매크로 패턴
        if tag == "ac:structured-macro" and attrs_dict.get("ac:name") == "code":
            self.current_block = {"language": None, "content": ""}
        elif tag == "ac:parameter" and self.current_block is not None:
            if attrs_dict.get("ac:name") == "language":
                pass  # handle_data에서 언어 추출
        elif tag == "ac:plain-text-body" and self.current_block is not None:
            self.in_code = True
            self.code_content = ""
        # 일반 pre/code 태그
        elif tag == "pre":
            self.current_block = {"language": None, "content": ""}
            self.in_code = True
            self.code_content = ""
        elif tag == "code" and self.current_block is None:
            # standalone code 태그
            lang = attrs_dict.get("class", "").replace("language-", "")
            self.current_block = {"language": lang if lang else None, "content": ""}
            self.in_code = True
            self.code_content = ""

    def handle_endtag(self, tag):
        if tag in ("ac:plain-text-body", "pre", "code") and self.current_block is not None:
            self.current_block["content"] = self.code_content.strip()
            if self.current_block["content"]:  # 빈 코드 블록 제외
                self.code_blocks.append(self.current_block)
            self.current_block = None
            self.in_code = False

    def handle_data(self, data):
        if self.in_code:
            self.code_content += data


@dataclass
class ExtractedCodeBlock:
    """추출된 코드 블록"""
    language: str | None
    content: str
    section: str | None = None


@dataclass
class StructuredIR:
    """Structured Intermediate Representation for RAG"""
    chunks: list[dict]  # 의미적 청크들
    sections: list[dict]  # 섹션 계층
    tables: list[dict]  # 구조화된 테이블
    code_blocks: list[dict]  # 코드 블록
    mentions: list[dict]  # @멘션


def generate_structured_ir(
    content_text: str,
    content_html: str,
    title: str,
    tables: list[ExtractedTable],
    sections: list[dict],
    mentions: list[ExtractedMention],
) -> dict:
    """
    HTML에서 Structured IR 생성

    RAG 최적화된 청킹을 위한 중간 표현 형식:
    - 테이블은 원자 단위로 보존
    - 코드 블록은 분리하지 않음
    - 섹션 계층 정보 포함
    """
    chunks = []
    chunk_id = 0

    # 코드 블록 추출
    code_extractor = CodeBlockExtractor()
    try:
        code_extractor.feed(content_html)
    except Exception:
        pass  # HTML 파싱 오류 무시
    code_blocks = code_extractor.code_blocks

    # 1. 섹션 헤더 청크
    for section in sections:
        chunk_id += 1
        chunks.append({
            "chunk_id": f"sec-{chunk_id:04d}",
            "type": "section_header",
            "level": section.get("level", 2),
            "content": section.get("title", ""),
        })

    # 2. 테이블 청크 (원자 단위 - 분리하지 않음)
    for i, table in enumerate(tables):
        chunk_id += 1
        # 테이블을 읽기 쉬운 텍스트로 변환
        table_text = _table_to_markdown(table)
        chunks.append({
            "chunk_id": f"tbl-{chunk_id:04d}",
            "type": "table",
            "headers": table.headers,
            "rows": table.rows,
            "table_type": table.table_type,
            "content": table_text,  # RAG용 텍스트 표현
            "row_count": len(table.rows),
        })

    # 3. 코드 블록 청크 (원자 단위)
    for i, code in enumerate(code_blocks):
        chunk_id += 1
        chunks.append({
            "chunk_id": f"code-{chunk_id:04d}",
            "type": "code_block",
            "language": code.get("language"),
            "content": code.get("content", ""),
        })

    # 4. 본문 단락 청크 (테이블, 코드 제외한 나머지)
    paragraphs = _split_into_paragraphs(content_text)
    for para in paragraphs:
        if len(para.strip()) > 50:  # 최소 청크 크기
            chunk_id += 1
            chunks.append({
                "chunk_id": f"para-{chunk_id:04d}",
                "type": "paragraph",
                "content": para.strip(),
            })

    # IR 구조 반환
    return {
        "title": title,
        "chunk_count": len(chunks),
        "chunks": chunks,
        "sections": [
            {"level": s.get("level", 2), "title": s.get("title", "")}
            for s in sections
        ],
        "tables": [
            {
                "headers": t.headers,
                "rows": t.rows,
                "table_type": t.table_type,
            }
            for t in tables
        ],
        "code_blocks": [
            {
                "language": c.get("language"),
                "content": c.get("content", ""),
            }
            for c in code_blocks
        ],
        "mentions": [
            {
                "user_id": m.user_id,
                "display_name": m.display_name,
                "context": m.context[:100] if m.context else "",
            }
            for m in mentions
        ],
    }


def _table_to_markdown(table: ExtractedTable) -> str:
    """테이블을 Markdown 형식으로 변환 (RAG용)"""
    lines = []
    if table.headers:
        lines.append("| " + " | ".join(table.headers) + " |")
        lines.append("|" + "|".join(["---"] * len(table.headers)) + "|")
    for row in table.rows:
        values = [str(row.get(h, "")).strip().replace("\n", " ") for h in table.headers]
        lines.append("| " + " | ".join(values) + " |")
    return "\n".join(lines)


def _split_into_paragraphs(text: str) -> list[str]:
    """본문을 단락으로 분리"""
    # 빈 줄 2개 이상으로 분리
    paragraphs = re.split(r"\n\n+", text)
    # 너무 긴 단락은 추가 분리
    result = []
    for para in paragraphs:
        if len(para) > 1000:
            # 문장 단위로 분리
            sentences = re.split(r"(?<=[.!?])\s+", para)
            current = ""
            for sent in sentences:
                if len(current) + len(sent) > 800:
                    if current:
                        result.append(current)
                    current = sent
                else:
                    current += (" " if current else "") + sent
            if current:
                result.append(current)
        else:
            result.append(para)
    return result


def extract_creator_info(creator: str) -> tuple[str | None, str | None]:
    """작성자에서 이름/팀 추출"""
    match = re.search(r"([가-힣]+)/([가-힣A-Za-z0-9]+(?:팀|스쿼드|실|부문|센터))", creator)
    if match:
        return match.group(1), match.group(2)
    return None, None


# =============================================================================
# Confluence Client
# =============================================================================
class ConfluenceFullClient:
    """전체 콘텐츠 크롤링 클라이언트"""

    CHECKPOINT_INTERVAL = 10  # 10페이지마다 체크포인트 저장

    def __init__(self, checkpoint_dir: Path | None = None, max_concurrent: int = 1, kb_id: str = ""):
        self.base_url = CONFLUENCE_BASE_URL
        self.headers = {
            "Authorization": f"Bearer {CONFLUENCE_PAT}",
            "Accept": "application/json",
        }
        _timeout = float(os.getenv("CONFLUENCE_CRAWL_TIMEOUT", "30"))
        self.client = httpx.AsyncClient(timeout=_timeout, verify=False, headers=self.headers)
        self.all_pages: list[FullPageContent] = []
        self.visited_pages: set[str] = set()  # 순환 참조 방지
        self.kb_id = kb_id  # KB ID for checkpoint validation

        # 병렬 크롤링 설정
        self._max_concurrent = max(1, max_concurrent)
        self._page_sem: asyncio.Semaphore | None = (
            asyncio.Semaphore(self._max_concurrent) if self._max_concurrent > 1 else None
        )

        # 체크포인트 설정
        self.checkpoint_dir = checkpoint_dir or OUTPUT_DIR
        self.checkpoint_file = self.checkpoint_dir / "checkpoint.json"
        self._pages_since_checkpoint = 0
        self._incremental_saved_count = 0  # 증분 저장된 페이지 수 (in current all_pages batch)
        self._total_pages_crawled = 0  # 총 크롤된 페이지 수 (메모리 해제 후에도 유지)
        self._shutdown_requested = False
        self._started_at = time.monotonic()
        self._runtime_stats: dict[str, Any] = {
            "pages_total": 0,
            "attachments_total": 0,
            "attachments_ocr_applied": 0,
            "attachments_ocr_skipped": 0,
            "attachments_ocr_deferred": 0,
            "pdf_pages_ocr_attempted": 0,
            "pdf_pages_ocr_deferred": 0,
            "ppt_slides_ocr_attempted": 0,
            "ppt_slides_ocr_deferred": 0,
            "native_text_chars_total": 0,
            "ocr_text_chars_total": 0,
            "attachment_ocr_mode": AttachmentParser.current_policy().attachment_ocr_mode,
        }

    def runtime_stats_path(self) -> Path:
        return self.checkpoint_dir / "crawl_runtime_stats.json"

    def _record_attachment_stats(self, attachment: AttachmentContent) -> None:
        stats = self._runtime_stats
        stats["attachments_total"] += 1
        stats["native_text_chars_total"] += attachment.native_text_chars or 0
        stats["ocr_text_chars_total"] += attachment.ocr_text_chars or 0

        if attachment.ocr_applied:
            stats["attachments_ocr_applied"] += 1
        elif attachment.ocr_units_deferred > 0:
            stats["attachments_ocr_deferred"] += 1
        elif attachment.ocr_skip_reason and attachment.ocr_skip_reason not in {"ocr_failed", "parse_error"}:
            stats["attachments_ocr_skipped"] += 1

        media_type = (attachment.media_type or "").lower()
        if "pdf" in media_type:
            stats["pdf_pages_ocr_attempted"] += attachment.ocr_units_attempted or 0
            stats["pdf_pages_ocr_deferred"] += attachment.ocr_units_deferred or 0
        elif "presentation" in media_type or "powerpoint" in media_type:
            stats["ppt_slides_ocr_attempted"] += attachment.ocr_units_attempted or 0
            stats["ppt_slides_ocr_deferred"] += attachment.ocr_units_deferred or 0

    def write_runtime_stats(self) -> None:
        payload = {
            **self._runtime_stats,
            "pages_total": self._total_pages_crawled,
            "elapsed_seconds": round(time.monotonic() - self._started_at, 2),
        }
        with open(self.runtime_stats_path(), "w", encoding="utf-8") as f:
            json.dump(payload, f, ensure_ascii=False, indent=2)

    def _s3_checkpoint_key(self) -> str | None:
        """S3 checkpoint key. None if S3 not configured."""
        bucket = os.getenv("KNOWLEDGE_FINGERPRINT_S3_BUCKET") or os.getenv("S3_UPLOAD_BUCKET")
        if not bucket or not self.kb_id:
            return None
        return f"knowledge/metadata/checkpoint_{self.kb_id}.json"

    def _s3_bucket(self) -> str | None:
        return os.getenv("KNOWLEDGE_FINGERPRINT_S3_BUCKET") or os.getenv("S3_UPLOAD_BUCKET")

    def _upload_checkpoint_s3(self, checkpoint_data: dict) -> None:
        """Upload checkpoint to S3 for persistence across pod restarts."""
        bucket = self._s3_bucket()
        key = self._s3_checkpoint_key()
        if not bucket or not key:
            return
        try:
            import boto3
            s3 = boto3.client("s3")
            s3.put_object(
                Bucket=bucket,
                Key=key,
                Body=json.dumps(checkpoint_data, ensure_ascii=False).encode("utf-8"),
                ContentType="application/json",
            )
            console.print(f"[dim]☁️ S3 체크포인트 저장: s3://{bucket}/{key}[/dim]")
        except Exception as e:
            console.print(f"[yellow]⚠️ S3 체크포인트 업로드 실패 (로컬은 정상): {e}[/yellow]")

    def _download_checkpoint_s3(self) -> dict | None:
        """Download checkpoint from S3 if local file doesn't exist."""
        bucket = self._s3_bucket()
        key = self._s3_checkpoint_key()
        if not bucket or not key:
            return None
        try:
            import boto3
            s3 = boto3.client("s3")
            resp = s3.get_object(Bucket=bucket, Key=key)
            data = json.loads(resp["Body"].read().decode("utf-8"))
            console.print(f"[green]☁️ S3에서 체크포인트 복원: s3://{bucket}/{key}[/green]")
            return data
        except Exception:
            return None

    def save_checkpoint(self, source_key: str) -> None:
        """현재 진행 상태를 체크포인트 파일에 저장 (로컬 + S3)"""
        checkpoint_data = {
            "source_key": source_key,
            "kb_id": self.kb_id,
            "visited_pages": list(self.visited_pages),
            "pages_count": self._total_pages_crawled,
            "last_page_id": self.all_pages[-1].page_id if self.all_pages else None,
            "last_page_title": self.all_pages[-1].title if self.all_pages else None,
            "saved_at": datetime.now(timezone.utc).isoformat(),
        }

        # 임시 파일에 먼저 저장 후 이동 (원자적 쓰기)
        temp_file = self.checkpoint_file.with_suffix(".tmp")
        with open(temp_file, "w", encoding="utf-8") as f:
            json.dump(checkpoint_data, f, ensure_ascii=False, indent=2)
        temp_file.rename(self.checkpoint_file)

        console.print(f"[dim]💾 체크포인트 저장: {len(self.visited_pages)}개 페이지[/dim]")

        # S3 백업 (Pod 재시작 후에도 resume 가능)
        self._upload_checkpoint_s3(checkpoint_data)

    def load_checkpoint(self, source_key: str) -> bool:
        """이전 체크포인트에서 상태 복원 (로컬 → S3 fallback)"""
        checkpoint_data = None

        if self.checkpoint_file.exists():
            try:
                with open(self.checkpoint_file, "r", encoding="utf-8") as f:
                    checkpoint_data = json.load(f)
            except Exception:
                checkpoint_data = None

        # 로컬 없으면 S3에서 복원 (Pod 재시작 후 resume 지원)
        if checkpoint_data is None:
            checkpoint_data = self._download_checkpoint_s3()
            if checkpoint_data is not None:
                # S3에서 복원한 데이터를 로컬에도 저장
                try:
                    with open(self.checkpoint_file, "w", encoding="utf-8") as f:
                        json.dump(checkpoint_data, f, ensure_ascii=False, indent=2)
                except Exception:
                    pass

        if checkpoint_data is None:
            return False

        try:

            # 같은 소스의 체크포인트인지 확인
            if checkpoint_data.get("source_key") != source_key:
                console.print(f"[yellow]⚠️ 체크포인트가 다른 소스({checkpoint_data.get('source_key')})의 것입니다. 무시합니다.[/yellow]")
                return False

            # KB ID 불일치 검증 (다른 컬렉션의 old checkpoint 방지)
            saved_kb_id = checkpoint_data.get("kb_id", "")
            if saved_kb_id and self.kb_id and saved_kb_id != self.kb_id:
                console.print(f"[yellow]⚠️ 체크포인트 KB 불일치 ({saved_kb_id} ≠ {self.kb_id}). 처음부터 시작합니다.[/yellow]")
                self.visited_pages.clear()
                return False

            # 방문한 페이지 복원 (증분 JSONL에서 이미 로드된 페이지와 merge)
            checkpoint_visited = set(checkpoint_data.get("visited_pages", []))
            self.visited_pages = self.visited_pages | checkpoint_visited
            self._total_pages_crawled = checkpoint_data.get("pages_count", 0)

            saved_at = checkpoint_data.get("saved_at", "알 수 없음")
            last_title = checkpoint_data.get("last_page_title", "알 수 없음")
            pages_count = checkpoint_data.get("pages_count", 0)

            console.print(Panel.fit(
                f"[bold green]♻️ 체크포인트에서 재개[/bold green]\n\n"
                f"📅 저장 시간: {saved_at}\n"
                f"📄 처리 완료: {pages_count}개 페이지\n"
                f"📝 마지막 페이지: {last_title}\n"
                f"⏭️ 건너뛸 페이지: {len(self.visited_pages)}개",
                border_style="green",
            ))

            return True
        except Exception as e:
            console.print(f"[red]체크포인트 로드 실패: {e}[/red]")
            return False

    def clear_checkpoint(self) -> None:
        """체크포인트 파일 삭제"""
        if self.checkpoint_file.exists():
            self.checkpoint_file.unlink()
            console.print("[dim]🗑️ 체크포인트 삭제됨[/dim]")

    def _get_incremental_path(self, source_key: str) -> Path:
        """증분 저장 JSONL 파일 경로"""
        safe_key = re.sub(r'[^\w]', '_', source_key)
        return self.checkpoint_dir / f"incremental_{safe_key}.jsonl"

    @property
    def shutdown_requested(self) -> bool:
        """종료 요청 여부."""
        return self._shutdown_requested

    def request_shutdown(self) -> None:
        """크롤링 중단 플래그 설정."""
        self._shutdown_requested = True

    def _truncate_partial_jsonl_tail(self, jsonl_path: Path) -> None:
        """JSONL append 실패 시 마지막 불완전 라인 제거."""
        if not jsonl_path.exists():
            return

        try:
            with open(jsonl_path, "rb+") as f:
                f.seek(0, os.SEEK_END)
                file_size = f.tell()
                if file_size == 0:
                    return

                f.seek(file_size - 1)
                if f.read(1) == b"\n":
                    return

                # 마지막 개행 문자 위치를 찾아 손상된 꼬리 라인을 제거
                window_size = min(file_size, 8192)
                f.seek(file_size - window_size)
                tail = f.read(window_size)
                last_newline = tail.rfind(b"\n")

                if last_newline == -1:
                    f.seek(0)
                    f.truncate(0)
                else:
                    truncate_pos = file_size - window_size + last_newline + 1
                    f.seek(truncate_pos)
                    f.truncate()

            console.print("[yellow]⚠️ 증분 파일의 마지막 불완전 라인을 복구했습니다.[/yellow]")
        except Exception as repair_error:
            console.print(f"[yellow]⚠️ 증분 파일 복구 실패: {repair_error}[/yellow]")

    def save_incremental(self, source_key: str) -> None:
        """새로 수집된 페이지를 JSONL 파일에 증분 저장"""
        new_pages = self.all_pages[self._incremental_saved_count:]
        if not new_pages:
            return

        jsonl_path = self._get_incremental_path(source_key)
        payload_lines = [
            json.dumps(self._page_to_dict(page), ensure_ascii=False) + "\n"
            for page in new_pages
        ]

        try:
            with open(jsonl_path, "a", encoding="utf-8") as f:
                f.writelines(payload_lines)
                f.flush()
                os.fsync(f.fileno())
        except Exception:
            self._truncate_partial_jsonl_tail(jsonl_path)
            raise

        # Free memory: all pages have been flushed to JSONL on disk.
        # _total_pages_crawled tracks the true count across flushes.
        self.all_pages.clear()
        self._incremental_saved_count = 0
        self.write_runtime_stats()
        console.print(f"[dim]💾 증분 저장: {len(new_pages)}페이지 추가 (총 {self._total_pages_crawled}페이지, 메모리 해제)[/dim]")

        # Mid-crawl S3 backup: 500페이지마다 또는 JSONL > 500MB
        try:
            jsonl_size_mb = jsonl_path.stat().st_size / (1024 * 1024)
            if self._total_pages_crawled % 500 == 0 or jsonl_size_mb > 500:
                self._upload_jsonl_s3(jsonl_path, source_key)
        except Exception as e:
            console.print(f"[dim]⚠️  S3 mid-crawl backup skipped: {e}[/dim]")

    def _upload_jsonl_s3(self, jsonl_path: Path, source_key: str) -> None:
        """Upload JSONL to S3 for mid-crawl backup."""
        try:
            import boto3
        except ImportError:
            return

        kb_id = os.getenv("KB_ID", source_key)
        bucket = os.getenv("S3_UPLOAD_BUCKET", "gs-retail-svc-dev-miso-files")
        s3_key = f"knowledge/backup/{kb_id}/incremental_{source_key}.jsonl"
        file_size_mb = jsonl_path.stat().st_size / (1024 * 1024)

        console.print(f"[dim]☁️  S3 백업 시작: {s3_key} ({file_size_mb:.0f}MB)...[/dim]")
        try:
            s3 = boto3.client("s3")
            s3.upload_file(str(jsonl_path), bucket, s3_key)
            console.print(f"[dim]☁️  S3 백업 완료: s3://{bucket}/{s3_key}[/dim]")
        except Exception as e:
            console.print(f"[yellow]⚠️  S3 백업 실패: {e}[/yellow]")

    def load_incremental(self, source_key: str) -> int:
        """JSONL 파일에서 이전에 저장된 페이지 로드.

        Only pages with non-empty content_text are marked as visited.
        Pages with empty content (e.g., index/TOC pages, failed extractions)
        are left unvisited so they get re-crawled on resume.
        """
        jsonl_path = self._get_incremental_path(source_key)
        if not jsonl_path.exists():
            return 0

        loaded = 0
        skipped_empty = 0
        try:
            with open(jsonl_path, "r", encoding="utf-8") as f:
                for line in f:
                    line = line.strip()
                    if not line:
                        continue
                    page_data = json.loads(line)
                    page_id = page_data.get("page_id", "")
                    if not page_id:
                        continue
                    # Only mark as visited if content was actually extracted.
                    # Empty-body pages (TOC, restricted) should be re-attempted.
                    content_text = page_data.get("content_text", "")
                    if content_text and len(content_text.strip()) > 0:
                        self.visited_pages.add(page_id)
                        loaded += 1
                    else:
                        skipped_empty += 1
        except Exception as e:
            console.print(f"[yellow]⚠️ 증분 파일 로드 중 오류: {e}[/yellow]")

        if loaded > 0 or skipped_empty > 0:
            console.print(
                f"[green]♻️ 증분 파일: {loaded}페이지 방문 처리, {skipped_empty}페이지 재수집 대상[/green]"
            )
            self._incremental_saved_count = 0
        return loaded

    def clear_incremental(self, source_key: str) -> None:
        """증분 저장 파일 삭제"""
        jsonl_path = self._get_incremental_path(source_key)
        if jsonl_path.exists():
            jsonl_path.unlink()

    def finalize_from_incremental(self, source_key: str) -> list[dict]:
        """증분 JSONL + 메모리의 all_pages를 합쳐서 최종 JSON 생성"""
        all_page_dicts: list[dict] = []
        seen_ids: set[str] = set()

        # 1. JSONL 파일에서 이전 페이지 로드
        jsonl_path = self._get_incremental_path(source_key)
        if jsonl_path.exists():
            with open(jsonl_path, "r", encoding="utf-8") as f:
                for line in f:
                    line = line.strip()
                    if not line:
                        continue
                    page_data = json.loads(line)
                    pid = page_data.get("page_id", "")
                    if pid and pid not in seen_ids:
                        all_page_dicts.append(page_data)
                        seen_ids.add(pid)

        # 2. 메모리의 all_pages 중 아직 저장 안 된 것 추가
        for p in self.all_pages:
            if p.page_id not in seen_ids:
                all_page_dicts.append(self._page_to_dict(p))
                seen_ids.add(p.page_id)

        return all_page_dicts

    @staticmethod
    def _page_to_dict(p: "FullPageContent") -> dict:
        """FullPageContent → dict 변환 (save_results와 동일 형식)"""
        return {
            "page_id": p.page_id,
            "title": p.title,
            "content_text": p.content_text,
            "content_html": p.content_html,
            "content_ir": p.content_ir,
            "content_preview": p.content_preview,
            "tables": [
                {"headers": t.headers, "rows": t.rows, "table_type": t.table_type, "row_count": len(t.rows)}
                for t in p.tables
            ],
            "code_blocks": p.code_blocks,
            "mentions": [
                {"user_id": m.user_id, "display_name": m.display_name, "email": m.email, "context": m.context}
                for m in p.mentions
            ],
            "sections": p.sections,
            "emails": [
                {"email": e.email, "display_name": e.display_name, "context": e.context}
                for e in p.emails
            ],
            "macros": [
                {"macro_type": m.macro_type, "title": m.title, "content": m.content, "parameters": m.parameters}
                for m in p.macros
            ],
            "labels": [{"name": l.name, "prefix": l.prefix} for l in p.labels],
            "comments": [
                {"comment_id": c.comment_id, "author": c.author, "author_email": c.author_email,
                 "content": c.content, "created_at": c.created_at, "parent_id": c.parent_id}
                for c in p.comments
            ],
            "creator": p.creator,
            "creator_name": p.creator_name,
            "creator_team": p.creator_team,
            "creator_email": p.creator_email,
            "last_modifier": p.last_modifier,
            "version": p.version,
            "url": p.url,
            "created_at": p.created_at,
            "updated_at": p.updated_at,
            "space_key": p.space_key,
            "ancestors": p.ancestors,
            "internal_links": [
                {"target_page_id": lnk.target_page_id, "target_url": lnk.target_url,
                 "anchor_text": lnk.anchor_text, "context": lnk.context}
                for lnk in p.internal_links
            ],
            "external_links": [
                {"target_url": lnk.target_url, "anchor_text": lnk.anchor_text, "context": lnk.context}
                for lnk in p.external_links
            ],
            "restrictions": [
                {"operation": r.operation, "restriction_type": r.restriction_type,
                 "name": r.name, "account_id": r.account_id}
                for r in p.restrictions
            ],
            "version_history": p.version_history,
            "attachments": [
                {"id": a.id, "filename": a.filename, "media_type": a.media_type, "file_size": a.file_size,
                 "download_path": a.download_path, "download_url": a.download_url,
                 "extracted_text": a.extracted_text,
                 "extracted_tables": a.extracted_tables, "ocr_confidence": a.ocr_confidence,
                 "parse_error": a.parse_error,
                 "has_visual_content": a.has_visual_content,
                 "visual_analysis_version": a.visual_analysis_version,
                 "ocr_mode": a.ocr_mode,
                 "ocr_applied": a.ocr_applied,
                 "ocr_skip_reason": a.ocr_skip_reason,
                 "ocr_units_attempted": a.ocr_units_attempted,
                 "ocr_units_extracted": a.ocr_units_extracted,
                 "ocr_units_deferred": a.ocr_units_deferred,
                 "native_text_chars": a.native_text_chars,
                 "ocr_text_chars": a.ocr_text_chars}
                for a in p.attachments
            ],
        }

    async def close(self):
        await self.client.aclose()

    async def _http_get_with_retry(
        self, url: str, params: dict | None = None, max_retries: int = 3,
    ) -> httpx.Response:
        """HTTP GET with exponential backoff for transient errors."""
        last_error: Exception | None = None
        for attempt in range(max_retries):
            if self.shutdown_requested:
                raise RuntimeError("Shutdown requested during HTTP retry")
            try:
                response = await self.client.get(url, params=params)
                response.raise_for_status()
                return response
            except (httpx.TimeoutException, httpx.ConnectError, httpx.PoolTimeout) as e:
                last_error = e
                if attempt < max_retries - 1:
                    wait = min(2 ** attempt * 2, 30)
                    console.print(
                        f"[yellow]⚠️ HTTP retry {attempt + 1}/{max_retries} "
                        f"({type(e).__name__}): {url[:80]}[/yellow]"
                    )
                    await asyncio.sleep(wait)
            except httpx.HTTPStatusError as e:
                if e.response.status_code in (429, 500, 502, 503, 504):
                    last_error = e
                    if attempt < max_retries - 1:
                        wait = min(2 ** attempt * 5, 60)
                        console.print(
                            f"[yellow]⚠️ HTTP {e.response.status_code} retry "
                            f"{attempt + 1}/{max_retries}: {url[:80]}[/yellow]"
                        )
                        await asyncio.sleep(wait)
                else:
                    raise
        if last_error is not None:
            raise last_error
        raise RuntimeError("Unreachable: HTTP retry exhausted without error")

    async def get_user_details(self, account_id: str) -> dict | None:
        """사용자 상세 정보 조회 (이메일 등)"""
        if not account_id:
            return None

        url = f"{self.base_url}/rest/api/user"
        params = {"accountId": account_id}

        try:
            response = await self.client.get(url, params=params)
            response.raise_for_status()
            data = response.json()
            return {
                "account_id": account_id,
                "display_name": data.get("displayName"),
                "email": data.get("email"),
                "profile_picture": data.get("profilePicture", {}).get("path"),
            }
        except Exception:
            # 사용자 조회 실패는 무시 (권한 문제 등)
            return None

    async def get_comments(self, page_id: str) -> list[ExtractedComment]:
        """페이지 댓글 조회"""
        url = f"{self.base_url}/rest/api/content/{page_id}/child/comment"
        params = {
            "expand": "body.storage,history.createdBy",
            "limit": 100,
        }
        comments = []

        try:
            response = await self.client.get(url, params=params)
            response.raise_for_status()
            data = response.json()

            for comment in data.get("results", []):
                comment_id = comment.get("id", "")
                history = comment.get("history", {})
                created_by = history.get("createdBy", {})

                # 댓글 본문 추출
                body_html = comment.get("body", {}).get("storage", {}).get("value", "")
                text_extractor = PlainTextExtractor()
                text_extractor.feed(body_html)
                content = text_extractor.get_text()

                comments.append(ExtractedComment(
                    comment_id=comment_id,
                    author=created_by.get("displayName", "Unknown"),
                    author_email=created_by.get("email"),
                    content=content,
                    created_at=history.get("createdDate", ""),
                    parent_id=comment.get("ancestors", [{}])[0].get("id") if comment.get("ancestors") else None,
                ))

        except Exception as e:
            console.print(f"[yellow]Warning: Could not fetch comments of {page_id}: {e}[/yellow]")

        return comments

    async def get_labels(self, page_id: str) -> list[ExtractedLabel]:
        """페이지 라벨(태그) 조회"""
        url = f"{self.base_url}/rest/api/content/{page_id}/label"
        labels = []

        try:
            response = await self.client.get(url)
            response.raise_for_status()
            data = response.json()

            for label in data.get("results", []):
                labels.append(ExtractedLabel(
                    name=label.get("name", ""),
                    prefix=label.get("prefix"),
                ))

        except Exception as e:
            console.print(f"[yellow]Warning: Could not fetch labels of {page_id}: {e}[/yellow]")

        return labels

    async def get_page_full(self, page_id: str) -> FullPageContent | None:
        """페이지 전체 내용 조회"""
        url = f"{self.base_url}/rest/api/content/{page_id}"
        _default_expand = (
            "body.storage,version,space,ancestors,"
            "history.createdBy,history.lastUpdated,metadata.labels,"
            "restrictions.read.restrictions.user,"
            "restrictions.read.restrictions.group,"
            "restrictions.update.restrictions.user,"
            "restrictions.update.restrictions.group"
        )
        params = {"expand": os.getenv("CONFLUENCE_CRAWL_EXPAND", _default_expand)}

        try:
            response = await self._http_get_with_retry(url, params=params)
            data = response.json()

            title = data.get("title", "Unknown")
            body_html = data.get("body", {}).get("storage", {}).get("value", "")

            # Plain text 추출
            text_extractor = PlainTextExtractor()
            text_extractor.feed(body_html)
            content_text = text_extractor.get_text()

            # 테이블 추출
            table_extractor = TableExtractor()
            table_extractor.feed(body_html)
            tables = table_extractor.tables

            # 멘션 추출
            mention_extractor = MentionExtractor()
            mention_extractor.feed(body_html)
            mentions = mention_extractor.mentions

            # 섹션 추출
            section_extractor = SectionExtractor()
            section_extractor.feed(body_html)
            sections = section_extractor.sections

            # 코드 블록 추출
            code_extractor = CodeBlockExtractor()
            try:
                code_extractor.feed(body_html)
            except Exception:
                pass  # HTML 파싱 오류 무시
            code_blocks = code_extractor.code_blocks

            # NEW: 이메일 링크 추출
            email_extractor = EmailExtractor()
            try:
                email_extractor.feed(body_html)
            except Exception:
                pass
            emails = email_extractor.emails

            # NEW: 매크로 추출 (expand, panel, note, info, warning, status 등)
            macro_extractor = MacroExtractor()
            try:
                macro_extractor.feed(body_html)
            except Exception:
                pass
            macros = macro_extractor.macros

            # Structured IR 생성 (RAG 최적화)
            content_ir = generate_structured_ir(
                content_text=content_text,
                content_html=body_html,
                title=title,
                tables=tables,
                sections=sections,
                mentions=mentions,
            )

            # 메타데이터
            history = data.get("history", {})
            created_by_data = history.get("createdBy", {})
            creator = created_by_data.get("displayName", "Unknown")
            creator_account_id = created_by_data.get("accountId")
            creator_name, creator_team = extract_creator_info(creator)
            created_at = history.get("createdDate", "")

            # NEW: 작성자 이메일 조회
            creator_email = None
            if creator_account_id:
                user_details = await self.get_user_details(creator_account_id)
                if user_details:
                    creator_email = user_details.get("email")

            last_updated = history.get("lastUpdated", {})
            last_modifier = last_updated.get("by", {}).get("displayName", creator)
            updated_at = last_updated.get("when", created_at)

            version = data.get("version", {}).get("number", 1)
            page_url = f"{self.base_url}/pages/viewpage.action?pageId={page_id}"

            # NEW: Space 정보
            space_key = data.get("space", {}).get("key")

            # NEW: Ancestors (상위 페이지 계층)
            ancestors = [
                {"id": a.get("id"), "title": a.get("title")}
                for a in data.get("ancestors", [])
            ]

            # NEW: Labels (태그) - API에서 직접 조회
            labels = await self.get_labels(page_id)

            # NEW: Comments (댓글) 조회
            comments = await self.get_comments(page_id)

            # NEW: 내부/외부 링크 추출
            link_extractor = LinkExtractor(base_url=self.base_url)
            try:
                link_extractor.feed(body_html)
            except Exception:
                pass
            internal_links = link_extractor.internal_links
            external_links = link_extractor.external_links

            # NEW: Restrictions (접근 제한) 파싱
            restrictions = []
            restrictions_data = data.get("restrictions", {})

            # Read 권한
            read_restrictions = restrictions_data.get("read", {}).get("restrictions", {})
            for user in read_restrictions.get("user", {}).get("results", []):
                restrictions.append(ExtractedRestriction(
                    operation="read",
                    restriction_type="user",
                    name=user.get("displayName", ""),
                    account_id=user.get("accountId"),
                ))
            for group in read_restrictions.get("group", {}).get("results", []):
                restrictions.append(ExtractedRestriction(
                    operation="read",
                    restriction_type="group",
                    name=group.get("name", ""),
                ))

            # Update 권한
            update_restrictions = restrictions_data.get("update", {}).get("restrictions", {})
            for user in update_restrictions.get("user", {}).get("results", []):
                restrictions.append(ExtractedRestriction(
                    operation="update",
                    restriction_type="user",
                    name=user.get("displayName", ""),
                    account_id=user.get("accountId"),
                ))
            for group in update_restrictions.get("group", {}).get("results", []):
                restrictions.append(ExtractedRestriction(
                    operation="update",
                    restriction_type="group",
                    name=group.get("name", ""),
                ))

            # NEW: 버전 이력 (추가 API 호출 없이 현재 버전 스냅샷만 기록)
            version_data = data.get("version", {})
            current_version_snapshot = {
                "number": version_data.get("number", version),
                "when": version_data.get("when", updated_at),
                "by": version_data.get("by", {}).get("displayName", last_modifier),
                "message": version_data.get("message", ""),
            }
            version_history = [current_version_snapshot]

            # NEW: 멘션에 이메일 정보 보강
            for mention in mentions:
                if mention.user_id:
                    user_details = await self.get_user_details(mention.user_id)
                    if user_details:
                        mention.email = user_details.get("email")
                        if not mention.display_name:
                            mention.display_name = user_details.get("display_name")

            return FullPageContent(
                page_id=page_id,
                title=title,
                content_text=content_text,
                content_html=body_html,
                content_preview=content_text[:200] + "..." if len(content_text) > 200 else content_text,
                content_ir=content_ir,
                tables=tables,
                mentions=mentions,
                sections=sections,
                code_blocks=code_blocks,
                creator=creator,
                creator_name=creator_name,
                creator_team=creator_team,
                creator_email=creator_email,  # NEW
                last_modifier=last_modifier,
                version=version,
                url=page_url,
                created_at=created_at,
                updated_at=updated_at,
                labels=labels,  # NEW
                comments=comments,  # NEW
                emails=emails,  # NEW
                macros=macros,  # NEW
                space_key=space_key,  # NEW
                ancestors=ancestors,  # NEW
                internal_links=internal_links,  # NEW
                external_links=external_links,  # NEW
                restrictions=restrictions,  # NEW
                version_history=version_history,  # NEW
            )

        except httpx.TimeoutException as e:
            console.print(f"[red]Page {page_id} TIMEOUT ({type(e).__name__})[/red]")
            return None
        except httpx.HTTPStatusError as e:
            status = e.response.status_code
            body = e.response.text[:200]
            console.print(f"[red]Page {page_id} HTTP {status}: {body}[/red]")
            return None
        except Exception as e:
            console.print(f"[red]Page {page_id} ERROR ({type(e).__name__}): {e}[/red]")
            return None

    async def get_attachments(self, page_id: str) -> list[dict]:
        """첨부파일 목록 조회"""
        url = f"{self.base_url}/rest/api/content/{page_id}/child/attachment"
        params = {"limit": 100}

        try:
            response = await self._http_get_with_retry(url, params=params)
            return response.json().get("results", [])
        except Exception as e:
            console.print(f"[yellow]Warning: Could not fetch attachments of {page_id}: {e}[/yellow]")
            return []

    async def download_attachment(self, attachment: dict, page_id: str) -> AttachmentContent:
        """첨부파일 다운로드 및 내용 추출 (완전 구현)"""
        att_id = attachment.get("id", "")
        filename = attachment.get("title", "")
        media_type = attachment.get("extensions", {}).get("mediaType", "unknown")
        file_size = attachment.get("extensions", {}).get("fileSize", 0)
        download_url = f"{self.base_url}{attachment.get('_links', {}).get('download', '')}"

        # 이미지 포함 파일 유형 감지
        has_visual = any(
            filename.lower().endswith(ext)
            for ext in (".pptx", ".ppt", ".pdf", ".png", ".jpg", ".jpeg", ".gif", ".bmp")
        )

        result = AttachmentContent(
            id=att_id,
            filename=filename,
            media_type=media_type,
            file_size=file_size,
            download_url=download_url,
            has_visual_content=has_visual,
            ocr_mode=AttachmentParser.current_policy().attachment_ocr_mode if has_visual else None,
        )

        # 파일 크기 제한 (50MB)
        if file_size > 50 * 1024 * 1024:
            result.parse_error = f"파일 크기 초과 ({file_size / 1024 / 1024:.1f}MB > 50MB)"
            self._record_attachment_stats(result)
            return result

        try:
            response = await self.client.get(download_url)
            response.raise_for_status()
            content = response.content
            status_emit = lambda message: print(f"[status] {message}")

            # 파일 저장
            safe_filename = re.sub(r'[^\w\-_\. ]', '_', filename)
            file_path = ATTACHMENTS_DIR / f"{page_id}_{safe_filename}"
            with open(file_path, "wb") as f:
                f.write(content)
            result.download_path = str(file_path)

            # 내용 추출 (타입별) - 완전 구현
            media_lower = media_type.lower()
            filename_lower = filename.lower()
            parse_result: AttachmentParseResult | None = None

            if "pdf" in media_lower or filename_lower.endswith(".pdf"):
                # PDF 파싱
                parse_result = AttachmentParser.parse_pdf(file_path, heartbeat_fn=status_emit)

            elif any(x in media_lower for x in ["spreadsheet", "excel", "xlsx", "xls"]) or \
                 any(filename_lower.endswith(ext) for ext in [".xlsx", ".xls", ".xlsm"]):
                # Excel 파싱
                parse_result = AttachmentParser.parse_excel(file_path)

            # PPT 분기를 Word보다 먼저 체크 (미디어 타입에 "document"가 포함되기 때문)
            elif any(x in media_lower for x in ["presentation", "powerpoint", "pptx", "ppt"]) or \
                 any(filename_lower.endswith(ext) for ext in [".pptx", ".ppt"]):
                # PPT 파싱
                parse_result = AttachmentParser.parse_ppt(file_path, heartbeat_fn=status_emit)

            elif any(x in media_lower for x in ["word", "docx", "doc"]) or \
                 any(filename_lower.endswith(ext) for ext in [".docx", ".doc"]):
                # Word 파싱 (주의: "document" 제외 - PPT 미디어 타입과 충돌 방지)
                parse_result = AttachmentParser.parse_word(file_path)

            elif "image" in media_lower or \
                 any(filename_lower.endswith(ext) for ext in [".png", ".jpg", ".jpeg", ".gif", ".bmp"]):
                # 이미지 파싱 (OCR 포함, async로 CPU-bound 오프로드)
                parse_result = await AttachmentParser.parse_image_async(file_path, content)

            elif filename_lower.endswith(".txt") or "text" in media_lower:
                # 텍스트 파일
                try:
                    text = content.decode("utf-8")
                    parse_result = AttachmentParseResult(
                        extracted_text=text,
                        extracted_tables=[],
                        confidence=1.0,
                        native_text_chars=AttachmentParser._text_chars(text),
                    )
                except UnicodeDecodeError:
                    text = content.decode("cp949", errors="replace")
                    parse_result = AttachmentParseResult(
                        extracted_text=text,
                        extracted_tables=[],
                        confidence=0.8,
                        native_text_chars=AttachmentParser._text_chars(text),
                    )

            elif filename_lower.endswith(".csv"):
                # CSV 파일
                try:
                    text = content.decode("utf-8")
                    parse_result = AttachmentParseResult(
                        extracted_text=text,
                        extracted_tables=[],
                        confidence=0.95,
                        native_text_chars=AttachmentParser._text_chars(text),
                    )
                except UnicodeDecodeError:
                    text = content.decode("cp949", errors="replace")
                    parse_result = AttachmentParseResult(
                        extracted_text=text,
                        extracted_tables=[],
                        confidence=0.8,
                        native_text_chars=AttachmentParser._text_chars(text),
                    )

            else:
                # 지원하지 않는 형식
                parse_result = AttachmentParseResult(
                    extracted_text=f"[지원하지 않는 형식: {media_type}]",
                    extracted_tables=[],
                    confidence=0.0,
                    ocr_skip_reason="unsupported_media_type",
                )

            if parse_result is not None:
                result.extracted_text = parse_result.extracted_text
                result.extracted_tables = parse_result.extracted_tables
                result.ocr_confidence = parse_result.confidence
                result.ocr_mode = parse_result.ocr_mode
                result.ocr_applied = parse_result.ocr_applied
                result.ocr_skip_reason = parse_result.ocr_skip_reason
                result.ocr_units_attempted = parse_result.ocr_units_attempted
                result.ocr_units_extracted = parse_result.ocr_units_extracted
                result.ocr_units_deferred = parse_result.ocr_units_deferred
                result.native_text_chars = parse_result.native_text_chars
                result.ocr_text_chars = parse_result.ocr_text_chars

            # 파서가 None 반환 시 빈 문자열로 fallback (인제스천 누락 방지)
            if result.extracted_text is None:
                result.extracted_text = ""

        except Exception as e:
            result.parse_error = str(e)

        self._record_attachment_stats(result)
        return result

    async def get_child_pages(self, page_id: str) -> list[str]:
        """하위 페이지 ID 목록"""
        url = f"{self.base_url}/rest/api/content/{page_id}/child/page"
        params = {"limit": 100}
        child_ids = []
        visited_urls: set[str] = set()  # 무한 루프 방지

        try:
            while url:
                # 이미 방문한 URL이면 중단 (무한 루프 방지)
                if url in visited_urls:
                    break
                visited_urls.add(url)

                response = await self._http_get_with_retry(url, params=params)
                data = response.json()

                results = data.get("results", [])
                # 결과가 비어있으면 중단
                if not results:
                    break

                for result in results:
                    child_ids.append(result.get("id"))

                links = data.get("_links", {})
                next_link = links.get("next")
                if next_link:
                    url = f"{self.base_url}{next_link}"
                    params = {}
                else:
                    url = None
        except Exception as e:
            console.print(f"[yellow]Warning: Could not fetch children of {page_id}: {e}[/yellow]")

        return child_ids

    async def get_all_descendant_page_ids_via_cql(self, root_page_id: str) -> set[str]:
        """CQL을 사용하여 root 아래 모든 descendant 페이지 ID를 수집."""
        url = f"{self.base_url}/rest/api/content/search"
        cql = f"ancestor={root_page_id} and type=page"
        start = 0
        limit = 100
        all_ids: set[str] = set()

        while True:
            if self.shutdown_requested:
                break
            params = {"cql": cql, "limit": limit, "start": start}
            try:
                resp = await self._http_get_with_retry(url, params=params)
                data = resp.json()
            except Exception as e:
                console.print(f"[yellow]CQL 검색 오류 (start={start}): {e}[/yellow]")
                break

            results = data.get("results", [])
            if not results:
                break

            for r in results:
                pid = r.get("id")
                if pid:
                    all_ids.add(str(pid))

            total_size = data.get("totalSize", 0)
            start += limit

            if start % 1000 == 0 or not results:
                console.print(f"[dim]CQL 열거: {len(all_ids)}/{total_size} 페이지[/dim]")

            if start >= total_size:
                break

        console.print(f"[green]CQL 열거 완료: {len(all_ids)} 페이지[/green]")
        return all_ids

    async def crawl_recursive(
        self,
        page_id: str,
        depth: int = 0,
        max_depth: int = 10,
        max_pages: int | None = None,
        download_attachments: bool = True,
        max_attachments_per_page: int = 20,
        progress: Progress | None = None,
        task_id: Any = None,
        source_key: str = "unknown",
    ) -> FullPageContent | None:
        """재귀적 크롤링 (병렬 지원)"""

        if self.shutdown_requested:
            return None

        # 이미 방문한 페이지: 콘텐츠 재수집 스킵, 하위 페이지는 탐색 (resume 지원)
        if page_id in self.visited_pages:
            if depth <= max_depth:
                try:
                    child_ids = await self.get_child_pages(page_id)
                    if child_ids:
                        await self._crawl_children(
                            child_ids, depth, max_depth, max_pages,
                            download_attachments, max_attachments_per_page,
                            progress, task_id, source_key,
                        )
                except Exception:
                    pass
            return None
        self.visited_pages.add(page_id)

        # 최대 페이지 수 제한
        if max_pages and self._total_pages_crawled >= max_pages:
            return None

        if depth > max_depth:
            return None

        # === Phase 1: 페이지 처리 (세마포어 보호, HTTP 집약 구간) ===
        page, child_ids = await self._process_single_page(
            page_id, download_attachments, max_attachments_per_page,
            progress, task_id, source_key,
        )

        # === Phase 2: 하위 페이지 크롤링 (세마포어 밖, 병렬 또는 순차) ===
        await self._crawl_children(
            child_ids, depth, max_depth, max_pages,
            download_attachments, max_attachments_per_page,
            progress, task_id, source_key,
        )

        return page

    async def _process_single_page(
        self,
        page_id: str,
        download_attachments: bool,
        max_attachments_per_page: int,
        progress: Progress | None,
        task_id: Any,
        source_key: str,
        skip_children: bool = False,
    ) -> tuple[FullPageContent | None, list[str]]:
        """단일 페이지 처리: 콘텐츠 조회 + 첨부파일 + 자식 ID 목록.

        세마포어가 설정된 경우 동시 HTTP 호출 수를 제한합니다.
        """
        if self._page_sem:
            async with self._page_sem:
                return await self._do_process_page(
                    page_id, download_attachments, max_attachments_per_page,
                    progress, task_id, source_key, skip_children=skip_children,
                )
        return await self._do_process_page(
            page_id, download_attachments, max_attachments_per_page,
            progress, task_id, source_key, skip_children=skip_children,
        )

    async def _do_process_page(
        self,
        page_id: str,
        download_attachments: bool,
        max_attachments_per_page: int,
        progress: Progress | None,
        task_id: Any,
        source_key: str,
        skip_children: bool = False,
    ) -> tuple[FullPageContent | None, list[str]]:
        """실제 페이지 처리 로직 (세마포어 내부에서 호출)."""
        page = await self.get_page_full(page_id)
        if not page:
            console.print(
                f"[yellow]페이지 {page_id} 조회 실패 (상세: 위 [red] 로그), 하위 페이지 탐색 계속[/yellow]"
            )
            child_ids = await self.get_child_pages(page_id)
            return None, child_ids

        # Per-page content validation
        text_len = len(page.content_text) if page.content_text else 0
        html_len = len(page.content_html) if page.content_html else 0
        if text_len == 0:
            if html_len == 0:
                console.print(
                    f"[yellow]WARNING: page {page_id} ({page.title[:50]}) "
                    f"has empty body (html=0, text=0) -- possible permission issue[/yellow]"
                )
            else:
                console.print(
                    f"[yellow]WARNING: page {page_id} ({page.title[:50]}) "
                    f"has HTML ({html_len} chars) but extracted text is empty[/yellow]"
                )

        self.all_pages.append(page)
        self._total_pages_crawled += 1

        # 주기적 체크포인트 + 증분 저장
        self._pages_since_checkpoint += 1
        if self._pages_since_checkpoint >= self.CHECKPOINT_INTERVAL:
            self.save_checkpoint(source_key)
            self.save_incremental(source_key)
            self._pages_since_checkpoint = 0

        if progress and task_id:
            progress.update(
                task_id,
                description=f"[cyan]({len(self.all_pages)}) {page.title[:40]}...[/cyan]"
            )

        # 첨부파일 처리 (shutdown 체크 포함, 병렬 다운로드+OCR)
        if download_attachments and not self.shutdown_requested:
            import gc
            attachments_meta = await self.get_attachments(page_id)
            target_attachments = attachments_meta[:max_attachments_per_page]

            if target_attachments:
                # 병렬 다운로드 (max 2 concurrent — OMP=2 × sem=2 = 4코어, OOM 방지)
                _att_sem = asyncio.Semaphore(2)

                async def _dl_one(meta: dict) -> AttachmentContent | None:
                    if self.shutdown_requested:
                        return None
                    async with _att_sem:
                        return await self.download_attachment(meta, page_id)

                results = await asyncio.gather(
                    *[_dl_one(m) for m in target_attachments],
                    return_exceptions=True,
                )
                for i, r in enumerate(results):
                    if isinstance(r, (Exception, BaseException)):
                        att_name = target_attachments[i].get("title", "unknown")
                        print(f"[Warning] 첨부파일 다운로드 실패 ({att_name}): {r}")
                page.attachments = [
                    r for r in results
                    if r is not None and not isinstance(r, (Exception, BaseException))
                ]

                image_count = sum(
                    1 for m in target_attachments
                    if "image" in m.get("extensions", {}).get("mediaType", "").lower()
                )
                if image_count > 0:
                    gc.collect()

        # 하위 페이지 ID 조회 (flat 모드에서는 스킵)
        child_ids = [] if skip_children else await self.get_child_pages(page_id)
        return page, child_ids

    async def _crawl_children(
        self,
        child_ids: list[str],
        depth: int,
        max_depth: int,
        max_pages: int | None,
        download_attachments: bool,
        max_attachments_per_page: int,
        progress: Progress | None,
        task_id: Any,
        source_key: str,
    ) -> None:
        """하위 페이지 크롤링 (병렬 또는 순차).

        _max_concurrent > 1 이면 asyncio.gather로 병렬 처리.
        세마포어(_page_sem)가 전체 동시 HTTP 호출 수를 제한하므로
        재귀적 gather가 중첩되어도 Confluence 서버 부하가 통제됩니다.
        """
        if not child_ids or self.shutdown_requested:
            return

        if self._max_concurrent > 1:
            # === 병렬 모드 ===
            tasks: list[asyncio.Task[FullPageContent | None]] = []
            for child_id in child_ids:
                if self.shutdown_requested:
                    break
                if max_pages and self._total_pages_crawled >= max_pages:
                    break
                tasks.append(
                    asyncio.ensure_future(
                        self.crawl_recursive(
                            child_id,
                            depth=depth + 1,
                            max_depth=max_depth,
                            max_pages=max_pages,
                            download_attachments=download_attachments,
                            max_attachments_per_page=max_attachments_per_page,
                            progress=progress,
                            task_id=task_id,
                            source_key=source_key,
                        )
                    )
                )

            if not tasks:
                return

            # gather + 인터럽트 안전: shutdown 시 대기 중인 태스크 취소
            try:
                results = await asyncio.gather(*tasks, return_exceptions=True)
            except asyncio.CancelledError:
                for t in tasks:
                    if not t.done():
                        t.cancel()
                raise

            for r in results:
                if isinstance(r, Exception) and not isinstance(r, asyncio.CancelledError):
                    console.print(f"[yellow]Child page crawl error: {r}[/yellow]")

        else:
            # === 순차 모드 (기존 동작) ===
            for child_id in child_ids:
                if self.shutdown_requested:
                    break
                if max_pages and self._total_pages_crawled >= max_pages:
                    break
                await self.crawl_recursive(
                    child_id,
                    depth=depth + 1,
                    max_depth=max_depth,
                    max_pages=max_pages,
                    download_attachments=download_attachments,
                    max_attachments_per_page=max_attachments_per_page,
                    progress=progress,
                    task_id=task_id,
                    source_key=source_key,
                )

    async def crawl_flat(
        self,
        page_ids: list[str],
        download_attachments: bool = True,
        max_attachments_per_page: int = 20,
        max_pages: int | None = None,
        progress: Progress | None = None,
        task_id: Any = None,
        source_key: str = "unknown",
    ) -> None:
        """미방문 페이지 목록을 flat하게 크롤링 (DFS 불필요)."""
        total = len(page_ids)
        for i, page_id in enumerate(page_ids):
            if self.shutdown_requested:
                break
            if max_pages and self._total_pages_crawled >= max_pages:
                break
            if page_id in self.visited_pages:
                continue

            self.visited_pages.add(page_id)

            await self._process_single_page(
                page_id, download_attachments, max_attachments_per_page,
                progress, task_id, source_key,
                skip_children=True,
            )

            if progress and task_id and (i + 1) % 50 == 0:
                progress.update(
                    task_id,
                    description=f"[cyan]flat 크롤링: {self._total_pages_crawled}/{total}[/cyan]",
                )


def save_results_from_jsonl(
    jsonl_path: Path,
    output_path: Path,
    source_info: dict | None = None,
) -> int:
    """JSONL을 스트리밍으로 JSON 변환 (메모리 효율적, 대규모 크롤링용).

    전체 페이지를 메모리에 로드하지 않고 JSONL → JSON을 line-by-line 스트리밍.
    통계도 스트리밍 중 계산하여 상수 메모리 사용.

    Returns:
        저장된 페이지 수
    """
    # 통계 카운터
    total_pages = 0
    total_tables = 0
    total_code_blocks = 0
    total_ir_chunks = 0
    total_attachments = 0
    parsed_stats = {"pdf": 0, "excel": 0, "word": 0, "ppt": 0, "image": 0, "other": 0, "failed": 0}
    total_extracted_text_length = 0
    total_labels = 0
    total_comments = 0
    total_emails = 0
    total_macros = 0
    total_internal_links = 0
    total_external_links = 0
    total_restrictions = 0
    total_version_history = 0
    seen_ids: set[str] = set()

    with open(output_path, "w", encoding="utf-8") as out_f:
        # JSON header (statistics placeholder — rewritten at end)
        # Write pages array directly for streaming
        out_f.write('{"crawled_at":"')
        out_f.write(datetime.now(timezone.utc).isoformat())
        out_f.write('","source_info":')
        out_f.write(json.dumps(source_info, ensure_ascii=False) if source_info else "null")
        out_f.write(',"pages":[')

        first = True
        with open(jsonl_path, "r", encoding="utf-8") as in_f:
            for line in in_f:
                line = line.strip()
                if not line:
                    continue
                try:
                    page = json.loads(line)
                except Exception:
                    continue
                if not isinstance(page, dict):
                    continue

                pid = page.get("page_id", "")
                if pid and pid in seen_ids:
                    continue
                if pid:
                    seen_ids.add(pid)

                # Write page
                if not first:
                    out_f.write(",")
                out_f.write(json.dumps(page, ensure_ascii=False))
                first = False
                total_pages += 1

                # Count page body text
                page_body = page.get("content_text", "")
                if page_body:
                    total_extracted_text_length += len(str(page_body))

                # Accumulate statistics
                total_tables += len(page.get("tables", []))
                total_code_blocks += len(page.get("code_blocks", []))
                total_ir_chunks += (page.get("content_ir") or {}).get("chunk_count", 0)
                total_labels += len(page.get("labels", []))
                total_comments += len(page.get("comments", []))
                total_emails += len(page.get("emails", []))
                total_macros += len(page.get("macros", []))
                total_internal_links += len(page.get("internal_links", []))
                total_external_links += len(page.get("external_links", []))
                total_restrictions += len(page.get("restrictions", []))
                total_version_history += len(page.get("version_history", []))

                for att in page.get("attachments", []):
                    total_attachments += 1
                    media_type = str(att.get("media_type", "")).lower()
                    if att.get("parse_error"):
                        parsed_stats["failed"] += 1
                    elif "pdf" in media_type:
                        parsed_stats["pdf"] += 1
                    elif "excel" in media_type or "spreadsheet" in media_type:
                        parsed_stats["excel"] += 1
                    elif "word" in media_type or "document" in media_type:
                        parsed_stats["word"] += 1
                    elif "presentation" in media_type or "powerpoint" in media_type:
                        parsed_stats["ppt"] += 1
                    elif "image" in media_type:
                        parsed_stats["image"] += 1
                    else:
                        parsed_stats["other"] += 1
                    extracted_text = att.get("extracted_text")
                    if extracted_text:
                        total_extracted_text_length += len(str(extracted_text))

        # Close pages array, write statistics, close root object
        out_f.write('],"statistics":')
        stats = {
            "total_pages": total_pages,
            "total_tables_in_pages": total_tables,
            "total_code_blocks": total_code_blocks,
            "total_ir_chunks": total_ir_chunks,
            "total_attachments": total_attachments,
            "attachment_parsing": parsed_stats,
            "total_extracted_text_chars": total_extracted_text_length,
            "total_labels": total_labels,
            "total_comments": total_comments,
            "total_emails": total_emails,
            "total_macros": total_macros,
            "total_internal_links": total_internal_links,
            "total_external_links": total_external_links,
            "total_restrictions": total_restrictions,
            "total_version_history": total_version_history,
        }
        out_f.write(json.dumps(stats, ensure_ascii=False))
        out_f.write("}")

    console.print(f"\n[green]✓ 스트리밍 저장 완료: {output_path} ({total_pages}페이지)[/green]")

    # 파싱 통계 출력
    console.print("\n[bold]📎 첨부파일 파싱 결과:[/bold]")
    for file_type, count in parsed_stats.items():
        if count > 0:
            icon = "✓" if file_type != "failed" else "✗"
            color = "green" if file_type != "failed" else "red"
            console.print(f"  [{color}]{icon}[/{color}] {file_type}: {count}개")
    console.print(f"  [cyan]총 추출 텍스트: {total_extracted_text_length:,}자[/cyan]")

    # Zero-content validation (streaming mode)
    if total_pages == 0:
        console.print(
            "\n[bold red]ERROR: Output has 0 pages (streaming). "
            "JSONL may be empty or corrupted. Try --fresh-full.[/bold red]"
        )
    elif total_extracted_text_length == 0:
        console.print(
            f"\n[bold red]ERROR: {total_pages} pages streamed but total extracted text is 0 chars. "
            f"Possible cause: Confluence returned empty body.storage for all pages.[/bold red]"
        )

    return total_pages


def save_results(
    pages: list[FullPageContent],
    output_path: Path,
    source_info: dict | None = None,
    page_dicts: list[dict] | None = None,
) -> None:
    """결과 저장 (첨부파일 테이블 포함)."""

    serialized_pages = page_dicts if page_dicts is not None else [ConfluenceFullClient._page_to_dict(page) for page in pages]

    # 첨부파일 파싱 통계
    parsed_stats = {"pdf": 0, "excel": 0, "word": 0, "ppt": 0, "image": 0, "other": 0, "failed": 0}
    total_extracted_text_length = 0

    for page in serialized_pages:
        # Count page body text
        page_body = page.get("content_text", "")
        if page_body:
            total_extracted_text_length += len(str(page_body))

        for attachment in page.get("attachments", []):
            media_type = str(attachment.get("media_type", "")).lower()
            if attachment.get("parse_error"):
                parsed_stats["failed"] += 1
            elif "pdf" in media_type:
                parsed_stats["pdf"] += 1
            elif "excel" in media_type or "spreadsheet" in media_type:
                parsed_stats["excel"] += 1
            elif "word" in media_type or "document" in media_type:
                parsed_stats["word"] += 1
            elif "presentation" in media_type or "powerpoint" in media_type:
                parsed_stats["ppt"] += 1
            elif "image" in media_type:
                parsed_stats["image"] += 1
            else:
                parsed_stats["other"] += 1

            extracted_text = attachment.get("extracted_text")
            if extracted_text:
                total_extracted_text_length += len(str(extracted_text))

    # NEW: 추가 통계
    total_labels = sum(len(page.get("labels", [])) for page in serialized_pages)
    total_comments = sum(len(page.get("comments", [])) for page in serialized_pages)
    total_emails = sum(len(page.get("emails", [])) for page in serialized_pages)
    total_macros = sum(len(page.get("macros", [])) for page in serialized_pages)
    total_internal_links = sum(len(page.get("internal_links", [])) for page in serialized_pages)
    total_external_links = sum(len(page.get("external_links", [])) for page in serialized_pages)
    total_restrictions = sum(len(page.get("restrictions", [])) for page in serialized_pages)
    total_version_history = sum(len(page.get("version_history", [])) for page in serialized_pages)

    result = {
        "crawled_at": datetime.now(timezone.utc).isoformat(),
        "source_info": source_info,
        "statistics": {
            "total_pages": len(serialized_pages),
            "total_tables_in_pages": sum(len(page.get("tables", [])) for page in serialized_pages),
            "total_code_blocks": sum(len(page.get("code_blocks", [])) for page in serialized_pages),
            "total_ir_chunks": sum(
                (page.get("content_ir") or {}).get("chunk_count", 0)
                for page in serialized_pages
            ),
            "total_attachments": sum(len(page.get("attachments", [])) for page in serialized_pages),
            "attachment_parsing": parsed_stats,
            "total_extracted_text_chars": total_extracted_text_length,
            # NEW: 추가 통계
            "total_labels": total_labels,
            "total_comments": total_comments,
            "total_emails": total_emails,
            "total_macros": total_macros,
            "total_internal_links": total_internal_links,
            "total_external_links": total_external_links,
            "total_restrictions": total_restrictions,
            "total_version_history": total_version_history,
        },
        "pages": serialized_pages,
    }

    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)

    console.print(f"\n[green]✓ 결과 저장: {output_path}[/green]")

    # 파싱 통계 출력
    console.print("\n[bold]📎 첨부파일 파싱 결과:[/bold]")
    for file_type, count in parsed_stats.items():
        if count > 0:
            icon = "✓" if file_type != "failed" else "✗"
            color = "green" if file_type != "failed" else "red"
            console.print(f"  [{color}]{icon}[/{color}] {file_type}: {count}개")
    console.print(f"  [cyan]총 추출 텍스트: {total_extracted_text_length:,}자[/cyan]")

    # Zero-content validation
    if len(serialized_pages) == 0:
        console.print(
            "\n[bold red]ERROR: Output has 0 pages. "
            "Possible causes: (1) --resume with stale checkpoint marking all pages as visited, "
            "(2) Confluence API permission issue, (3) all pages returned empty body. "
            "Try running with --fresh-full to start clean.[/bold red]"
        )
    elif total_extracted_text_length == 0:
        pages_with_body = sum(1 for p in serialized_pages if p.get("content_text"))
        console.print(
            f"\n[bold red]ERROR: {len(serialized_pages)} pages collected but total extracted text is 0 chars "
            f"(pages with body: {pages_with_body}). "
            f"Possible cause: Confluence returns metadata but empty body.storage "
            f"(PAT permission issue or view restriction on pages).[/bold red]"
        )

    # NEW: 추가 데이터 통계 출력
    has_extra = any([total_labels, total_comments, total_emails, total_macros,
                     total_internal_links, total_external_links, total_restrictions])
    if has_extra:
        console.print("\n[bold]📋 추가 수집 데이터:[/bold]")
        if total_labels > 0:
            console.print(f"  [green]✓[/green] 라벨(태그): {total_labels}개")
        if total_comments > 0:
            console.print(f"  [green]✓[/green] 댓글: {total_comments}개")
        if total_emails > 0:
            console.print(f"  [green]✓[/green] 이메일 링크: {total_emails}개")
        if total_macros > 0:
            console.print(f"  [green]✓[/green] 매크로: {total_macros}개")
        if total_internal_links > 0:
            console.print(f"  [green]✓[/green] 내부 링크: {total_internal_links}개")
        if total_external_links > 0:
            console.print(f"  [green]✓[/green] 외부 링크: {total_external_links}개")
        if total_restrictions > 0:
            console.print(f"  [yellow]🔒[/yellow] 접근 제한: {total_restrictions}개")


# =============================================================================
# Main
# =============================================================================
def show_available_sources():
    """사용 가능한 지식 소스 표시"""
    table = Table(title="📚 사용 가능한 지식 소스")
    table.add_column("Key", style="cyan")
    table.add_column("Name", style="green")
    table.add_column("Page ID", style="yellow")

    for key, info in KNOWLEDGE_SOURCES.items():
        table.add_row(key, info["name"], info["page_id"])

    console.print(table)


async def crawl_space(
    page_id: str,
    source_name: str,
    source_key: str,
    max_pages: int | None,
    download_attachments: bool,
    max_attachments_per_page: int = 20,
    resume: bool = False,
    max_concurrent: int = 1,
    kb_id: str = "",
) -> CrawlSpaceResult:
    """단일 스페이스 크롤링"""
    policy = AttachmentParser.configure_run(source_key)

    resume_mode = "♻️ 재개 모드" if resume else "🆕 새로 시작"

    console.print(Panel.fit(
        f"[bold blue]Confluence Full Content Crawler[/bold blue]\n\n"
        f"📁 Source: {source_name}\n"
        f"📄 Root Page ID: {page_id}\n"
        f"📊 Max Pages: {'무제한' if max_pages is None else max_pages}\n"
        f"📎 Attachments: {'✓ 다운로드 및 파싱' if download_attachments else '✗ 스킵'}\n"
        f"⚡ Concurrency: {max_concurrent}\n"
        f"🔄 Mode: {resume_mode}\n"
        f"🧠 OCR Policy: {policy.attachment_ocr_mode} "
        f"(pdf≤{policy.ocr_max_pdf_pages}, ppt≤{policy.ocr_max_ppt_slides}, slide_render={'on' if policy.slide_render_enabled else 'off'})",
        border_style="blue",
    ))

    client = ConfluenceFullClient(max_concurrent=max_concurrent, kb_id=kb_id)
    interrupted = False
    previous_sigterm = signal.getsignal(signal.SIGTERM)
    previous_sigint = signal.getsignal(signal.SIGINT)

    # SIGTERM/SIGINT 핸들러 등록 (안전한 종료)
    def _graceful_shutdown(signum, frame):
        nonlocal interrupted
        if interrupted:
            return

        interrupted = True
        client.request_shutdown()
        sig_name = "SIGTERM" if signum == signal.SIGTERM else "SIGINT"
        console.print(f"\n[yellow]⚠️ {sig_name} 수신 - 현재 작업 완료 후 안전 종료합니다...[/yellow]")

    signal.signal(signal.SIGTERM, _graceful_shutdown)
    signal.signal(signal.SIGINT, _graceful_shutdown)

    # 체크포인트에서 재개
    if resume:
        # 증분 JSONL에서 이전 페이지 복원
        loaded = client.load_incremental(source_key)
        if client.load_checkpoint(source_key):
            console.print("[green]✓ 이전 진행 상태에서 재개합니다.[/green]")
        elif loaded > 0:
            console.print(f"[green]✓ 증분 파일에서 {loaded}페이지 복원, 이어서 진행합니다.[/green]")
        else:
            console.print("[yellow]⚠️ 체크포인트가 없습니다. 처음부터 시작합니다.[/yellow]")
    else:
        # 새로 시작할 때 이전 체크포인트 + 증분 파일 삭제
        client.clear_checkpoint()
        client.clear_incremental(source_key)

    try:
        with Progress(
            SpinnerColumn(),
            TextColumn("[progress.description]{task.description}"),
            console=console,
        ) as progress:
            task = progress.add_task("[cyan]크롤링 시작...[/cyan]", total=None)

            if resume and client.visited_pages:
                # === DFS resume: visited_pages 기반으로 방문 페이지 스킵 + 미방문 하위 탐색 ===
                # NOTE: CQL ancestor 쿼리는 일부 페이지만 반환하는 Confluence 제한이 있어
                #       DFS crawl_recursive를 사용합니다. visited_pages에 이미 있는 페이지는
                #       콘텐츠 재수집을 스킵하고 하위 페이지만 탐색합니다.
                console.print(
                    f"[green]체크포인트에서 재개: {len(client.visited_pages)}페이지 방문 완료, 미방문 하위 페이지 탐색[/green]"
                )
                await client.crawl_recursive(
                    page_id,
                    max_depth=10,
                    max_pages=max_pages,
                    download_attachments=download_attachments,
                    max_attachments_per_page=max_attachments_per_page,
                    progress=progress,
                    task_id=task,
                    source_key=source_key,
                )
            else:
                # === 기존 DFS 크롤링 (새로 시작 또는 체크포인트 없음) ===
                await client.crawl_recursive(
                    page_id,
                    max_depth=10,
                    max_pages=max_pages,
                    download_attachments=download_attachments,
                    max_attachments_per_page=max_attachments_per_page,
                    progress=progress,
                    task_id=task,
                    source_key=source_key,
                )

            if client.shutdown_requested:
                progress.update(task, description="[yellow]중단 요청 처리됨[/yellow]")
            else:
                progress.update(task, description="[green]완료![/green]")

        # Resume guard: detect stale state where all pages were already visited
        if resume and client.visited_pages and len(client.all_pages) == 0:
            jsonl_path = client._get_incremental_path(source_key)
            jsonl_exists = jsonl_path.exists() and jsonl_path.stat().st_size > 0
            if not jsonl_exists:
                console.print(
                    "\n[bold red]WARNING: --resume produced 0 new pages and no incremental JSONL data exists. "
                    "This likely means a previous crawl completed fully and the resume checkpoint is stale. "
                    "Use --fresh-full to start a clean crawl.[/bold red]"
                )

        # 완료 전 마지막 증분 저장
        client.save_incremental(source_key)

        if client.shutdown_requested:
            client.save_checkpoint(source_key)
            page_dicts = client.finalize_from_incremental(source_key)
            client.write_runtime_stats()
            if download_attachments:
                AttachmentParser.cleanup_ocr()
            console.print("[green]✓ 안전하게 저장 완료. --resume으로 재개 가능합니다.[/green]")
            return CrawlSpaceResult(
                pages=client.all_pages,
                page_dicts=page_dicts,
                interrupted=True,
            )

        # NOTE: checkpoint 삭제는 main()에서 JSON save 완료 후 수행
        # (이전: 여기서 삭제 → finalize 중 OOM 시 checkpoint 유실 → resume 불가)

        # 결과 출력 (all_pages may be empty after incremental flush; use total counter)
        console.print(f"\n[green]✓ {client._total_pages_crawled}개 페이지 수집[/green]")

        total_tables = sum(len(p.tables) for p in client.all_pages)
        total_mentions = sum(len(p.mentions) for p in client.all_pages)
        total_attachments = sum(len(p.attachments) for p in client.all_pages)
        # NEW: 추가 통계
        total_labels = sum(len(p.labels) for p in client.all_pages)
        total_comments = sum(len(p.comments) for p in client.all_pages)
        total_emails = sum(len(p.emails) for p in client.all_pages)
        total_macros = sum(len(p.macros) for p in client.all_pages)

        console.print(f"[green]✓ {total_tables}개 테이블 추출 (HTML)[/green]")
        console.print(f"[green]✓ {total_mentions}개 @멘션 추출[/green]")
        console.print(f"[green]✓ {total_attachments}개 첨부파일 처리[/green]")
        # NEW: 추가 통계
        total_internal_links = sum(len(p.internal_links) for p in client.all_pages)
        total_external_links = sum(len(p.external_links) for p in client.all_pages)
        total_restrictions = sum(len(p.restrictions) for p in client.all_pages)

        # NEW: 추가 통계 출력
        if total_labels > 0:
            console.print(f"[green]✓ {total_labels}개 라벨(태그) 수집[/green]")
        if total_comments > 0:
            console.print(f"[green]✓ {total_comments}개 댓글 수집[/green]")
        if total_emails > 0:
            console.print(f"[green]✓ {total_emails}개 이메일 링크 추출[/green]")
        if total_macros > 0:
            console.print(f"[green]✓ {total_macros}개 매크로 추출[/green]")
        if total_internal_links > 0:
            console.print(f"[green]✓ {total_internal_links}개 내부 링크 추출[/green]")
        if total_external_links > 0:
            console.print(f"[green]✓ {total_external_links}개 외부 링크 추출[/green]")
        if total_restrictions > 0:
            console.print(f"[yellow]🔒 {total_restrictions}개 접근 제한 감지[/yellow]")

        # 테이블 유형별 통계
        table_types: dict[str, int] = {}
        owner_tables_count = 0
        for page in client.all_pages:
            for table in page.tables:
                tt = table.table_type or "general"
                table_types[tt] = table_types.get(tt, 0) + 1
                if tt == "owner_table":
                    owner_tables_count += 1

        if table_types:
            console.print("\n[bold]📊 테이블 유형:[/bold]")
            for tt, count in sorted(table_types.items(), key=lambda x: -x[1]):
                icon = "👤" if tt == "owner_table" else "📋"
                console.print(f"  {icon} {tt}: {count}개")

        if owner_tables_count > 0:
            console.print(f"\n[bold green]🎯 담당자 테이블 {owner_tables_count}개 발견![/bold green]")

        # OCR 메모리 정리 (첨부파일 처리 완료 후)
        if download_attachments:
            AttachmentParser.cleanup_ocr()

        # 스트리밍 모드 결정: JSONL > 500MB이면 finalize 스킵 (OOM 방지)
        jsonl_path = client._get_incremental_path(source_key)
        jsonl_size_mb = jsonl_path.stat().st_size / (1024 * 1024) if jsonl_path.exists() else 0

        if jsonl_size_mb > 500:
            console.print(f"[cyan]📦 스트리밍 모드: JSONL {jsonl_size_mb:.0f}MB > 500MB, finalize 스킵[/cyan]")
            page_dicts: list[dict] = []
            use_streaming = True
        else:
            page_dicts = client.finalize_from_incremental(source_key)
            use_streaming = False

        client.write_runtime_stats()

        return CrawlSpaceResult(
            pages=client.all_pages,
            page_dicts=page_dicts,
            interrupted=False,
            jsonl_path=str(jsonl_path) if use_streaming else "",
            source_key=source_key,
        )

    finally:
        try:
            client.write_runtime_stats()
        except Exception:
            pass
        signal.signal(signal.SIGTERM, previous_sigterm)
        signal.signal(signal.SIGINT, previous_sigint)
        await client.close()


async def main():
    import argparse

    parser = argparse.ArgumentParser(
        description="Confluence Full Content Crawler - 지식 체계 구축용",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # 인프라부문 샘플 10개
  uv run python scripts/confluence_full_crawler.py --source infra --sample 10

  # 홈쇼핑AX부문 전체 크롤링
  uv run python scripts/confluence_full_crawler.py --source homeshopping --full

  # 직접 페이지 ID 지정
  uv run python scripts/confluence_full_crawler.py --page-id 318803690 --full

  # 여러 소스 순차 크롤링 (모든 지식 소스)
  uv run python scripts/confluence_full_crawler.py --all-sources --full

  # 중단된 크롤링 재개 (체크포인트 사용)
  uv run python scripts/confluence_full_crawler.py --all-sources --full --resume

  # 사용 가능한 소스 목록
  uv run python scripts/confluence_full_crawler.py --list-sources
        """,
    )
    parser.add_argument("--source", type=str, choices=list(KNOWLEDGE_SOURCES.keys()),
                        help="지식 소스 선택 (infra, homeshopping)")
    parser.add_argument("--page-id", type=str, help="직접 페이지 ID 지정")
    parser.add_argument("--all-sources", action="store_true", help="모든 지식 소스 순차 크롤링")
    parser.add_argument("--list-sources", action="store_true", help="사용 가능한 소스 목록")
    parser.add_argument("--sample", type=int, default=10, help="샘플 페이지 수 (기본: 10)")
    parser.add_argument("--full", action="store_true", help="전체 크롤링 (페이지 수 제한 없음)")
    parser.add_argument("--no-attachments", action="store_true", help="첨부파일 다운로드 스킵")
    parser.add_argument("--max-attachments", type=int, default=20,
                        help="페이지당 최대 첨부파일 수 (기본: 20)")
    parser.add_argument("--max-concurrent", type=int, default=3,
                        help="동시 페이지 처리 수 (기본: 3, OMP_NUM_THREADS=2 기준 최적)")
    parser.add_argument("--resume", action="store_true",
                        help="이전 크롤링 중단 지점부터 재개")
    parser.add_argument("--fresh-full", action="store_true",
                        help="이전 체크포인트/JSONL 모두 삭제 후 전체 크롤링 (새 컬렉션용)")
    parser.add_argument("--kb-id", type=str, default="",
                        help="KB ID (체크포인트 검증용)")
    parser.add_argument(
        "--upload-s3",
        action="store_true",
        help="크롤링 결과 JSON을 S3를 통해 업로드",
    )
    parser.add_argument(
        "--s3-prefix",
        type=str,
        default="",
        help="S3 object key prefix (default: KNOWLEDGE_CRAWL_UPLOAD_PREFIX or knowledge/crawls/confluence)",
    )
    parser.add_argument(
        "--s3-bucket",
        type=str,
        default="",
        help="S3 bucket name (표시용; default: S3_UPLOAD_BUCKET or gs-retail-svc-dev-miso-files)",
    )
    args = parser.parse_args()

    # 소스 목록 표시
    if args.list_sources:
        show_available_sources()
        return

    # 크롤링 대상 결정
    sources_to_crawl: list[tuple[str, str, str]] = []  # (page_id, name, key)

    if args.all_sources:
        for key, info in KNOWLEDGE_SOURCES.items():
            sources_to_crawl.append((info["page_id"], info["name"], key))
    elif args.source:
        info = KNOWLEDGE_SOURCES[args.source]
        sources_to_crawl.append((info["page_id"], info["name"], args.source))
    elif args.page_id:
        sources_to_crawl.append((args.page_id, f"Custom Page ({args.page_id})", f"custom_{args.page_id}"))
    else:
        console.print("[yellow]⚠️ 크롤링 대상을 지정하세요.[/yellow]\n")
        show_available_sources()
        console.print("\n예시: uv run python scripts/confluence_full_crawler.py --source infra --sample 10")
        return

    # --fresh-full: 체크포인트/증분 파일 삭제 후 전체 크롤링
    if args.fresh_full:
        args.full = True
        args.resume = False
        # 모든 소스의 체크포인트 + 증분 파일 삭제
        checkpoint_file = OUTPUT_DIR / "checkpoint.json"
        if checkpoint_file.exists():
            checkpoint_file.unlink()
            console.print("[yellow]🗑️ --fresh-full: 체크포인트 삭제[/yellow]")
        for jsonl in OUTPUT_DIR.glob("incremental_*.jsonl"):
            jsonl.unlink()
            console.print(f"[yellow]🗑️ --fresh-full: {jsonl.name} 삭제[/yellow]")

    max_pages = None if args.full else args.sample
    download_attachments = not args.no_attachments
    upload_enabled = args.upload_s3 or (os.getenv("KNOWLEDGE_CRAWL_UPLOAD_ENABLED", "false").lower() == "true")
    s3_prefix = (args.s3_prefix or os.getenv("KNOWLEDGE_CRAWL_UPLOAD_PREFIX", "knowledge/crawls/confluence")).strip().strip("/")
    s3_bucket = (args.s3_bucket or os.getenv("S3_UPLOAD_BUCKET", "gs-retail-svc-dev-miso-files")).strip()
    run_id = datetime.now(timezone.utc).strftime("%Y%m%dT%H%M%SZ")

    upload_manifest: list[dict[str, Any]] = []
    doc_adapter = None

    # 크롤링 실행
    all_results: list[FullPageContent] = []
    all_page_dicts: list[dict] = []
    processed_sources: list[tuple[str, str, str]] = []
    streaming_jsonl_paths: list[str] = []  # 스트리밍 모드로 처리된 JSONL 경로

    if upload_enabled:
        try:
            from src.infrastructure.adapters.document_http_adapter import DocumentHttpAdapter

            doc_adapter = DocumentHttpAdapter()
            console.print(
                f"[cyan]☁️  S3 upload enabled[/cyan] (prefix={s3_prefix}, run_id={run_id})"
            )
        except Exception as e:
            console.print(f"[yellow]⚠️  Failed to init S3 document adapter: {e}[/yellow]")
            upload_enabled = False

    for page_id, source_name, source_key in sources_to_crawl:
        console.print(f"\n{'='*60}")
        crawl_result = await crawl_space(
            page_id=page_id,
            source_name=source_name,
            source_key=source_key,
            max_pages=max_pages,
            download_attachments=download_attachments,
            max_attachments_per_page=args.max_attachments,
            resume=args.resume,
            max_concurrent=args.max_concurrent,
            kb_id=args.kb_id,
        )
        pages = crawl_result.pages
        if not crawl_result.jsonl_path:
            # 일반 모드: 메모리에 축적
            all_results.extend(pages)
            all_page_dicts.extend(crawl_result.page_dicts)
        else:
            # 스트리밍 모드: all_results/all_page_dicts에 추가하지 않음 (OOM 방지)
            streaming_jsonl_paths.append(crawl_result.jsonl_path)
        processed_sources.append((page_id, source_name, source_key))

        # 개별 소스 저장
        safe_name = re.sub(r'[^\w]', '_', source_name)
        output_path = OUTPUT_DIR / f"crawl_{safe_name}.json"
        try:
            if crawl_result.jsonl_path:
                # 스트리밍 모드: JSONL → JSON 직접 변환 (메모리 효율적)
                save_results_from_jsonl(
                    Path(crawl_result.jsonl_path),
                    output_path,
                    source_info={
                        "page_id": page_id,
                        "name": source_name,
                        "key": source_key,
                    },
                )
            else:
                save_results(
                    pages,
                    output_path,
                    source_info={
                        "page_id": page_id,
                        "name": source_name,
                        "key": source_key,
                    },
                    page_dicts=crawl_result.page_dicts,
                )
        except Exception as save_err:
            console.print(f"[yellow]⚠️  JSON 저장 실패 ({source_key}): {save_err}[/yellow]")
            console.print(f"[dim]증분 JSONL은 보존됩니다: {OUTPUT_DIR / f'incremental_{source_key}.jsonl'}[/dim]")

        # JSON 저장 완료 후 체크포인트 삭제 (OOM 시 resume 가능하도록)
        if crawl_result.source_key and not crawl_result.interrupted:
            try:
                checkpoint_file = OUTPUT_DIR / "checkpoint.json"
                if checkpoint_file.exists():
                    import json as _json
                    with open(checkpoint_file) as _f:
                        cp_data = _json.load(_f)
                    if cp_data.get("source_key") == crawl_result.source_key:
                        checkpoint_file.unlink()
                        console.print(f"[dim]🗑️ 체크포인트 삭제 (저장 완료): {crawl_result.source_key}[/dim]")
            except Exception:
                pass

        if upload_enabled and doc_adapter:
            try:
                # Versioned + latest objects for simple retrieval.
                for slot in (run_id, "latest"):
                    object_key = f"{s3_prefix}/{source_key}/{slot}/crawl_{source_key}.json"
                    res = await doc_adapter.upload_object_from_path(
                        key=object_key,
                        file_path=output_path,
                        content_type="application/json",
                    )
                    s3_uri = f"s3://{s3_bucket}/{res.key}" if s3_bucket else res.key
                    upload_manifest.append(
                        {
                            "source_key": source_key,
                            "slot": slot,
                            "local_path": str(output_path),
                            "object_key": object_key,
                            "stored_key": res.key,
                            "s3_uri": s3_uri,
                            "size_bytes": res.size_bytes,
                            "etag": res.etag,
                        }
                    )
                console.print(
                    f"[green]☁️  Uploaded[/green] {source_key}: "
                    f"s3://{s3_bucket}/{s3_prefix}/{source_key}/{run_id}/crawl_{source_key}.json"
                )
            except Exception as e:
                console.print(f"[yellow]⚠️  S3 upload failed ({source_key}): {e}[/yellow]")

        if crawl_result.interrupted:
            console.print("[yellow]⚠️ 중단 요청이 감지되어 남은 소스 크롤링을 건너뜁니다.[/yellow]")
            break

    # 전체 통합 결과 저장 (단일 소스도 항상 생성)
    if processed_sources:
        combined_output = OUTPUT_DIR / "crawl_combined.json"
        if len(processed_sources) == 1:
            pid, name, key = processed_sources[0]
            combined_source_info: dict[str, Any] = {
                "page_id": pid,
                "name": name,
                "key": key,
            }
        else:
            combined_source_info = {
                "sources": [
                    {"page_id": pid, "name": name, "key": key}
                    for pid, name, key in processed_sources
                ],
            }
        try:
            if streaming_jsonl_paths and len(streaming_jsonl_paths) == 1:
                # 단일 소스 스트리밍: JSONL → JSON 변환
                save_results_from_jsonl(
                    Path(streaming_jsonl_paths[0]),
                    combined_output,
                    source_info=combined_source_info,
                )
            elif streaming_jsonl_paths:
                # 다중 소스 스트리밍: 개별 JSON은 이미 저장됨, combined는 스킵
                console.print("[dim]ℹ️  다중 소스 스트리밍 모드: combined JSON 스킵 (개별 JSON 참조)[/dim]")
            else:
                save_results(
                    all_results,
                    combined_output,
                    source_info=combined_source_info,
                    page_dicts=all_page_dicts,
                )
        except Exception as save_err:
            console.print(f"[yellow]⚠️  Combined JSON 저장 실패: {save_err}[/yellow]")
        if upload_enabled and doc_adapter:
            try:
                for slot in (run_id, "latest"):
                    object_key = f"{s3_prefix}/combined/{slot}/crawl_combined.json"
                    res = await doc_adapter.upload_object_from_path(
                        key=object_key,
                        file_path=combined_output,
                        content_type="application/json",
                    )
                    s3_uri = f"s3://{s3_bucket}/{res.key}" if s3_bucket else res.key
                    upload_manifest.append(
                        {
                            "source_key": "combined",
                            "slot": slot,
                            "local_path": str(combined_output),
                            "object_key": object_key,
                            "stored_key": res.key,
                            "s3_uri": s3_uri,
                            "size_bytes": res.size_bytes,
                            "etag": res.etag,
                        }
                    )
            except Exception as e:
                console.print(f"[yellow]⚠️  S3 upload failed (combined): {e}[/yellow]")

    if upload_enabled and upload_manifest:
        manifest_path = OUTPUT_DIR / "crawl_upload_manifest.json"
        try:
            manifest_path.write_text(
                json.dumps(
                    {
                        "run_id": run_id,
                        "s3_prefix": s3_prefix,
                        "s3_bucket": s3_bucket,
                        "uploaded": upload_manifest,
                    },
                    ensure_ascii=False,
                    indent=2,
                ),
                encoding="utf-8",
            )
        except Exception as e:
            console.print(f"[yellow]⚠️  Failed to write upload manifest: {e}[/yellow]")
        finally:
            if doc_adapter:
                try:
                    await doc_adapter.close()
                except Exception:
                    pass

    # 최종 요약
    console.print(f"\n{'='*60}")

    _streaming_page_count = 0
    if streaming_jsonl_paths:
        # 스트리밍 모드: JSONL에서 통계만 추출 (전체 로드 없이)
        summary_pages = []
        for _jp in streaming_jsonl_paths:
            try:
                with open(_jp, "r", encoding="utf-8") as _f:
                    for _line in _f:
                        if _line.strip():
                            _streaming_page_count += 1
            except Exception:
                pass
        console.print(f"[dim]📊 스트리밍 모드: {_streaming_page_count:,}페이지 (상세 통계 생략)[/dim]")
    else:
        summary_pages = all_page_dicts or [ConfluenceFullClient._page_to_dict(page) for page in all_results]

    # NEW: 전체 통계 계산
    final_labels = sum(len(page.get("labels", [])) for page in summary_pages)
    final_comments = sum(len(page.get("comments", [])) for page in summary_pages)
    final_emails = sum(len(page.get("emails", [])) for page in summary_pages)
    final_macros = sum(len(page.get("macros", [])) for page in summary_pages)
    final_internal_links = sum(len(page.get("internal_links", [])) for page in summary_pages)
    final_external_links = sum(len(page.get("external_links", [])) for page in summary_pages)
    final_restrictions = sum(len(page.get("restrictions", [])) for page in summary_pages)

    _total_page_count = _streaming_page_count if streaming_jsonl_paths else len(summary_pages)
    summary_text = (
        f"[bold green]✅ 크롤링 완료![/bold green]\n\n"
        f"📄 총 페이지: {_total_page_count:,}개\n"
        f"📊 총 테이블: {sum(len(page.get('tables', [])) for page in summary_pages):,}개\n"
        f"📎 총 첨부파일: {sum(len(page.get('attachments', [])) for page in summary_pages):,}개\n"
        f"👤 총 멘션: {sum(len(page.get('mentions', [])) for page in summary_pages):,}개\n"
    )

    # NEW: 추가 통계 (값이 있는 경우만)
    if final_labels > 0:
        summary_text += f"🏷️  총 라벨: {final_labels:,}개\n"
    if final_comments > 0:
        summary_text += f"💬 총 댓글: {final_comments:,}개\n"
    if final_emails > 0:
        summary_text += f"📧 총 이메일: {final_emails:,}개\n"
    if final_macros > 0:
        summary_text += f"🧩 총 매크로: {final_macros:,}개\n"
    if final_internal_links > 0:
        summary_text += f"🔗 총 내부링크: {final_internal_links:,}개\n"
    if final_external_links > 0:
        summary_text += f"🌐 총 외부링크: {final_external_links:,}개\n"
    if final_restrictions > 0:
        summary_text += f"🔒 접근제한 문서: {final_restrictions:,}건\n"

    summary_text += f"\n📁 결과 저장: {OUTPUT_DIR}"

    console.print(Panel.fit(summary_text, border_style="green"))


if __name__ == "__main__":
    try:
        asyncio.run(main())
    except KeyboardInterrupt:
        console.print("[yellow]⚠️ 사용자 중단[/yellow]")
        sys.exit(130)
    except Exception as exc:
        console.print(f"[red]❌ 크롤링 실패: {exc}[/red]")
        import traceback
        traceback.print_exc()
        sys.exit(1)
