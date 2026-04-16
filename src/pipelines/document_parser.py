"""Document parser for PDF, DOCX, PPTX, XLSX, and TXT files.

Extracted and simplified from oreo-ecosystem attachment_parser.py.
Enhanced with:
- PDF table extraction using PyMuPDF page.find_tables()
- Scanned PDF detection (pages with no extractable text)
- Image extraction from PDF (page.get_images() + doc.extract_image())
- Image extraction from PPTX (shape.image.blob for PICTURE shapes)
- OCR/CV Pipeline routing for extracted images
"""

from __future__ import annotations

import io
import logging
import shutil
import subprocess
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from src.config.weights import weights as _w

_EXT_PPTX = ".pptx"

logger = logging.getLogger(__name__)


@dataclass
class ParseResult:
    """Enhanced document parse result with text, tables, and OCR data."""

    text: str = ""
    tables: list[list[list[str]]] = field(default_factory=list)
    ocr_text: str = ""  # text extracted via OCR from images/scanned pages
    images_processed: int = 0
    visual_analyses: list[dict[str, Any]] = field(default_factory=list)
    file_modified_at: str = ""  # ISO timestamp from file metadata (PDF modDate, PPTX modified)

    @property
    def full_text(self) -> str:
        """Combined text from all sources."""
        parts = [self.text]
        if self.ocr_text:
            parts.append(f"\n[OCR Extracted Text]\n{self.ocr_text}")
        for table in self.tables:
            parts.append(f"\n[Table]\n{_table_to_markdown(table)}")
        return "\n".join(p for p in parts if p.strip())


MAX_FILE_SIZE = _w.pipeline.max_file_size_mb * 1024 * 1024  # default 200MB


def parse_file(filepath: str | Path) -> str:
    """Parse a local file and return its text content.

    Supports: PDF, DOCX, PPTX, XLSX, TXT/MD/CSV/JSON/XML/YAML.
    Returns empty string for unsupported types or on errors.
    """
    path = Path(filepath)
    if not path.exists() or not path.is_file():
        logger.warning("File not found: %s", path)
        return ""
    file_size = path.stat().st_size
    if file_size > MAX_FILE_SIZE:
        raise ValueError(f"File too large: {file_size / 1e6:.0f}MB (max {MAX_FILE_SIZE / 1e6:.0f}MB)")
    data = path.read_bytes()
    return parse_bytes(data, path.name)


def parse_file_enhanced(filepath: str | Path) -> ParseResult:
    """Parse a local file with enhanced output including tables and OCR.

    Returns a ParseResult with text, tables, OCR text, and visual analyses.
    """
    path = Path(filepath)
    if not path.exists() or not path.is_file():
        logger.warning("File not found: %s", path)
        return ParseResult()
    file_size = path.stat().st_size
    if file_size > MAX_FILE_SIZE:
        raise ValueError(f"File too large: {file_size / 1e6:.0f}MB (max {MAX_FILE_SIZE / 1e6:.0f}MB)")
    data = path.read_bytes()
    return parse_bytes_enhanced(data, path.name)


def parse_bytes(data: bytes, filename: str) -> str:
    """Parse file bytes and return plain text content."""
    if not data:
        raise ValueError(f"Empty file: {filename}")

    ext = Path(filename).suffix.lower()

    if ext == ".pdf":
        return _parse_pdf(data, filename)
    elif ext == ".docx":
        return _parse_docx(data, filename)
    elif ext == ".doc":
        logger.warning("Legacy .doc format is not supported (use .docx): %s", filename)
        return ""
    elif ext == _EXT_PPTX:
        return _parse_pptx(data, filename)
    elif ext == ".ppt":
        converted = _convert_ppt_to_pptx(data, filename)
        if converted is None:
            return ""
        return _parse_pptx(converted, filename)
    elif ext in (".xlsx", ".xls", ".xlsm"):
        return _parse_xlsx(data, filename)
    elif ext in (".txt", ".md", ".csv", ".json", ".xml", ".yaml", ".yml"):
        return _parse_text(data)
    else:
        logger.warning("Unsupported file type: %s", ext)
        return ""


def parse_bytes_enhanced(data: bytes, filename: str) -> ParseResult:
    """Parse file bytes with enhanced output including tables, OCR, and images."""
    ext = Path(filename).suffix.lower()

    if ext == ".pdf":
        return _parse_pdf_enhanced(data, filename)
    elif ext == _EXT_PPTX:
        return _parse_pptx_enhanced(data, filename)
    elif ext == ".ppt":
        converted = _convert_ppt_to_pptx(data, filename)
        if converted is None:
            return ParseResult(text=f"[Error] Failed to convert .ppt file: {filename}. LibreOffice (soffice) is required.")
        return _parse_pptx_enhanced(converted, filename)
    elif ext in (".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".tiff"):
        return _parse_image(data, filename)
    else:
        # For other types, wrap plain text in ParseResult
        text = parse_bytes(data, filename)
        return ParseResult(text=text)


# ---------------------------------------------------------------------------
# Legacy format conversion
# ---------------------------------------------------------------------------


def _convert_ppt_to_pptx(data: bytes, filename: str) -> bytes | None:
    """Convert legacy .ppt to .pptx using LibreOffice CLI.

    Returns the .pptx file bytes on success, or None if conversion fails
    (e.g. LibreOffice not installed).
    """
    soffice = shutil.which("soffice")
    if soffice is None:
        logger.warning(
            "LibreOffice (soffice) not found on PATH. Cannot convert .ppt file: %s. "
            "Install LibreOffice to enable .ppt support.",
            filename,
        )
        return None

    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            ppt_path = Path(tmpdir) / Path(filename).name
            ppt_path.write_bytes(data)

            result = subprocess.run(
                [
                    soffice,
                    "--headless",
                    "--convert-to",
                    "pptx",
                    "--outdir",
                    tmpdir,
                    str(ppt_path),
                ],
                capture_output=True,
                timeout=_w.timeouts.httpx_default,
            )

            if result.returncode != 0:
                stderr = result.stderr.decode("utf-8", errors="replace").strip()
                logger.warning(
                    "LibreOffice conversion failed for %s (exit %d): %s",
                    filename,
                    result.returncode,
                    stderr,
                )
                return None

            pptx_path = ppt_path.with_suffix(_EXT_PPTX)
            if not pptx_path.exists():
                logger.warning(
                    "LibreOffice conversion produced no output for %s", filename
                )
                return None

            converted_bytes = pptx_path.read_bytes()
            logger.info(
                "Converted .ppt to .pptx: %s (%d bytes)", filename, len(converted_bytes)
            )
            return converted_bytes

    except subprocess.TimeoutExpired:
        logger.warning(
            "LibreOffice conversion timed out (30s) for %s", filename
        )
        return None
    except Exception as e:  # noqa: BLE001
        logger.warning("LibreOffice conversion error for %s: %s", filename, e)
        return None


# ---------------------------------------------------------------------------
# Individual parsers
# ---------------------------------------------------------------------------


def _parse_pdf(data: bytes, _filename: str) -> str:
    """Parse PDF using pymupdf (PyMuPDF)."""
    import pymupdf

    try:
        doc = pymupdf.open(stream=data, filetype="pdf")
    except Exception as e:
        raise ValueError(f"PDF open failed (encrypted or corrupt?): {e}") from e
    texts = []
    for page_num, page in enumerate(doc):
        text = page.get_text()
        if text.strip():
            texts.append(f"[Page {page_num + 1}]\n{text}")
    doc.close()
    return "\n\n".join(texts)


def _extract_pdf_date(raw: str) -> str:
    """Parse PDF date format (D:YYYYMMDDHHmmSS+TZ) to ISO 8601."""
    if not raw:
        return ""
    import re as _re
    m = _re.match(r"D:(\d{4})(\d{2})(\d{2})(\d{2})(\d{2})(\d{2})", raw)
    if m:
        return f"{m.group(1)}-{m.group(2)}-{m.group(3)}T{m.group(4)}:{m.group(5)}:{m.group(6)}"
    return ""


def _extract_page_tables(page, tables: list) -> None:
    """Extract tables from a PDF page."""
    for table in page.find_tables():
        table_data = table.extract()
        if table_data:
            cleaned = [
                [str(cell) if cell is not None else "" for cell in row]
                for row in table_data
            ]
            tables.append(cleaned)


def _extract_page_images(page, doc, extracted_images: list) -> None:
    """Extract images from a PDF page."""
    for img_info in page.get_images(full=True):
        xref = img_info[0]
        img_data = doc.extract_image(xref)
        if img_data and img_data.get("image"):
            img_bytes = img_data["image"]
            if 1024 < len(img_bytes) < 10_000_000:
                extracted_images.append(img_bytes)


def _extract_pdf_page_heading(page) -> str:
    """PDF 페이지에서 폰트 크기 기반으로 heading(제목) 추출.

    페이지 상단 20% 영역에서 가장 큰 폰트의 텍스트를 heading으로 판단.
    """
    try:
        blocks = page.get_text("dict", flags=0)["blocks"]
        if not blocks:
            return ""
        page_height = page.rect.height
        heading_zone = page_height * 0.2  # 상단 20%

        candidates = []
        for block in blocks:
            if block.get("type") != 0:  # text block만
                continue
            for line in block.get("lines", []):
                # 상단 영역만
                if line["bbox"][1] > heading_zone:
                    continue
                for span in line.get("spans", []):
                    text = span.get("text", "").strip()
                    size = span.get("size", 0)
                    flags = span.get("flags", 0)
                    is_bold = bool(flags & 2**4)  # bit 4 = bold
                    if text and len(text) >= 2:
                        # 볼드이면 크기 보너스
                        effective_size = size * 1.2 if is_bold else size
                        candidates.append((effective_size, text))

        if not candidates:
            return ""
        # 가장 큰 폰트 텍스트
        candidates.sort(key=lambda x: -x[0])
        best_size, best_text = candidates[0]
        # 본문 크기(보통 10~12pt) 대비 충분히 커야 heading
        if best_size < 13:
            return ""
        return best_text[:100]
    except Exception:  # noqa: BLE001
        return ""


def _classify_pdf_page(
    page, page_num: int, filename: str, doc,
    has_broken_cmap_fn, is_garbled_fn,
    texts: list, scanned_pages: list, tables: list, extracted_images: list,
) -> None:
    """Classify a single PDF page and extract text/tables/images."""
    has_broken_fonts = has_broken_cmap_fn(page)
    text = page.get_text()
    page_label = page_num + 1

    if has_broken_fonts or (text.strip() and is_garbled_fn(text)):
        scanned_pages.append(page_label)
        if has_broken_fonts:
            logger.info("Broken CMap font on page %d of %s, routing to OCR", page_label, filename)
        elif text.strip():
            logger.info("Garbled text on page %d of %s, routing to OCR", page_label, filename)
    elif text.strip():
        # PDF heading 추출 — 폰트 크기 기반
        heading = _extract_pdf_page_heading(page)
        if heading:
            texts.append(f"## {heading}\n\n[Page {page_label}]\n{text}")
        else:
            texts.append(f"[Page {page_label}]\n{text}")
    else:
        scanned_pages.append(page_label)

    _extract_page_tables(page, tables)
    _extract_page_images(page, doc, extracted_images)


def _has_broken_cmap_fonts(page, doc) -> bool:
    """Detect fonts with broken ToUnicode CMaps (common in PowerPoint PDF exports).

    Type0/Identity-H fonts with large embedded TrueType but very few Unicode
    mappings indicate stripped CMap tables — all Korean chars map to a single
    character (e.g. 폐 or 손).
    """
    import re as _re
    for xref, _ext, ftype, _basefont, _name, encoding in page.get_fonts():
        if ftype != "Type0" or encoding != "Identity-H":
            continue
        if _check_font_broken_cmap(doc, xref, _re):
            return True
    return False


def _check_font_broken_cmap(doc, xref: int, _re) -> bool:
    """Check if a single Type0 font has a broken CMap."""
    try:
        obj = doc.xref_object(xref)
        tounicode_match = _re.search(r"/ToUnicode (\d+) 0 R", obj)
        if not tounicode_match:
            return True
        cmap = doc.xref_stream(int(tounicode_match.group(1))).decode("latin-1", errors="replace")
        bfchar_count = sum(int(m) for m in _re.findall(r"(\d+)\s+beginbfchar", cmap))
        bfrange_count = sum(int(m) for m in _re.findall(r"(\d+)\s+beginbfrange", cmap))
        total_mappings = bfchar_count + bfrange_count

        font_size = _get_embedded_font_size(doc, obj, _re)
        if font_size is None:
            return False
        return font_size > 10_000 and total_mappings < 20
    except Exception:  # noqa: BLE001
        return False


def _get_embedded_font_size(doc, font_obj: str, _re) -> int | None:
    """Get embedded font file size from a CID font descriptor chain."""
    desc_match = _re.search(r"/DescendantFonts\s+(\d+)\s+0\s+R", font_obj)
    if not desc_match:
        return None
    desc_arr = doc.xref_object(int(desc_match.group(1)))
    cidfont_xref = _re.search(r"(\d+)\s+0\s+R", desc_arr)
    if not cidfont_xref:
        return None
    cidfont = doc.xref_object(int(cidfont_xref.group(1)))
    fd_match = _re.search(r"/FontDescriptor\s+(\d+)\s+0\s+R", cidfont)
    if not fd_match:
        return None
    fd = doc.xref_object(int(fd_match.group(1)))
    ff2_match = _re.search(r"/FontFile2\s+(\d+)\s+0\s+R", fd)
    if not ff2_match:
        return None
    return len(doc.xref_stream(int(ff2_match.group(1))))


def _is_garbled_text(text: str) -> bool:
    """Detect garbled text from broken font mapping in PDFs."""
    if not text or len(text.strip()) < 10:
        return False
    clean = text.strip().replace(" ", "").replace("\n", "")
    if not clean:
        return False
    from collections import Counter
    char_counts = Counter(clean)
    total = len(clean)

    top1_ratio = char_counts.most_common(1)[0][1] / total
    if top1_ratio > 0.25:
        return True

    cjk_counts = [(ch, cnt) for ch, cnt in char_counts.most_common(10)
                   if '\u4e00' <= ch <= '\u9fff' or '\uac00' <= ch <= '\ud7a3']
    if len(cjk_counts) >= 2:
        top3_cjk_total = sum(cnt for _, cnt in cjk_counts[:3])
        if top3_cjk_total / total > 0.4:
            return True

    unique_ratio = len(set(clean)) / total
    if unique_ratio < 0.08 and total > 30:
        return True
    return False


def _parse_pdf_enhanced(data: bytes, filename: str) -> ParseResult:
    """Enhanced PDF parsing with table extraction, scanned page detection, and image OCR."""
    import pymupdf

    try:
        doc = pymupdf.open(stream=data, filetype="pdf")
    except Exception as e:
        raise ValueError(f"PDF open failed (encrypted or corrupt?): {e}") from e

    file_modified_at = _extract_pdf_date(doc.metadata.get("modDate", "")) or \
                       _extract_pdf_date(doc.metadata.get("creationDate", ""))

    texts = []
    tables: list[list[list[str]]] = []
    scanned_pages: list[int] = []
    extracted_images: list[bytes] = []

    for page_num, page in enumerate(doc):
        _classify_pdf_page(
            page, page_num, filename, doc,
            lambda p: _has_broken_cmap_fonts(p, doc), _is_garbled_text,
            texts, scanned_pages, tables, extracted_images,
        )

    doc.close()

    ocr_texts = []
    visual_analyses = []
    total_images = 0

    if scanned_pages:
        logger.info("Detected %d scanned/garbled pages in %s: %s", len(scanned_pages), filename, scanned_pages)
        doc2 = pymupdf.open(stream=data, filetype="pdf")
        for page_num in scanned_pages:
            page = doc2[page_num - 1]
            pix = page.get_pixmap(dpi=300)
            img_bytes = pix.tobytes("png")
            logger.info("Rendered page %d: %dx%d (%d bytes)", page_num, pix.width, pix.height, len(img_bytes))
            page_ocr_text, page_analyses = _process_images_ocr([img_bytes])
            if page_ocr_text.strip():
                # Replace [Image 1 OCR] with [Page N OCR] for page renders
                page_ocr_clean = page_ocr_text.replace("[Image 1 OCR] ", "").strip()
                ocr_texts.append(f"[Page {page_num} OCR] {page_ocr_clean}")
                total_images += 1
            visual_analyses.extend(page_analyses)
        doc2.close()

    # Process embedded images (non-page) — keep [Image N OCR] tags
    if extracted_images:
        img_ocr_text, img_analyses = _process_images_ocr(extracted_images)
        if img_ocr_text.strip():
            ocr_texts.append(img_ocr_text)
            total_images += len(extracted_images)
        visual_analyses.extend(img_analyses)

    return ParseResult(
        text="\n\n".join(texts),
        tables=tables,
        ocr_text="\n".join(ocr_texts),
        images_processed=total_images,
        visual_analyses=visual_analyses,
        file_modified_at=file_modified_at,
    )


def _format_docx_paragraph(para) -> str | None:
    """Format a DOCX paragraph, returning None if empty."""
    if not para.text.strip():
        return None
    if para.style and para.style.name.startswith("Heading"):
        level = para.style.name[-1] if para.style.name[-1].isdigit() else "1"
        return f"{'#' * int(level)} {para.text}"
    return para.text


def _extract_table_data(table) -> list[list[str]]:
    """Extract table rows as list of lists."""
    return [[cell.text.strip() for cell in row.cells] for row in table.rows]


def _parse_docx(data: bytes, _filename: str) -> str:
    """Parse DOCX using python-docx."""
    from docx import Document

    try:
        doc = Document(io.BytesIO(data))
    except Exception as e:
        raise ValueError(f"DOCX open failed (corrupt?): {e}") from e

    texts = [t for p in doc.paragraphs if (t := _format_docx_paragraph(p))]

    for table in doc.tables:
        table_data = _extract_table_data(table)
        if table_data:
            texts.append(_table_to_markdown(table_data))

    return "\n\n".join(texts)


def _iter_pptx_shapes(shapes, _depth: int = 0):
    """Recursively yield shapes from PPTX, handling grouped shapes."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    if _depth > 10:
        return
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_pptx_shapes(shape.shapes, _depth + 1)


def _extract_slide_title(slide) -> str:
    """PPTX 슬라이드 제목 추출 (title placeholder 또는 첫 번째 큰 텍스트)."""
    # 1. Title placeholder에서 추출 (가장 정확)
    if slide.shapes.title and slide.shapes.title.text.strip():
        return slide.shapes.title.text.strip()[:100]
    # 2. placeholder type 13(TITLE) 또는 15(CENTER_TITLE) 탐색
    try:
        for shape in slide.placeholders:
            if shape.placeholder_format.idx in (0, 13, 15):
                if shape.text.strip():
                    return shape.text.strip()[:100]
    except Exception:  # noqa: BLE001
        pass
    return ""


def _extract_slide_text(slide, slide_num: int) -> str | None:
    """Extract text from a single PPTX slide."""
    title = _extract_slide_title(slide)
    if title:
        slide_texts = [f"## {title}\n\n[Slide {slide_num}]"]
    else:
        slide_texts = [f"[Slide {slide_num}]"]
    for shape in _iter_pptx_shapes(slide.shapes):
        if hasattr(shape, "text") and shape.text.strip():
            slide_texts.append(shape.text)
        if shape.has_table:
            table_data = _extract_table_data(shape.table)
            if table_data:
                slide_texts.append(_table_to_markdown(table_data))
    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            slide_texts.append(f"[Notes] {notes_text}")
    return "\n".join(slide_texts) if len(slide_texts) > 1 else None


def _parse_pptx(data: bytes, _filename: str) -> str:
    """Parse PPTX using python-pptx."""
    from pptx import Presentation

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as e:
        raise ValueError(f"PPTX open failed (corrupt?): {e}") from e

    texts = [t for i, s in enumerate(prs.slides, 1) if (t := _extract_slide_text(s, i))]
    return "\n\n".join(texts)


def _extract_pptx_modified_date(prs) -> str:
    """Extract modification date from PPTX core properties."""
    try:
        modified = prs.core_properties.modified
        return modified.isoformat() if modified else ""
    except Exception:  # noqa: BLE001
        return ""


def _process_pptx_shape(shape, slide_texts: list, tables: list, images: list) -> None:
    """Process a single PPTX shape: extract text, table, or image."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if hasattr(shape, "text") and shape.text.strip():
        slide_texts.append(shape.text)
    if shape.has_table:
        table_data = _extract_table_data(shape.table)
        if table_data:
            tables.append(table_data)
            slide_texts.append(_table_to_markdown(table_data))
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        img_bytes = shape.image.blob
        if img_bytes and 1024 < len(img_bytes) < 10_000_000:
            images.append(img_bytes)


def _process_enhanced_slide(slide, slide_num: int, tables: list, images: list) -> str | None:
    """Process a single slide for enhanced parsing, collecting tables and images."""
    title = _extract_slide_title(slide)
    if title:
        slide_texts = [f"## {title}\n\n[Slide {slide_num}]"]
    else:
        slide_texts = [f"[Slide {slide_num}]"]
    for shape in _iter_pptx_shapes(slide.shapes):
        _process_pptx_shape(shape, slide_texts, tables, images)
    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            slide_texts.append(f"[Notes] {notes_text}")
    return "\n".join(slide_texts) if len(slide_texts) > 1 else None


def _parse_pptx_enhanced(data: bytes, _filename: str) -> ParseResult:
    """Enhanced PPTX parsing with image extraction and OCR routing."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    file_modified_at = _extract_pptx_modified_date(prs)

    texts = []
    tables: list[list[list[str]]] = []
    extracted_images: list[bytes] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        text = _process_enhanced_slide(slide, slide_num, tables, extracted_images)
        if text:
            texts.append(text)

    visual_analyses = []
    ocr_text = ""
    if extracted_images:
        ocr_text, visual_analyses = _process_images_ocr(extracted_images)

    return ParseResult(
        text="\n\n".join(texts),
        tables=tables,
        ocr_text=ocr_text,
        images_processed=len(extracted_images),
        visual_analyses=visual_analyses,
        file_modified_at=file_modified_at,
    )


def _parse_image(data: bytes, _filename: str) -> ParseResult:
    """Parse image file through OCR/CV pipeline."""
    ocr_text, visual_analyses = _process_images_ocr([data])
    return ParseResult(
        ocr_text=ocr_text,
        images_processed=1,
        visual_analyses=visual_analyses,
    )


def _parse_xlsx(data: bytes, _filename: str) -> str:
    """Parse XLSX using openpyxl."""
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    texts = []

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        texts.append(f"[Sheet: {sheet_name}]")
        sheet_data = []
        for row in sheet.iter_rows(values_only=True):
            row_data = [str(cell) if cell is not None else "" for cell in row]
            if any(cell.strip() for cell in row_data):
                sheet_data.append(row_data)
        if sheet_data:
            texts.append(_table_to_markdown(sheet_data))

    wb.close()
    return "\n\n".join(texts)


def _parse_text(data: bytes) -> str:
    """Parse text files with UTF-8 / EUC-KR fallback."""
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        return data.decode("euc-kr", errors="ignore")


# ---------------------------------------------------------------------------
# OCR / CV Pipeline routing
# ---------------------------------------------------------------------------


def _resize_image(raw: bytes, scale: float = 0.75) -> bytes | None:
    """Resize image as fallback when OCR fails on the original."""
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(raw))
        new_w = max(int(img.width * scale), 32)
        new_h = max(int(img.height * scale), 32)
        img = img.resize((new_w, new_h), Image.LANCZOS)
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue()
    except Exception:  # noqa: BLE001
        return None


def _ocr_request(client, endpoint: str, payload: bytes) -> dict | None:
    """Send OCR request. Returns parsed JSON or None on failure."""
    import base64
    b64_image = base64.b64encode(payload).decode("utf-8")
    try:
        resp = client.post(endpoint, json={"image": b64_image}, timeout=_w.timeouts.httpx_ocr)
        resp.raise_for_status()
        return resp.json()
    except Exception:  # noqa: BLE001
        return None


def _prepare_image_bytes(img_bytes: bytes, index: int) -> bytes | None:
    """Validate, decode, and convert image to PNG. Returns None if invalid."""
    from PIL import Image
    _MIN_W, _MIN_H = 20, 20

    try:
        _img = Image.open(io.BytesIO(img_bytes))
    except Exception:  # noqa: BLE001
        logger.debug("Skipping image %d: cannot decode", index)
        return None
    if _img.width < _MIN_W or _img.height < _MIN_H:
        logger.debug("Skipping image %d: too small (%dx%d)", index, _img.width, _img.height)
        return None
    if _img.mode not in ("RGB", "L"):
        _img = _img.convert("RGB")
    _buf = io.BytesIO()
    _img.save(_buf, format="PNG")
    return _buf.getvalue()


def _extract_text_from_boxes(boxes: list[dict], min_conf: float, index: int) -> str:
    """Extract text from OCR boxes with confidence filtering."""
    filtered_texts = []
    dropped = 0
    for box in boxes:
        conf = box.get("confidence", 0.0)
        box_text = box.get("text", "")
        if conf >= min_conf and box_text.strip():
            filtered_texts.append(box_text)
        elif box_text.strip():
            dropped += 1
    if dropped > 0:
        logger.info(
            "OCR confidence filter: kept %d, dropped %d (< %.0f%%) for image %d",
            len(filtered_texts), dropped, min_conf * 100, index,
        )
    return " ".join(filtered_texts)


def _extract_text_fallback(result: dict) -> str:
    """Extract text from OCR result when no boxes are available."""
    text_lines = result.get("texts", result.get("result", []))
    if isinstance(text_lines, list):
        return " ".join(str(t) for t in text_lines if t)
    if isinstance(text_lines, str):
        return text_lines
    return str(text_lines)


def _extract_vision_analysis(
    result: dict, index: int,
    ocr_texts: list[str], visual_analyses: list[dict[str, Any]],
) -> None:
    """Extract vision analysis (shapes/arrows/mappings) from OCR result."""
    shapes = result.get("shapes", [])
    arrows = result.get("arrows", [])
    mappings = result.get("text_shape_mappings", [])
    if not shapes and not arrows and not mappings:
        return

    visual_analyses.append({
        "image_index": index,
        "shapes": shapes,
        "arrows": arrows,
        "text_shape_mappings": mappings,
        "shape_count": result.get("shape_count", len(shapes)),
        "arrow_count": result.get("arrow_count", len(arrows)),
        "ocr_confidence": result.get("ocr_confidence", 0.0),
    })

    desc_parts: list[str] = []
    if shapes:
        shape_types = [s.get("type", "unknown") for s in shapes]
        desc_parts.append(f"Shapes: {', '.join(shape_types)}")
    if arrows:
        desc_parts.append(f"Arrows: {len(arrows)} connections")
    if mappings:
        for m in mappings:
            shape_label = m.get("shape_type", "shape")
            texts_in = m.get("texts", m.get("text", ""))
            if texts_in:
                desc_parts.append(f"  [{shape_label}] {texts_in}")
    if desc_parts:
        diagram_desc = "\n".join(desc_parts)
        ocr_texts.append(f"[Image {index} Diagram]\n{diagram_desc}")


def _process_single_image_ocr(
    client,
    endpoint: str,
    img_bytes: bytes,
    index: int,
    min_conf: float,
    vision_enabled: bool,
    ocr_texts: list[str],
    visual_analyses: list[dict[str, Any]],
) -> None:
    """Process a single image through OCR with retry on failure."""
    try:
        prepared = _prepare_image_bytes(img_bytes, index)
        if prepared is None:
            return

        result = _ocr_request(client, endpoint, prepared)
        if result is None:
            resized = _resize_image(prepared)
            if resized:
                logger.info("Retrying image %d with resized version", index)
                result = _ocr_request(client, endpoint, resized)

        if result is None:
            logger.warning("PaddleOCR failed for image %d after retry, skipping", index)
            return

        boxes = result.get("boxes", [])
        text = (
            _extract_text_from_boxes(boxes, min_conf, index)
            if boxes
            else _extract_text_fallback(result)
        )

        if text.strip():
            ocr_texts.append(f"[Image {index} OCR] {text}")
            logger.info("OCR success for image %d: %d chars", index, len(text))

        if vision_enabled:
            _extract_vision_analysis(result, index, ocr_texts, visual_analyses)

    except Exception as e:  # noqa: BLE001
        logger.warning("PaddleOCR API failed for image %d: %s", index, e)


def _process_images_ocr(
    images: list[bytes],
) -> tuple[str, list[dict[str, Any]]]:
    """Route extracted images through PaddleOCR Docker API server."""
    import httpx
    import os

    if not images:
        return "", []

    ocr_texts: list[str] = []
    visual_analyses: list[dict[str, Any]] = []
    base_url = os.getenv("PADDLEOCR_API_URL", "http://localhost:8866")
    if base_url.endswith("/ocr") or base_url.endswith("/analyze"):
        base_url = base_url.rsplit("/", 1)[0]

    vision_enabled = _w.ocr.enable_vision_analysis
    endpoint = f"{base_url}/analyze" if vision_enabled else f"{base_url}/ocr"
    min_conf = float(os.getenv("OCR_MIN_CONFIDENCE", "0.65"))
    client = httpx.Client(timeout=_w.timeouts.httpx_ocr)

    for i, img_bytes in enumerate(images):
        _process_single_image_ocr(
            client, endpoint, img_bytes, i + 1, min_conf, vision_enabled,
            ocr_texts, visual_analyses,
        )

    client.close()
    return "\n".join(ocr_texts), visual_analyses


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _table_to_markdown(data: list[list[str]], max_rows: int = 50) -> str:
    """Convert table data to markdown format."""
    if not data:
        return ""
    lines = []
    header = data[0]
    lines.append("| " + " | ".join(header) + " |")
    lines.append("| " + " | ".join(["---"] * len(header)) + " |")
    for row in data[1:max_rows]:
        padded_row = row + [""] * (len(header) - len(row))
        padded_row = padded_row[: len(header)]
        lines.append("| " + " | ".join(padded_row) + " |")
    if len(data) > max_rows:
        lines.append(f"... ({len(data) - max_rows} rows omitted)")
    return "\n".join(lines)
