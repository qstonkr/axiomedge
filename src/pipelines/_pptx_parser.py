"""Document parser -- PPTX parsing functions.

Extracted from document_parser.py for SRP.
"""

from __future__ import annotations

import io
import logging
from collections.abc import Iterator

from . import _parser_utils
from ._parser_utils import (
    ParseResult,
    _extract_table_data,
    _table_to_markdown,
)

logger = logging.getLogger(__name__)


def _iter_pptx_shapes(shapes, _depth: int = 0) -> Iterator:
    """Recursively yield shapes from PPTX, handling grouped shapes."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    if _depth > 10:
        return
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_pptx_shapes(shape.shapes, _depth + 1)


def _extract_slide_title(slide) -> str:
    """PPTX 슬라이드 제목 추출 (title placeholder 또는 첫 번째 큰 텍스트)."""
    # 1. Title placeholder에서 추출 (가장 정확)
    if slide.shapes.title and slide.shapes.title.text.strip():
        return slide.shapes.title.text.strip()[:100]
    # 2. placeholder type 13(TITLE) 또는 15(CENTER_TITLE) 탐색
    try:
        for shape in slide.placeholders:
            if shape.placeholder_format.idx in (0, 13, 15):
                if shape.text.strip():
                    return shape.text.strip()[:100]
    except (RuntimeError, OSError, ValueError, TypeError, KeyError, AttributeError):
        pass
    return ""


def _extract_slide_text(slide, slide_num: int) -> str | None:
    """Extract text from a single PPTX slide."""
    title = _extract_slide_title(slide)
    if title:
        slide_texts = [f"## {title}\n\n[Slide {slide_num}]"]
    else:
        slide_texts = [f"[Slide {slide_num}]"]
    for shape in _iter_pptx_shapes(slide.shapes):
        if hasattr(shape, "text") and shape.text.strip():
            slide_texts.append(shape.text)
        if shape.has_table:
            table_data = _extract_table_data(shape.table)
            if table_data:
                slide_texts.append(_table_to_markdown(table_data))
    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            slide_texts.append(f"[Notes] {notes_text}")
    return "\n".join(slide_texts) if len(slide_texts) > 1 else None


def _parse_pptx(data: bytes, _filename: str) -> str:
    """Parse PPTX using python-pptx."""
    from pptx import Presentation

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as e:
        raise ValueError(f"PPTX open failed (corrupt?): {e}") from e

    texts = [t for i, s in enumerate(prs.slides, 1) if (t := _extract_slide_text(s, i))]
    return "\n\n".join(texts)


def _extract_pptx_modified_date(prs) -> str:
    """Extract modification date from PPTX core properties."""
    try:
        modified = prs.core_properties.modified
        return modified.isoformat() if modified else ""
    except (RuntimeError, OSError, ValueError, TypeError, KeyError, AttributeError):
        return ""


def _process_pptx_shape(shape, slide_texts: list, tables: list, images: list) -> None:
    """Process a single PPTX shape: extract text, table, or image."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if hasattr(shape, "text") and shape.text.strip():
        slide_texts.append(shape.text)
    if shape.has_table:
        table_data = _extract_table_data(shape.table)
        if table_data:
            tables.append(table_data)
            slide_texts.append(_table_to_markdown(table_data))
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        img_bytes = shape.image.blob
        if img_bytes and 1024 < len(img_bytes) < 10_000_000:
            images.append(img_bytes)


def _process_enhanced_slide(slide, slide_num: int, tables: list, images: list) -> str | None:
    """Process a single slide for enhanced parsing, collecting tables and images."""
    title = _extract_slide_title(slide)
    if title:
        slide_texts = [f"## {title}\n\n[Slide {slide_num}]"]
    else:
        slide_texts = [f"[Slide {slide_num}]"]
    for shape in _iter_pptx_shapes(slide.shapes):
        _process_pptx_shape(shape, slide_texts, tables, images)
    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            slide_texts.append(f"[Notes] {notes_text}")
    return "\n".join(slide_texts) if len(slide_texts) > 1 else None


def _parse_pptx_enhanced(data: bytes, _filename: str) -> ParseResult:
    """Enhanced PPTX parsing with image extraction and OCR routing."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    file_modified_at = _extract_pptx_modified_date(prs)

    texts = []
    tables: list[list[list[str]]] = []
    extracted_images: list[bytes] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        text = _process_enhanced_slide(slide, slide_num, tables, extracted_images)
        if text:
            texts.append(text)

    visual_analyses = []
    ocr_text = ""
    if extracted_images:
        ocr_text, visual_analyses = _parser_utils._process_images_ocr(extracted_images)

    return ParseResult(
        text="\n\n".join(texts),
        tables=tables,
        ocr_text=ocr_text,
        images_processed=len(extracted_images),
        visual_analyses=visual_analyses,
        file_modified_at=file_modified_at,
    )
