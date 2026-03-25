"""Document parser for PDF, DOCX, PPTX, XLSX, and TXT files.

Extracted and simplified from oreo-ecosystem attachment_parser.py.
Enhanced with:
- PDF table extraction using PyMuPDF page.find_tables()
- Scanned PDF detection (pages with no extractable text)
- Image extraction from PDF (page.get_images() + doc.extract_image())
- Image extraction from PPTX (shape.image.blob for PICTURE shapes)
- OCR/CV Pipeline routing for extracted images
"""

from __future__ import annotations

import io
import logging
import shutil
import subprocess
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)


@dataclass
class ParseResult:
    """Enhanced document parse result with text, tables, and OCR data."""

    text: str = ""
    tables: list[list[list[str]]] = field(default_factory=list)
    ocr_text: str = ""  # text extracted via OCR from images/scanned pages
    images_processed: int = 0
    visual_analyses: list[dict[str, Any]] = field(default_factory=list)

    @property
    def full_text(self) -> str:
        """Combined text from all sources."""
        parts = [self.text]
        if self.ocr_text:
            parts.append(f"\n[OCR Extracted Text]\n{self.ocr_text}")
        for table in self.tables:
            parts.append(f"\n[Table]\n{_table_to_markdown(table)}")
        return "\n".join(p for p in parts if p.strip())


from src.config_weights import weights as _w

MAX_FILE_SIZE = _w.pipeline.max_file_size_mb * 1024 * 1024  # default 200MB


def parse_file(filepath: str | Path) -> str:
    """Parse a local file and return its text content.

    Supports: PDF, DOCX, PPTX, XLSX, TXT/MD/CSV/JSON/XML/YAML.
    Returns empty string for unsupported types or on errors.
    """
    path = Path(filepath)
    if not path.exists() or not path.is_file():
        logger.warning("File not found: %s", path)
        return ""
    file_size = path.stat().st_size
    if file_size > MAX_FILE_SIZE:
        raise ValueError(f"File too large: {file_size / 1e6:.0f}MB (max {MAX_FILE_SIZE / 1e6:.0f}MB)")
    data = path.read_bytes()
    return parse_bytes(data, path.name)


def parse_file_enhanced(filepath: str | Path) -> ParseResult:
    """Parse a local file with enhanced output including tables and OCR.

    Returns a ParseResult with text, tables, OCR text, and visual analyses.
    """
    path = Path(filepath)
    if not path.exists() or not path.is_file():
        logger.warning("File not found: %s", path)
        return ParseResult()
    file_size = path.stat().st_size
    if file_size > MAX_FILE_SIZE:
        raise ValueError(f"File too large: {file_size / 1e6:.0f}MB (max {MAX_FILE_SIZE / 1e6:.0f}MB)")
    data = path.read_bytes()
    return parse_bytes_enhanced(data, path.name)


def parse_bytes(data: bytes, filename: str) -> str:
    """Parse file bytes and return plain text content."""
    if not data:
        raise ValueError(f"Empty file: {filename}")

    ext = Path(filename).suffix.lower()

    if ext == ".pdf":
        return _parse_pdf(data, filename)
    elif ext == ".docx":
        return _parse_docx(data, filename)
    elif ext == ".doc":
        logger.warning("Legacy .doc format is not supported (use .docx): %s", filename)
        return ""
    elif ext == ".pptx":
        return _parse_pptx(data, filename)
    elif ext == ".ppt":
        converted = _convert_ppt_to_pptx(data, filename)
        if converted is None:
            return ""
        return _parse_pptx(converted, filename)
    elif ext in (".xlsx", ".xls", ".xlsm"):
        return _parse_xlsx(data, filename)
    elif ext in (".txt", ".md", ".csv", ".json", ".xml", ".yaml", ".yml"):
        return _parse_text(data)
    else:
        logger.warning("Unsupported file type: %s", ext)
        return ""


def parse_bytes_enhanced(data: bytes, filename: str) -> ParseResult:
    """Parse file bytes with enhanced output including tables, OCR, and images."""
    ext = Path(filename).suffix.lower()

    if ext == ".pdf":
        return _parse_pdf_enhanced(data, filename)
    elif ext == ".pptx":
        return _parse_pptx_enhanced(data, filename)
    elif ext == ".ppt":
        converted = _convert_ppt_to_pptx(data, filename)
        if converted is None:
            return ParseResult(text=f"[Error] Failed to convert .ppt file: {filename}. LibreOffice (soffice) is required.")
        return _parse_pptx_enhanced(converted, filename)
    elif ext in (".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".tiff"):
        return _parse_image(data, filename)
    else:
        # For other types, wrap plain text in ParseResult
        text = parse_bytes(data, filename)
        return ParseResult(text=text)


# ---------------------------------------------------------------------------
# Legacy format conversion
# ---------------------------------------------------------------------------


def _convert_ppt_to_pptx(data: bytes, filename: str) -> bytes | None:
    """Convert legacy .ppt to .pptx using LibreOffice CLI.

    Returns the .pptx file bytes on success, or None if conversion fails
    (e.g. LibreOffice not installed).
    """
    soffice = shutil.which("soffice")
    if soffice is None:
        logger.warning(
            "LibreOffice (soffice) not found on PATH. Cannot convert .ppt file: %s. "
            "Install LibreOffice to enable .ppt support.",
            filename,
        )
        return None

    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            ppt_path = Path(tmpdir) / Path(filename).name
            ppt_path.write_bytes(data)

            result = subprocess.run(
                [
                    soffice,
                    "--headless",
                    "--convert-to",
                    "pptx",
                    "--outdir",
                    tmpdir,
                    str(ppt_path),
                ],
                capture_output=True,
                timeout=30,
            )

            if result.returncode != 0:
                stderr = result.stderr.decode("utf-8", errors="replace").strip()
                logger.warning(
                    "LibreOffice conversion failed for %s (exit %d): %s",
                    filename,
                    result.returncode,
                    stderr,
                )
                return None

            pptx_path = ppt_path.with_suffix(".pptx")
            if not pptx_path.exists():
                logger.warning(
                    "LibreOffice conversion produced no output for %s", filename
                )
                return None

            converted_bytes = pptx_path.read_bytes()
            logger.info(
                "Converted .ppt to .pptx: %s (%d bytes)", filename, len(converted_bytes)
            )
            return converted_bytes

    except subprocess.TimeoutExpired:
        logger.warning(
            "LibreOffice conversion timed out (30s) for %s", filename
        )
        return None
    except Exception as e:
        logger.warning("LibreOffice conversion error for %s: %s", filename, e)
        return None


# ---------------------------------------------------------------------------
# Individual parsers
# ---------------------------------------------------------------------------


def _parse_pdf(data: bytes, filename: str) -> str:
    """Parse PDF using pymupdf (PyMuPDF)."""
    import pymupdf

    try:
        doc = pymupdf.open(stream=data, filetype="pdf")
    except Exception as e:
        raise ValueError(f"PDF open failed (encrypted or corrupt?): {e}") from e
    texts = []
    for page_num, page in enumerate(doc):
        text = page.get_text()
        if text.strip():
            texts.append(f"[Page {page_num + 1}]\n{text}")
    doc.close()
    return "\n\n".join(texts)


def _parse_pdf_enhanced(data: bytes, filename: str) -> ParseResult:
    """Enhanced PDF parsing with table extraction, scanned page detection, and image OCR."""
    import pymupdf

    try:
        doc = pymupdf.open(stream=data, filetype="pdf")
    except Exception as e:
        raise ValueError(f"PDF open failed (encrypted or corrupt?): {e}") from e
    texts = []
    tables: list[list[list[str]]] = []
    scanned_pages: list[int] = []
    extracted_images: list[bytes] = []

    def _is_garbled_text(text: str) -> bool:
        """Detect garbled text from broken font mapping in PDFs.

        Patterns: repeated single chars (손손손), CJK chars with no meaning,
        very low unique char ratio, etc.
        """
        if not text or len(text.strip()) < 10:
            return False
        stripped = text.strip()
        # High repetition ratio: same char appears > 30% of text
        from collections import Counter
        char_counts = Counter(stripped.replace(" ", "").replace("\n", ""))
        if char_counts:
            most_common_ratio = char_counts.most_common(1)[0][1] / max(len(stripped), 1)
            if most_common_ratio > 0.3:
                return True
        # Very low unique char ratio (garbled = few unique chars repeated)
        unique_ratio = len(set(stripped)) / max(len(stripped), 1)
        if unique_ratio < 0.05 and len(stripped) > 50:
            return True
        return False

    for page_num, page in enumerate(doc):
        # Text extraction
        text = page.get_text()
        if text.strip() and not _is_garbled_text(text):
            texts.append(f"[Page {page_num + 1}]\n{text}")
        else:
            # Scanned page OR garbled text → needs OCR
            scanned_pages.append(page_num + 1)
            if text.strip():
                logger.info("Garbled text detected on page %d of %s, routing to OCR", page_num + 1, filename)

        # Table extraction using PyMuPDF find_tables()
        page_tables = page.find_tables()
        for table in page_tables:
            table_data = table.extract()
            if table_data:
                # Convert None cells to empty strings
                cleaned = [
                    [str(cell) if cell is not None else "" for cell in row]
                    for row in table_data
                ]
                tables.append(cleaned)

        # Image extraction from PDF pages
        image_list = page.get_images(full=True)
        for img_info in image_list:
            xref = img_info[0]
            img_data = doc.extract_image(xref)
            if img_data and img_data.get("image"):
                img_bytes = img_data["image"]
                # Only process reasonably sized images (> 1KB, < 10MB)
                if 1024 < len(img_bytes) < 10_000_000:
                    extracted_images.append(img_bytes)

    doc.close()

    # Process scanned pages via OCR
    ocr_texts = []
    if scanned_pages:
        logger.info("Detected %d scanned pages in %s: %s", len(scanned_pages), filename, scanned_pages)
        # Re-open and render scanned pages as images for OCR
        doc2 = pymupdf.open(stream=data, filetype="pdf")
        for page_num in scanned_pages:
            page = doc2[page_num - 1]
            # Render page at 200 DPI for OCR
            pix = page.get_pixmap(dpi=200)
            img_bytes = pix.tobytes("png")
            extracted_images.append(img_bytes)
        doc2.close()

    # Route extracted images through OCR/CV pipeline
    visual_analyses = []
    if extracted_images:
        ocr_text, analyses = _process_images_ocr(extracted_images)
        if ocr_text:
            ocr_texts.append(ocr_text)
        visual_analyses = analyses

    return ParseResult(
        text="\n\n".join(texts),
        tables=tables,
        ocr_text="\n".join(ocr_texts),
        images_processed=len(extracted_images),
        visual_analyses=visual_analyses,
    )


def _parse_docx(data: bytes, filename: str) -> str:
    """Parse DOCX using python-docx."""
    from docx import Document

    try:
        doc = Document(io.BytesIO(data))
    except Exception as e:
        raise ValueError(f"DOCX open failed (corrupt?): {e}") from e
    texts = []

    for para in doc.paragraphs:
        if para.text.strip():
            if para.style and para.style.name.startswith("Heading"):
                level = para.style.name[-1] if para.style.name[-1].isdigit() else "1"
                texts.append(f"{'#' * int(level)} {para.text}")
            else:
                texts.append(para.text)

    for table in doc.tables:
        table_data = []
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            table_data.append(row_data)
        if table_data:
            texts.append(_table_to_markdown(table_data))

    return "\n\n".join(texts)


def _parse_pptx(data: bytes, filename: str) -> str:
    """Parse PPTX using python-pptx."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as e:
        raise ValueError(f"PPTX open failed (corrupt?): {e}") from e
    texts = []

    def _iter_shapes(shapes, _depth: int = 0):
        if _depth > 10:
            return
        for shape in shapes:
            yield shape
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from _iter_shapes(shape.shapes, _depth + 1)

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = [f"[Slide {slide_num}]"]
        for shape in _iter_shapes(slide.shapes):
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                if table_data:
                    slide_texts.append(_table_to_markdown(table_data))
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_texts.append(f"[Notes] {notes_text}")
        if len(slide_texts) > 1:
            texts.append("\n".join(slide_texts))

    return "\n\n".join(texts)


def _parse_pptx_enhanced(data: bytes, filename: str) -> ParseResult:
    """Enhanced PPTX parsing with image extraction and OCR routing."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(data))
    texts = []
    tables: list[list[list[str]]] = []
    extracted_images: list[bytes] = []

    def _iter_shapes(shapes, _depth: int = 0):
        if _depth > 10:
            return
        for shape in shapes:
            yield shape
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from _iter_shapes(shape.shapes, _depth + 1)

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = [f"[Slide {slide_num}]"]
        for shape in _iter_shapes(slide.shapes):
            # Text frames
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)

            # Tables
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                if table_data:
                    tables.append(table_data)
                    slide_texts.append(_table_to_markdown(table_data))

            # Image extraction from PICTURE shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_bytes = shape.image.blob
                if img_bytes and 1024 < len(img_bytes) < 10_000_000:
                    extracted_images.append(img_bytes)

        # Slide notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_texts.append(f"[Notes] {notes_text}")

        if len(slide_texts) > 1:
            texts.append("\n".join(slide_texts))

    # Route extracted images through OCR/CV pipeline
    visual_analyses = []
    ocr_text = ""
    if extracted_images:
        ocr_text, visual_analyses = _process_images_ocr(extracted_images)

    return ParseResult(
        text="\n\n".join(texts),
        tables=tables,
        ocr_text=ocr_text,
        images_processed=len(extracted_images),
        visual_analyses=visual_analyses,
    )


def _parse_image(data: bytes, filename: str) -> ParseResult:
    """Parse image file through OCR/CV pipeline."""
    ocr_text, visual_analyses = _process_images_ocr([data])
    return ParseResult(
        ocr_text=ocr_text,
        images_processed=1,
        visual_analyses=visual_analyses,
    )


def _parse_xlsx(data: bytes, filename: str) -> str:
    """Parse XLSX using openpyxl."""
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    texts = []

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        texts.append(f"[Sheet: {sheet_name}]")
        sheet_data = []
        for row in sheet.iter_rows(values_only=True):
            row_data = [str(cell) if cell is not None else "" for cell in row]
            if any(cell.strip() for cell in row_data):
                sheet_data.append(row_data)
        if sheet_data:
            texts.append(_table_to_markdown(sheet_data))

    wb.close()
    return "\n\n".join(texts)


def _parse_text(data: bytes) -> str:
    """Parse text files with UTF-8 / EUC-KR fallback."""
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        return data.decode("euc-kr", errors="ignore")


# ---------------------------------------------------------------------------
# OCR / CV Pipeline routing
# ---------------------------------------------------------------------------


def _process_images_ocr(
    images: list[bytes],
) -> tuple[str, list[dict[str, Any]]]:
    """Route extracted images through PaddleOCR Docker API server.

    When ``weights.ocr.enable_vision_analysis`` is True, calls ``/analyze``
    instead of ``/ocr`` to get shape/arrow detection alongside OCR text.

    Returns:
        (ocr_text, visual_analyses)
    """
    import base64
    import httpx
    import os

    ocr_texts: list[str] = []
    visual_analyses: list[dict[str, Any]] = []
    base_url = os.getenv("PADDLEOCR_API_URL", "http://localhost:8866")
    # Strip trailing path if user set full URL like "http://host:8866/ocr"
    if base_url.endswith("/ocr") or base_url.endswith("/analyze"):
        base_url = base_url.rsplit("/", 1)[0]

    vision_enabled = _w.ocr.enable_vision_analysis
    endpoint = f"{base_url}/analyze" if vision_enabled else f"{base_url}/ocr"

    if not images:
        return "", []

    client = httpx.Client(timeout=60.0)

    def _resize_image(raw: bytes, scale: float = 0.75) -> bytes | None:
        """Resize image as fallback when OCR fails on the original."""
        try:
            from PIL import Image
            import io
            img = Image.open(io.BytesIO(raw))
            new_w = max(int(img.width * scale), 32)
            new_h = max(int(img.height * scale), 32)
            img = img.resize((new_w, new_h), Image.LANCZOS)
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            return buf.getvalue()
        except Exception:
            return None

    def _ocr_request(payload: bytes) -> dict | None:
        """Send OCR request with one retry using resized image on failure."""
        b64_image = base64.b64encode(payload).decode("utf-8")
        try:
            resp = client.post(endpoint, json={"image": b64_image}, timeout=60.0)
            resp.raise_for_status()
            return resp.json()
        except Exception:
            return None

    # Pre-filter: skip images too small to contain meaningful text
    MIN_W, MIN_H = 20, 20

    for i, img_bytes in enumerate(images):
        try:
            from PIL import Image
            import io
            _img = Image.open(io.BytesIO(img_bytes))
            if _img.width < MIN_W or _img.height < MIN_H:
                logger.debug("Skipping image %d: too small (%dx%d)", i + 1, _img.width, _img.height)
                continue

            result = _ocr_request(img_bytes)

            # Retry with resized image on failure
            if result is None:
                resized = _resize_image(img_bytes)
                if resized:
                    logger.info("Retrying image %d with resized version", i + 1)
                    result = _ocr_request(resized)

            if result is None:
                logger.warning("PaddleOCR failed for image %d after retry, skipping", i + 1)
                continue

            # Parse OCR text (same keys for both /ocr and /analyze)
            text_lines = result.get("texts", result.get("result", []))
            if isinstance(text_lines, list):
                text = " ".join(str(t) for t in text_lines if t)
            elif isinstance(text_lines, str):
                text = text_lines
            else:
                text = str(text_lines)

            if text.strip():
                ocr_texts.append(f"[Image {i + 1} OCR] {text}")
                logger.info("OCR success for image %d: %d chars", i + 1, len(text))

            # When vision analysis is enabled, capture shapes/arrows/mappings
            if vision_enabled:
                shapes = result.get("shapes", [])
                arrows = result.get("arrows", [])
                mappings = result.get("text_shape_mappings", [])
                if shapes or arrows or mappings:
                    analysis: dict[str, Any] = {
                        "image_index": i + 1,
                        "shapes": shapes,
                        "arrows": arrows,
                        "text_shape_mappings": mappings,
                        "shape_count": result.get("shape_count", len(shapes)),
                        "arrow_count": result.get("arrow_count", len(arrows)),
                        "ocr_confidence": result.get("ocr_confidence", 0.0),
                    }
                    visual_analyses.append(analysis)

                    # Build a structured textual description of the diagram
                    desc_parts: list[str] = []
                    if shapes:
                        shape_types = [s.get("type", "unknown") for s in shapes]
                        desc_parts.append(f"Shapes: {', '.join(shape_types)}")
                    if arrows:
                        desc_parts.append(f"Arrows: {len(arrows)} connections")
                    if mappings:
                        for m in mappings:
                            shape_label = m.get("shape_type", "shape")
                            texts_in = m.get("texts", m.get("text", ""))
                            if texts_in:
                                desc_parts.append(f"  [{shape_label}] {texts_in}")
                    if desc_parts:
                        diagram_desc = "\n".join(desc_parts)
                        ocr_texts.append(f"[Image {i + 1} Diagram]\n{diagram_desc}")

        except Exception as e:
            logger.warning("PaddleOCR API failed for image %d: %s", i + 1, e)

    client.close()
    return "\n".join(ocr_texts), visual_analyses


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _table_to_markdown(data: list[list[str]], max_rows: int = 50) -> str:
    """Convert table data to markdown format."""
    if not data:
        return ""
    lines = []
    header = data[0]
    lines.append("| " + " | ".join(header) + " |")
    lines.append("| " + " | ".join(["---"] * len(header)) + " |")
    for row in data[1:max_rows]:
        padded_row = row + [""] * (len(header) - len(row))
        padded_row = padded_row[: len(header)]
        lines.append("| " + " | ".join(padded_row) + " |")
    if len(data) > max_rows:
        lines.append(f"... ({len(data) - max_rows} rows omitted)")
    return "\n".join(lines)
