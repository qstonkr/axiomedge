"""PPT/PPTX parsing mixin for AttachmentParser.

Handles PowerPoint text extraction: OLE2 legacy .ppt, .pptx slide
content extraction, and orchestrates OCR stages via _PptOcrMixin.
"""

from __future__ import annotations

import logging
from pathlib import Path

from .models import AttachmentParseResult
from ._attachment_helpers import (
    _get_ocr_feature_flags,
    _parse_ppt_ole_records,
    _should_ocr_ppt,
    _text_chars,
    _try_catppt_extract,
    _try_libreoffice_ppt_convert,
)
from ._ppt_ocr import _PptOcrMixin

logger = logging.getLogger(__name__)


class _PptParserMixin(_PptOcrMixin):
    """PPT/PPTX parsing methods for AttachmentParser."""

    # -----------------------------------------------------------------
    # PPT OLE2 helpers
    # -----------------------------------------------------------------

    @staticmethod
    def _extract_ppt_olefile(file_path: Path) -> str | None:
        """Pure Python .ppt 텍스트 추출 (OLE2 PowerPoint 레코드 파싱)"""
        try:
            import olefile
        except ImportError:
            return None

        import struct

        try:
            ole = olefile.OleFileIO(str(file_path))
        except (
            RuntimeError, OSError, ValueError,
            TypeError, KeyError, AttributeError,
        ):
            return None

        try:
            if not ole.exists("PowerPoint Document"):
                return None

            raw = ole.openstream("PowerPoint Document").read()
            text_parts = _parse_ppt_ole_records(raw, struct)
            return "\n\n".join(text_parts) if text_parts else None
        finally:
            ole.close()

    @classmethod
    def _parse_legacy_ppt(
        cls, file_path: Path, heartbeat_fn=None,
    ) -> AttachmentParseResult:
        """레거시 .ppt (OLE2) 파일에서 텍스트 추출

        Strategy:
        1. LibreOffice headless로 .ppt -> .pptx 변환
        2. Fallback: catppt -> olefile 텍스트 추출
        """
        result = _try_libreoffice_ppt_convert(file_path, heartbeat_fn)
        if result is not None:
            return result

        result = _try_catppt_extract(file_path)
        if result is not None:
            return result

        # 3차: olefile pure Python 추출 (로컬 환경)
        text = cls._extract_ppt_olefile(file_path)
        if text:
            return AttachmentParseResult(
                extracted_text=text,
                extracted_tables=[],
                confidence=0.5,
                native_text_chars=_text_chars(text),
            )

        return AttachmentParseResult(
            extracted_text=(
                "[.ppt 파싱 실패: "
                "LibreOffice/catppt/olefile 모두 실패]"
            ),
            extracted_tables=[],
            confidence=0.0,
            ocr_mode=cls.current_policy().attachment_ocr_mode,
            ocr_skip_reason="parse_error",
        )

    # -----------------------------------------------------------------
    # PPT (.pptx) slide content extraction
    # -----------------------------------------------------------------

    @classmethod
    def _iter_shapes(cls, shapes):
        """GroupShape 포함 재귀 shape 탐색."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        for shape in shapes:
            yield shape
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from cls._iter_shapes(shape.shapes)

    @classmethod
    def _extract_ppt_slide_content(cls, slide, slide_num: int):
        """Extract text, tables, and image shapes from a slide.

        Returns (slide_texts, tables, image_shapes).
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        slide_texts: list[str] = []
        tables: list[dict] = []
        image_shapes: list[tuple[int, bytes]] = []

        for shape in cls._iter_shapes(slide.shapes):
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)

            if shape.has_table:
                tbl = cls._extract_pptx_table(
                    shape.table, slide_num,
                )
                if tbl is not None:
                    tables.append(tbl)

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                cls._collect_image_shape(
                    shape, slide_num, image_shapes,
                )

        # 슬라이드 노트 추출
        if slide.has_notes_slide:
            notes_text = (
                slide.notes_slide.notes_text_frame.text.strip()
            )
            if notes_text:
                slide_texts.append(f"[Notes] {notes_text}")

        return slide_texts, tables, image_shapes

    @staticmethod
    def _extract_pptx_table(table, slide_num: int) -> dict | None:
        """Extract a single table from a PPTX shape."""
        rows_data = []
        for row in table.rows:
            row_values = [cell.text.strip() for cell in row.cells]
            rows_data.append(row_values)

        if not rows_data:
            return None

        headers = rows_data[0]
        data_rows = rows_data[1:] if len(rows_data) > 1 else []
        return {
            "slide": slide_num,
            "headers": headers,
            "rows": [
                dict(zip(headers, row))
                for row in data_rows
                if len(row) == len(headers)
            ],
        }

    @staticmethod
    def _collect_image_shape(
        shape, slide_num: int, image_shapes: list,
    ):
        """Collect image bytes from a picture shape."""
        try:
            image_bytes = shape.image.blob
            if len(image_bytes) > 10_000:  # 10KB 이상만
                image_shapes.append((slide_num, image_bytes))
        except (
            RuntimeError, OSError, ValueError,
            TypeError, KeyError, AttributeError,
        ):
            pass

    # -----------------------------------------------------------------
    # Result builder + skip reason
    # -----------------------------------------------------------------

    @classmethod
    def _determine_ppt_skip_reason(
        cls, policy, should_ocr: bool, ocr_units_deferred: int,
    ) -> str | None:
        """Determine the OCR skip reason for the PPT result."""
        if policy.attachment_ocr_mode == "off":
            return "disabled"
        if policy.attachment_ocr_mode == "auto" and not should_ocr:
            return "native_text_sufficient"
        if ocr_units_deferred > 0:
            return "budget_exceeded"
        return None

    @classmethod
    def _build_ppt_result(
        cls, full_text, tables, policy, should_ocr,
        ocr_units_attempted, ocr_units_extracted,
        ocr_units_deferred, native_text_chars, ocr_text_chars,
    ) -> AttachmentParseResult:
        """Build the final AttachmentParseResult for PPT."""
        return AttachmentParseResult(
            extracted_text=full_text,
            extracted_tables=tables,
            confidence=0.85 if full_text.strip() else 0.0,
            ocr_mode=policy.attachment_ocr_mode,
            ocr_applied=ocr_units_extracted > 0,
            ocr_skip_reason=cls._determine_ppt_skip_reason(
                policy, should_ocr, ocr_units_deferred,
            ),
            ocr_units_attempted=ocr_units_attempted,
            ocr_units_extracted=ocr_units_extracted,
            ocr_units_deferred=ocr_units_deferred,
            native_text_chars=native_text_chars,
            ocr_text_chars=ocr_text_chars,
        )

    # -----------------------------------------------------------------
    # Main parse_ppt entry point
    # -----------------------------------------------------------------

    @classmethod
    def parse_ppt(
        cls, file_path: Path, heartbeat_fn=None,
    ) -> AttachmentParseResult:
        """PPT에서 슬라이드 텍스트 추출

        .pptx: python-pptx, .ppt: catppt
        """
        try:
            if str(file_path).lower().endswith(".ppt"):
                return cls._parse_legacy_ppt(
                    file_path, heartbeat_fn=heartbeat_fn,
                )
            from pptx import Presentation

            policy = cls.current_policy()
            prs = Presentation(file_path)
            text_parts: list[str] = []
            tables: list[dict] = []
            image_shapes: list[tuple[int, bytes]] = []

            # Stage 1: Extract text, tables, images
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts, slide_tables, slide_images = (
                    cls._extract_ppt_slide_content(
                        slide, slide_num,
                    )
                )
                tables.extend(slide_tables)
                image_shapes.extend(slide_images)
                if slide_texts:
                    text_parts.append(
                        f"[Slide {slide_num}]\n"
                        + "\n".join(slide_texts),
                    )

            native_text = "\n\n".join(text_parts)
            native_text_chars = cls._text_chars(native_text)
            cls._emit_status(
                heartbeat_fn,
                f"native_extract ppt slides={len(prs.slides)}"
                f" chars={native_text_chars}",
            )

            should_ocr = _should_ocr_ppt(policy, native_text_chars)
            ocr_preprocess, ocr_postprocess = (
                _get_ocr_feature_flags()
            )

            ocr_units_attempted = 0
            ocr_units_extracted = 0
            ocr_units_deferred = 0
            ocr_text_chars = 0
            extracted_slides: set[int] = set()

            # Stage 2: Slide rendering OCR
            if (
                should_ocr
                and policy.slide_render_enabled
                and file_path
            ):
                (
                    slide_rendered, render_texts,
                    ocr_units_attempted, ocr_units_extracted,
                    ocr_units_deferred, ocr_text_chars,
                    extracted_slides,
                ) = cls._render_and_ocr_slides(
                    file_path, policy, heartbeat_fn,
                    ocr_preprocess, ocr_postprocess,
                )
                text_parts.extend(render_texts)
            else:
                slide_rendered = False

            # Stage 3: Shape-by-shape OCR fallback
            if should_ocr and not slide_rendered:
                (
                    shape_texts, shape_attempted,
                    shape_extracted, shape_deferred,
                    shape_chars,
                ) = cls._shape_ocr_pass(
                    image_shapes, policy, len(prs.slides),
                    heartbeat_fn, ocr_preprocess,
                    ocr_postprocess, extracted_slides,
                )
                text_parts.extend(shape_texts)
                ocr_units_attempted += shape_attempted
                ocr_units_extracted += shape_extracted
                ocr_units_deferred += shape_deferred
                ocr_text_chars += shape_chars

            full_text = "\n\n".join(text_parts)

            # Stage 4: PDF fallback for empty results
            full_text, tables, ocr_text_chars = (
                cls._apply_pdf_fallback_if_needed(
                    should_ocr, full_text, tables,
                    ocr_text_chars, file_path, heartbeat_fn,
                )
            )

            if ocr_units_deferred > 0:
                cls._emit_status(
                    heartbeat_fn,
                    f"ocr_skipped_budget ppt "
                    f"deferred={ocr_units_deferred}",
                )

            return cls._build_ppt_result(
                full_text, tables, policy, should_ocr,
                ocr_units_attempted, ocr_units_extracted,
                ocr_units_deferred, native_text_chars,
                ocr_text_chars,
            )

        except (
            RuntimeError, OSError, ValueError,
            TypeError, KeyError, AttributeError,
        ) as e:
            return AttachmentParseResult(
                extracted_text=f"[PPT 파싱 오류: {e}]",
                extracted_tables=[],
                confidence=0.0,
                ocr_mode=cls.current_policy().attachment_ocr_mode,
                ocr_skip_reason="parse_error",
            )
