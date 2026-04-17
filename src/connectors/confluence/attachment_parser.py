"""Attachment content extraction for Confluence crawled files.

Handles PDF, Excel, Word (.doc/.docx), PowerPoint (.ppt/.pptx), and image files.
Uses PaddleOCR for image-based content extraction with subprocess isolation to
defend against SIGSEGV crashes in PaddleOCR's C++ inference engine.

Originally extracted from ``scripts/confluence_full_crawler.py``.
"""

from __future__ import annotations

import asyncio
import io
import logging
import os
from pathlib import Path
from typing import Any

from src.config.weights import weights as _w
from .models import AttachmentOCRPolicy, AttachmentParseResult

logger = logging.getLogger(__name__)

from ._attachment_helpers import (  # noqa: E402
    _DEFAULT_ATTACHMENT_OCR_MODE,
    _DEFAULT_OCR_MAX_IMAGES_PER_ATTACHMENT,
    _DEFAULT_OCR_MAX_PDF_PAGES,
    _DEFAULT_OCR_MAX_PPT_SLIDES,
    _DEFAULT_OCR_MIN_TEXT_CHARS,
    _SOURCE_ATTACHMENT_OCR_DEFAULTS,
    _apply_ocr_postprocess,
    _decode_ole_text,
    _downscale_image,
    _env_bool,
    OCR_MAX_ASPECT_RATIO,
    OCR_MIN_DIMENSION,
    _env_int,  # noqa: F401 — re-exported for test backward compat
    _filter_ocr_noise,
    _get_ocr_feature_flags,
    _get_ocr_postprocess_flag,
    _image_result_no_ocr,
    _ocr_worker_fn,
    _pad_extreme_aspect_ratio,
    _parse_ppt_ole_records,
    _postprocess_slide_text,
    _preprocess_shape_image,
    _preprocess_slide_image,
    _resolve_bool_field,
    _resolve_int_field,
    _resolve_ocr_mode,
    _should_ocr_ppt,
    _try_catppt_extract,
    _try_cli_doc_extract,
    _try_layout_ocr,
    _try_libreoffice_ppt_convert,
    _try_slide_layout_ocr,
)

# =============================================================================
# Attachment Parsers
# =============================================================================
class AttachmentParser:
    """첨부파일 내용 추출기"""

    _active_source_key = ""
    _ocr_policy = AttachmentOCRPolicy(
        attachment_ocr_mode=_DEFAULT_ATTACHMENT_OCR_MODE,
        ocr_min_text_chars=_DEFAULT_OCR_MIN_TEXT_CHARS,
        ocr_max_pdf_pages=_DEFAULT_OCR_MAX_PDF_PAGES,
        ocr_max_ppt_slides=_DEFAULT_OCR_MAX_PPT_SLIDES,
        ocr_max_images_per_attachment=_DEFAULT_OCR_MAX_IMAGES_PER_ATTACHMENT,
        slide_render_enabled=_env_bool("KNOWLEDGE_SLIDE_RENDER_ENABLED", True),
        layout_analysis_enabled=_env_bool("KNOWLEDGE_LAYOUT_ANALYSIS_ENABLED", True),
    )

    @classmethod
    def configure_run(
        cls, source_key: str, overrides: dict[str, Any] | None = None,
    ) -> AttachmentOCRPolicy:
        """Resolve run-local OCR policy once per crawl source."""
        overrides = overrides or {}
        source_defaults = _SOURCE_ATTACHMENT_OCR_DEFAULTS.get(source_key, {})
        legacy_slide_render = _env_bool("KNOWLEDGE_SLIDE_RENDER_ENABLED", True)
        legacy_layout_analysis = _env_bool("KNOWLEDGE_LAYOUT_ANALYSIS_ENABLED", True)

        cls._active_source_key = source_key
        cls._ocr_policy = AttachmentOCRPolicy(
            attachment_ocr_mode=_resolve_ocr_mode(overrides, source_defaults),
            ocr_min_text_chars=_resolve_int_field(
                overrides, source_defaults,
                "ocr_min_text_chars", "KNOWLEDGE_CRAWL_OCR_MIN_TEXT_CHARS",
                _DEFAULT_OCR_MIN_TEXT_CHARS,
            ),
            ocr_max_pdf_pages=_resolve_int_field(
                overrides, source_defaults,
                "ocr_max_pdf_pages", "KNOWLEDGE_CRAWL_OCR_MAX_PDF_PAGES",
                _DEFAULT_OCR_MAX_PDF_PAGES,
            ),
            ocr_max_ppt_slides=_resolve_int_field(
                overrides, source_defaults,
                "ocr_max_ppt_slides", "KNOWLEDGE_CRAWL_OCR_MAX_PPT_SLIDES",
                _DEFAULT_OCR_MAX_PPT_SLIDES,
            ),
            ocr_max_images_per_attachment=_resolve_int_field(
                overrides, source_defaults,
                "ocr_max_images_per_attachment",
                "KNOWLEDGE_CRAWL_OCR_MAX_IMAGES_PER_ATTACHMENT",
                _DEFAULT_OCR_MAX_IMAGES_PER_ATTACHMENT,
            ),
            slide_render_enabled=_resolve_bool_field(
                overrides, source_defaults,
                "slide_render_enabled", "KNOWLEDGE_CRAWL_SLIDE_RENDER_ENABLED",
                legacy_slide_render,
            ),
            layout_analysis_enabled=_resolve_bool_field(
                overrides, source_defaults,
                "layout_analysis_enabled", "KNOWLEDGE_CRAWL_LAYOUT_ANALYSIS_ENABLED",
                legacy_layout_analysis,
            ),
        )
        return cls._ocr_policy

    @classmethod
    def current_policy(cls) -> AttachmentOCRPolicy:
        return cls._ocr_policy

    @staticmethod
    def _emit_status(heartbeat_fn, message: str) -> None:
        if heartbeat_fn:
            heartbeat_fn(message)
        else:
            logger.debug("[status] %s", message)

    @staticmethod
    def _text_chars(value: str | None) -> int:
        return len(value.strip()) if value and value.strip() else 0

    # =================================================================
    # PDF helpers
    # =================================================================

    @classmethod
    def _ocr_pdf_page(cls, page, page_num: int, total_pages: int, policy, heartbeat_fn):
        """OCR a single textless PDF page via image rendering.

        Returns (ocr_text, success_bool).
        """
        import fitz

        cls._emit_status(
            heartbeat_fn,
            f"ocr_processing pdf page={page_num}/{total_pages}",
        )
        zoom = 2.0  # 144 DPI (72 * 2)
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat)
        png_bytes = pix.tobytes("png")
        del pix  # Free native memory immediately

        postprocess = _get_ocr_postprocess_flag()

        ocr_text = cls._ocr_slide_image(
            png_bytes, page_num,
            preprocess=True,
            layout_analysis=policy.layout_analysis_enabled,
            postprocess=postprocess,
        )
        return ocr_text

    @staticmethod
    def _extract_pdf_page_tables(page, page_num: int) -> list[dict]:
        """Extract tables from a single PDF page using PyMuPDF."""
        tables = []
        try:
            page_tables = page.find_tables()
            for table in page_tables:
                table_data = table.extract()
                if table_data and len(table_data) > 1:
                    headers = table_data[0]
                    rows = table_data[1:]
                    tables.append({
                        "page": page_num,
                        "headers": headers,
                        "rows": [
                            dict(zip(headers, row))
                            for row in rows
                            if len(row) == len(headers)
                        ],
                    })
        except (RuntimeError, OSError, ValueError, TypeError, KeyError, AttributeError):
            pass  # 테이블 추출 실패 시 건너뛰기
        return tables

    @staticmethod
    def _compute_pdf_confidence(
        has_text: bool, ocr_units_extracted: int,
    ) -> float:
        """Compute PDF parse confidence score."""
        if has_text and ocr_units_extracted == 0:
            return 0.9
        if has_text:
            return 0.7
        return 0.0

    @classmethod
    def parse_pdf(cls, file_path: Path, heartbeat_fn=None) -> AttachmentParseResult:
        """PDF에서 텍스트와 테이블 추출

        Strategy:
        1. PyMuPDF 텍스트 레이어 추출 시도
        2. 텍스트가 빈 페이지 → 이미지 렌더링 → PaddleOCR fallback
           (이미지 기반 PDF — PPT를 PDF로 내보낸 파일 등)
        """
        try:
            policy = cls.current_policy()
            import fitz  # PyMuPDF  # noqa: F811

            doc = fitz.open(file_path)
            text_parts: list[str] = []
            tables: list[dict] = []
            native_text_chars = 0
            total_pages = len(doc)
            textless_pages = 0
            ocr_counters = {
                "attempted": 0, "extracted": 0, "deferred": 0, "chars": 0,
            }

            cls._emit_status(heartbeat_fn, f"native_extract pdf pages={total_pages}")

            for page_num, page in enumerate(doc, 1):
                page_text = page.get_text("text")
                if page_text.strip():
                    text_parts.append(f"[Page {page_num}]\n{page_text}")
                    native_text_chars += cls._text_chars(page_text)
                else:
                    textless_pages += 1
                    cls._process_textless_pdf_page(
                        page, page_num, total_pages, policy, heartbeat_fn,
                        text_parts, ocr_counters,
                    )

                tables.extend(cls._extract_pdf_page_tables(page, page_num))

                if heartbeat_fn and page_num % 5 == 0:
                    heartbeat_fn(f"pdf_ocr: {page_num}/{total_pages}")

            doc.close()

            full_text = "\n\n".join(text_parts)
            if ocr_counters["extracted"] > 0:
                logger.info(
                    "[PDF OCR] %d/%d pages used OCR fallback",
                    ocr_counters["extracted"], total_pages,
                )
            if ocr_counters["deferred"] > 0:
                cls._emit_status(
                    heartbeat_fn,
                    f"ocr_skipped_budget pdf deferred={ocr_counters['deferred']}",
                )

            confidence = cls._compute_pdf_confidence(
                bool(full_text.strip()), ocr_counters["extracted"],
            )

            skip_reason = None
            if textless_pages > 0 and policy.attachment_ocr_mode == "off":
                skip_reason = "disabled"
            elif ocr_counters["deferred"] > 0:
                skip_reason = "budget_exceeded"

            return AttachmentParseResult(
                extracted_text=full_text,
                extracted_tables=tables,
                confidence=confidence,
                ocr_mode=policy.attachment_ocr_mode,
                ocr_applied=ocr_counters["extracted"] > 0,
                ocr_skip_reason=skip_reason,
                ocr_units_attempted=ocr_counters["attempted"],
                ocr_units_extracted=ocr_counters["extracted"],
                ocr_units_deferred=ocr_counters["deferred"],
                native_text_chars=native_text_chars,
                ocr_text_chars=ocr_counters["chars"],
            )

        except (RuntimeError, OSError, ValueError, TypeError, KeyError, AttributeError) as e:
            return AttachmentParseResult(
                extracted_text=f"[PDF 파싱 오류: {e}]",
                extracted_tables=[],
                confidence=0.0,
                ocr_mode=cls.current_policy().attachment_ocr_mode,
                ocr_skip_reason="parse_error",
            )

    @classmethod
    def _process_textless_pdf_page(
        cls, page, page_num, total_pages, policy, heartbeat_fn,
        text_parts, ocr_counters,
    ):
        """Handle a textless PDF page — apply OCR if allowed by policy.

        Note: This mutates text_parts and ocr_counters in place.
        """
        if policy.attachment_ocr_mode == "off":
            return
        if ocr_counters["attempted"] >= policy.ocr_max_pdf_pages:
            ocr_counters["deferred"] += 1
            return
        try:
            ocr_counters["attempted"] += 1
            ocr_text = cls._ocr_pdf_page(
                page, page_num, total_pages, policy, heartbeat_fn,
            )
            if ocr_text and ocr_text.strip():
                text_parts.append(f"[Page {page_num}]\n{ocr_text}")
                ocr_counters["extracted"] += 1
                ocr_counters["chars"] += cls._text_chars(ocr_text)
        except (RuntimeError, OSError, ValueError, TypeError, KeyError, AttributeError) as ocr_err:
            logger.warning("[PDF OCR] Page %d OCR failed: %s", page_num, ocr_err)

    # =================================================================
    # Excel helpers
    # =================================================================

    @staticmethod
    def _process_excel_sheet(sheet, sheet_name: str) -> tuple[list[dict], list[str]]:
        """Extract table data and text lines from a single Excel sheet.

        Returns (tables, text_parts).
        """
        tables: list[dict] = []
        text_parts: list[str] = []

        rows_data = []
        for row in sheet.iter_rows(values_only=True):
            row_values = [str(cell) if cell is not None else "" for cell in row]
            if any(v.strip() for v in row_values):
                rows_data.append(row_values)

        if not rows_data:
            return tables, text_parts

        headers = rows_data[0]
        data_rows = rows_data[1:] if len(rows_data) > 1 else []

        tables.append({
            "sheet": sheet_name,
            "headers": headers,
            "rows": [
                dict(zip(headers, row))
                for row in data_rows
                if len(row) == len(headers)
            ],
            "row_count": len(data_rows),
        })

        text_parts.append(f"[Sheet: {sheet_name}]")
        text_parts.append(" | ".join(headers))
        for row in data_rows[:10]:
            text_parts.append(" | ".join(row))
        if len(data_rows) > 10:
            text_parts.append(f"... 외 {len(data_rows) - 10}행")

        return tables, text_parts

    @staticmethod
    def parse_excel(file_path: Path) -> AttachmentParseResult:
        """Excel에서 시트 데이터 추출"""
        try:
            from openpyxl import load_workbook

            wb = load_workbook(file_path, read_only=True, data_only=True)
            text_parts: list[str] = []
            tables: list[dict] = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_tables, sheet_texts = AttachmentParser._process_excel_sheet(
                    sheet, sheet_name,
                )
                tables.extend(sheet_tables)
                text_parts.extend(sheet_texts)

            wb.close()

            full_text = "\n".join(text_parts)
            confidence = 0.95 if tables else 0.0

            return AttachmentParseResult(
                extracted_text=full_text,
                extracted_tables=tables,
                confidence=confidence,
                native_text_chars=AttachmentParser._text_chars(full_text),
            )

        except (RuntimeError, OSError, ValueError, TypeError, KeyError, AttributeError) as e:
            return AttachmentParseResult(
                extracted_text=f"[Excel 파싱 오류: {e}]",
                extracted_tables=[],
                confidence=0.0,
                ocr_skip_reason="parse_error",
            )

    # =================================================================
    # Word helpers
    # =================================================================

    @staticmethod
    def _extract_doc_olefile(file_path: Path) -> str | None:
        """Pure Python .doc 텍스트 추출 (olefile OLE2 스트림 파싱)"""
        try:
            import olefile
        except ImportError:
            return None

        try:
            ole = olefile.OleFileIO(str(file_path))
        except (RuntimeError, OSError, ValueError, TypeError, KeyError, AttributeError):
            return None

        try:
            if not ole.exists("WordDocument"):
                return None

            raw = ole.openstream("WordDocument").read()
            if len(raw) < 20:
                return None

            return _decode_ole_text(raw)
        finally:
            ole.close()

    @staticmethod
    def _parse_legacy_doc(file_path: Path) -> AttachmentParseResult:
        """레거시 .doc (OLE2) 파일에서 텍스트 추출 (antiword → catdoc → olefile fallback)"""
        import shutil

        # 1차: antiword (테이블 구조 보존, Docker 환경)
        result = _try_cli_doc_extract(shutil.which("antiword"), file_path, confidence=0.7)
        if result is not None:
            return result

        # 2차: catdoc (antiword 없을 때, Docker 환경)
        result = _try_cli_doc_extract(
            shutil.which("catdoc"), file_path, confidence=0.6, extra_args=["-w"],
        )
        if result is not None:
            return result

        # 3차: olefile pure Python 추출 (로컬 환경, 시스템 도구 없을 때)
        text = AttachmentParser._extract_doc_olefile(file_path)
        if text:
            return AttachmentParseResult(
                extracted_text=text,
                extracted_tables=[],
                confidence=0.5,
                native_text_chars=AttachmentParser._text_chars(text),
            )

        return AttachmentParseResult(
            extracted_text="[.doc 파싱 실패: antiword/catdoc 미설치 및 olefile 추출 실패]",
            extracted_tables=[],
            confidence=0.0,
            ocr_skip_reason="parse_error",
        )

    @staticmethod
    def _extract_word_tables(doc) -> list[dict]:
        """Extract tables from a python-docx Document object."""
        tables = []
        for idx, table in enumerate(doc.tables, 1):
            rows_data = []
            for row in table.rows:
                row_values = [cell.text.strip() for cell in row.cells]
                rows_data.append(row_values)

            if rows_data:
                headers = rows_data[0]
                data_rows = rows_data[1:] if len(rows_data) > 1 else []
                tables.append({
                    "table_index": idx,
                    "headers": headers,
                    "rows": [
                        dict(zip(headers, row))
                        for row in data_rows
                        if len(row) == len(headers)
                    ],
                })
        return tables

    @staticmethod
    def parse_word(file_path: Path) -> AttachmentParseResult:
        """Word에서 텍스트와 테이블 추출 (.docx: python-docx, .doc: antiword/catdoc)"""
        try:
            if str(file_path).lower().endswith(".doc"):
                return AttachmentParser._parse_legacy_doc(file_path)
            from docx import Document

            doc = Document(file_path)
            text_parts = [p.text for p in doc.paragraphs if p.text.strip()]
            tables = AttachmentParser._extract_word_tables(doc)

            full_text = "\n\n".join(text_parts)
            confidence = 0.9 if full_text.strip() else 0.0

            return AttachmentParseResult(
                extracted_text=full_text,
                extracted_tables=tables,
                confidence=confidence,
                native_text_chars=AttachmentParser._text_chars(full_text),
            )

        except (RuntimeError, OSError, ValueError, TypeError, KeyError, AttributeError) as e:
            return AttachmentParseResult(
                extracted_text=f"[Word 파싱 오류: {e}]",
                extracted_tables=[],
                confidence=0.0,
                ocr_skip_reason="parse_error",
            )

    # =================================================================
    # PPT OLE2 helpers
    # =================================================================

    @staticmethod
    def _extract_ppt_olefile(file_path: Path) -> str | None:
        """Pure Python .ppt 텍스트 추출 (OLE2 PowerPoint 레코드 파싱)"""
        try:
            import olefile
        except ImportError:
            return None

        import struct

        try:
            ole = olefile.OleFileIO(str(file_path))
        except (RuntimeError, OSError, ValueError, TypeError, KeyError, AttributeError):
            return None

        try:
            if not ole.exists("PowerPoint Document"):
                return None

            raw = ole.openstream("PowerPoint Document").read()
            text_parts = _parse_ppt_ole_records(raw, struct)
            return "\n\n".join(text_parts) if text_parts else None
        finally:
            ole.close()

    @staticmethod
    def _parse_legacy_ppt(file_path: Path, heartbeat_fn=None) -> AttachmentParseResult:
        """레거시 .ppt (OLE2) 파일에서 텍스트 추출

        Strategy:
        1. LibreOffice headless로 .ppt → .pptx 변환 → parse_ppt 재사용 (OCR 포함)
        2. Fallback: catppt → olefile 텍스트 추출
        """
        result = _try_libreoffice_ppt_convert(file_path, heartbeat_fn)
        if result is not None:
            return result

        result = _try_catppt_extract(file_path)
        if result is not None:
            return result

        # 3차: olefile pure Python 추출 (로컬 환경)
        text = AttachmentParser._extract_ppt_olefile(file_path)
        if text:
            return AttachmentParseResult(
                extracted_text=text,
                extracted_tables=[],
                confidence=0.5,
                native_text_chars=AttachmentParser._text_chars(text),
            )

        return AttachmentParseResult(
            extracted_text="[.ppt 파싱 실패: LibreOffice/catppt/olefile 모두 실패]",
            extracted_tables=[],
            confidence=0.0,
            ocr_mode=AttachmentParser.current_policy().attachment_ocr_mode,
            ocr_skip_reason="parse_error",
        )

    # =================================================================
    # PPT (.pptx) parsing — split into stages
    # =================================================================

    @staticmethod
    def _iter_shapes(shapes):
        """GroupShape 포함 재귀 shape 탐색."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        for shape in shapes:
            yield shape
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from AttachmentParser._iter_shapes(shape.shapes)

    @classmethod
    def _extract_ppt_slide_content(cls, slide, slide_num: int):
        """Extract text, tables, and image shapes from a single slide.

        Returns (slide_texts, tables, image_shapes).
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        slide_texts: list[str] = []
        tables: list[dict] = []
        image_shapes: list[tuple[int, bytes]] = []

        for shape in cls._iter_shapes(slide.shapes):
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)

            if shape.has_table:
                tbl = cls._extract_pptx_table(shape.table, slide_num)
                if tbl is not None:
                    tables.append(tbl)

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                cls._collect_image_shape(shape, slide_num, image_shapes)

        # 슬라이드 노트 추출
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_texts.append(f"[Notes] {notes_text}")

        return slide_texts, tables, image_shapes

    @staticmethod
    def _extract_pptx_table(table, slide_num: int) -> dict | None:
        """Extract a single table from a PPTX shape."""
        rows_data = []
        for row in table.rows:
            row_values = [cell.text.strip() for cell in row.cells]
            rows_data.append(row_values)

        if not rows_data:
            return None

        headers = rows_data[0]
        data_rows = rows_data[1:] if len(rows_data) > 1 else []
        return {
            "slide": slide_num,
            "headers": headers,
            "rows": [
                dict(zip(headers, row))
                for row in data_rows
                if len(row) == len(headers)
            ],
        }

    @staticmethod
    def _collect_image_shape(shape, slide_num: int, image_shapes: list):
        """Collect image bytes from a picture shape if large enough."""
        try:
            image_bytes = shape.image.blob
            if len(image_bytes) > 10_000:  # 10KB 이상만
                image_shapes.append((slide_num, image_bytes))
        except (RuntimeError, OSError, ValueError, TypeError, KeyError, AttributeError):
            pass

    @classmethod
    def _render_and_ocr_slides(
        cls, file_path: Path, policy, heartbeat_fn,
        ocr_preprocess: bool, ocr_postprocess: bool,
    ) -> tuple[bool, list[str], int, int, int, int, set[int]]:
        """Render PPTX slides via LibreOffice and OCR each image.

        Returns (slide_rendered, text_parts, attempted, extracted,
                 deferred, ocr_chars, extracted_slides).
        """
        text_parts: list[str] = []
        ocr_units_attempted = 0
        ocr_units_extracted = 0
        ocr_units_deferred = 0
        ocr_text_chars = 0
        extracted_slides: set[int] = set()

        try:
            from scripts.slide_renderer import render_slides_as_images

            rendered_slides = render_slides_as_images(Path(str(file_path)))
            if not rendered_slides:
                return False, text_parts, 0, 0, 0, 0, extracted_slides

            logger.info(
                "[OCR] Slide rendering: %d slides from %s",
                len(rendered_slides), file_path,
            )
            for slide_num, png_bytes in rendered_slides:
                if ocr_units_attempted >= policy.ocr_max_ppt_slides:
                    ocr_units_deferred += 1
                    continue
                ocr_units_attempted += 1
                cls._emit_status(
                    heartbeat_fn,
                    f"ocr_processing ppt slide={slide_num}/{len(rendered_slides)}",
                )
                ocr_text = cls._ocr_slide_image(
                    png_bytes, slide_num,
                    preprocess=ocr_preprocess,
                    layout_analysis=policy.layout_analysis_enabled,
                    postprocess=ocr_postprocess,
                )
                if ocr_text:
                    text_parts.append(f"[Slide {slide_num} OCR]\n{ocr_text}")
                    ocr_units_extracted += 1
                    ocr_text_chars += cls._text_chars(ocr_text)
                    extracted_slides.add(slide_num)
                if heartbeat_fn and slide_num % 5 == 0:
                    heartbeat_fn(
                        f"slide_render_ocr: {slide_num}/{len(rendered_slides)}",
                    )

            return (
                True, text_parts, ocr_units_attempted,
                ocr_units_extracted, ocr_units_deferred,
                ocr_text_chars, extracted_slides,
            )
        except (RuntimeError, OSError, ValueError, TypeError, KeyError, AttributeError) as render_err:
            logger.warning(
                "[OCR] Slide rendering failed, falling back to shape OCR: %s",
                render_err,
            )
            return False, text_parts, 0, 0, 0, 0, extracted_slides

    @classmethod
    def _ocr_single_shape_image(
        cls, slide_num: int, image_bytes: bytes, policy,
        ocr_preprocess: bool, ocr_postprocess: bool,
        total_slides: int, heartbeat_fn,
        attempted_slides: set[int],
    ) -> tuple[str | None, float, bool]:
        """OCR a single shape image with preprocessing and layout analysis.

        Returns (ocr_text, ocr_conf, timed_out).
        """
        from PIL import Image

        img = Image.open(io.BytesIO(image_bytes))
        img = cls._resize_image_if_needed(img)
        if img is None:
            return None, 0.0, False

        if slide_num not in attempted_slides:
            attempted_slides.add(slide_num)
            cls._emit_status(
                heartbeat_fn,
                f"ocr_processing ppt slide={slide_num}/{total_slides}",
            )

        img_original = img.copy()
        if img_original.mode != "RGB":
            img_original = img_original.convert("RGB")

        img = _preprocess_shape_image(img, ocr_preprocess)

        img_buffer = io.BytesIO()
        img.save(img_buffer, format="PNG")
        png_bytes = img_buffer.getvalue()

        # Layout analysis (PP-Structure) — uses original color image
        ocr_text, ocr_conf = _try_layout_ocr(img_original, policy)

        if not ocr_text:
            with cls._ocr_lock:
                ocr_text, ocr_conf, _ = cls._ocr_extract_safe(
                    png_bytes, f"slide_{slide_num}",
                )

        if ocr_text and ocr_postprocess:
            ocr_text, ocr_conf = _apply_ocr_postprocess(ocr_text, ocr_conf)

        timed_out = ocr_text is None
        return ocr_text, ocr_conf, timed_out

    @classmethod
    def _process_one_shape_ocr(
        cls, slide_num, image_bytes, policy, prs_slides_count, heartbeat_fn,
        ocr_preprocess, ocr_postprocess, attempted_slides, extracted_slides,
    ) -> dict | None:
        """Process OCR for a single shape image. Returns None if OCR unavailable."""
        if slide_num not in attempted_slides and len(attempted_slides) >= policy.ocr_max_ppt_slides:
            return {"deferred": 1}
        ocr = cls._get_ocr_instance()
        if ocr is None:
            return None
        result = {"attempted": 0, "extracted": 0, "deferred": 0, "chars": 0, "text": None, "timed_out_item": None}
        try:
            if slide_num not in attempted_slides:
                result["attempted"] = 1
            ocr_text, ocr_conf, timed_out = cls._ocr_single_shape_image(
                slide_num, image_bytes, policy, ocr_preprocess, ocr_postprocess,
                prs_slides_count, heartbeat_fn, attempted_slides,
            )
            if ocr_text and ocr_conf > 0.3:
                result["text"] = f"[Slide {slide_num} Image OCR]\n{ocr_text}"
                if slide_num not in extracted_slides:
                    extracted_slides.add(slide_num)
                    result["extracted"] = 1
                result["chars"] = cls._text_chars(ocr_text)
            elif timed_out:
                result["timed_out_item"] = (slide_num, image_bytes)
        except (RuntimeError, OSError, ValueError, TypeError, KeyError, AttributeError) as ocr_err:
            logger.warning("[OCR Warning] Slide %d image: %s", slide_num, ocr_err)
        return result

    @staticmethod
    def _accumulate_ocr_result(
        item: dict, totals: dict, text_parts: list[str],
        timed_out_images: list[tuple[int, bytes]] | None = None,
    ) -> None:
        """Accumulate a single OCR item result into running totals."""
        totals["attempted"] += item.get("attempted", 0)
        totals["extracted"] += item.get("extracted", 0)
        totals["deferred"] += item.get("deferred", 0)
        totals["chars"] += item.get("chars", 0)
        if item.get("text"):
            text_parts.append(item["text"])
        if timed_out_images is not None and item.get("timed_out_item"):
            timed_out_images.append(item["timed_out_item"])

    @classmethod
    def _shape_ocr_pass(
        cls, image_shapes, policy, prs_slides_count, heartbeat_fn,
        ocr_preprocess: bool, ocr_postprocess: bool,
        extracted_slides: set[int],
    ):
        """Run shape-by-shape OCR on collected image shapes.

        Returns (text_parts, attempted, extracted, deferred, ocr_chars).
        """
        text_parts: list[str] = []
        timed_out_images: list[tuple[int, bytes]] = []
        totals = {"attempted": 0, "extracted": 0, "deferred": 0, "chars": 0}
        ocr_processed = 0
        ocr_total = len(image_shapes)
        attempted_slides: set[int] = set()

        for slide_num, image_bytes in image_shapes:
            item_result = cls._process_one_shape_ocr(
                slide_num, image_bytes, policy, prs_slides_count, heartbeat_fn,
                ocr_preprocess, ocr_postprocess, attempted_slides, extracted_slides,
            )
            if item_result is None:
                break  # OCR not available
            cls._accumulate_ocr_result(item_result, totals, text_parts, timed_out_images)
            ocr_processed += 1
            if heartbeat_fn and ocr_processed % 10 == 0:
                heartbeat_fn(f"ocr: {ocr_processed}/{ocr_total} images, slide_{slide_num}")

        # Retry timed-out images
        retry_results = cls._retry_timed_out_images(
            timed_out_images, policy, ocr_postprocess,
            attempted_slides, extracted_slides, heartbeat_fn,
        )
        text_parts.extend(retry_results["text_parts"])
        totals["attempted"] += retry_results["attempted"]
        totals["extracted"] += retry_results["extracted"]
        totals["deferred"] += retry_results["deferred"]
        totals["chars"] += retry_results["chars"]

        return (
            text_parts, totals["attempted"], totals["extracted"],
            totals["deferred"], totals["chars"],
        )

    @classmethod
    def _retry_timed_out_images(
        cls, timed_out_images, policy, ocr_postprocess,
        attempted_slides, extracted_slides, heartbeat_fn,
    ) -> dict:
        """Retry OCR on images that timed out during the first pass."""
        result = {
            "text_parts": [],
            "attempted": 0,
            "extracted": 0,
            "deferred": 0,
            "chars": 0,
        }
        if not timed_out_images:
            return result

        logger.info(
            "[OCR Retry] %d timed-out images, retrying sequentially...",
            len(timed_out_images),
        )
        if heartbeat_fn:
            heartbeat_fn(f"ocr_retry: {len(timed_out_images)} images to retry")

        for slide_num, png_bytes in timed_out_images:
            item = cls._retry_one_image(
                slide_num, png_bytes, policy, ocr_postprocess,
                attempted_slides, extracted_slides,
            )
            cls._accumulate_ocr_result(item, result, result["text_parts"])

        return result

    @classmethod
    def _retry_one_image(cls, slide_num, png_bytes, policy, ocr_postprocess, attempted_slides, extracted_slides) -> dict:
        """Retry OCR on a single timed-out image."""
        r = {"attempted": 0, "extracted": 0, "deferred": 0, "chars": 0, "text": None}
        if slide_num not in attempted_slides and len(attempted_slides) >= policy.ocr_max_ppt_slides:
            r["deferred"] = 1
            return r
        try:
            if slide_num not in attempted_slides:
                attempted_slides.add(slide_num)
                r["attempted"] = 1
            with cls._ocr_lock:
                ocr_text, ocr_conf, _ = cls._ocr_extract_safe(png_bytes, f"retry_slide_{slide_num}")
            if ocr_text and ocr_postprocess:
                ocr_text, ocr_conf = _apply_ocr_postprocess(ocr_text, ocr_conf)
            if ocr_text and ocr_conf > 0.3:
                r["text"] = f"[Slide {slide_num} Image OCR]\n{ocr_text}"
                if slide_num not in extracted_slides:
                    extracted_slides.add(slide_num)
                    r["extracted"] = 1
                r["chars"] = cls._text_chars(ocr_text)
                logger.info("[OCR Retry] slide_%d: OK (%d chars)", slide_num, len(ocr_text))
            else:
                logger.info("[OCR Retry] slide_%d: still failed", slide_num)
        except (RuntimeError, OSError, ValueError, TypeError, KeyError, AttributeError) as retry_err:
            logger.warning("[OCR Retry] slide_%d: error - %s", slide_num, retry_err)
        return r

    @classmethod
    def _ppt_pdf_fallback(cls, file_path: Path, heartbeat_fn):
        """Convert PPTX to PDF via LibreOffice, then OCR the PDF.

        Returns (text, tables, ocr_text_chars) or None if fallback fails.
        """
        import subprocess
        import tempfile

        from scripts.slide_renderer import _find_soffice

        soffice = _find_soffice()
        if not soffice:
            return None

        try:
            with tempfile.TemporaryDirectory(prefix="pptx_pdf_fallback_") as tmpdir:
                lo_profile = os.path.join(tmpdir, "lo_profile")
                os.makedirs(lo_profile, exist_ok=True)

                result = subprocess.run(
                    [
                        soffice, "--headless", "--norestore",
                        f"-env:UserInstallation=file://{lo_profile}",
                        "--convert-to", "pdf",
                        "--outdir", tmpdir,
                        str(file_path),
                    ],
                    capture_output=True, text=True, timeout=_w.timeouts.subprocess_ocr_cli,
                )
                if result.returncode != 0:
                    return None

                pdf_files = list(Path(tmpdir).glob("*.pdf"))
                if not pdf_files:
                    return None

                logger.info(
                    "[PPT] PDF fallback: converted %s -> %s",
                    file_path.name, pdf_files[0].name,
                )
                pdf_result = cls.parse_pdf(pdf_files[0], heartbeat_fn=heartbeat_fn)
                if pdf_result.extracted_text.strip():
                    return (
                        pdf_result.extracted_text,
                        pdf_result.extracted_tables,
                        pdf_result.ocr_text_chars,
                    )
        except (RuntimeError, OSError, ValueError, TypeError, KeyError, AttributeError) as fb_err:
            logger.warning("[PPT] PDF fallback failed: %s", fb_err)

        return None

    @classmethod
    def _determine_ppt_skip_reason(
        cls, policy, should_ocr: bool, ocr_units_deferred: int,
    ) -> str | None:
        """Determine the OCR skip reason for the PPT result."""
        if policy.attachment_ocr_mode == "off":
            return "disabled"
        if policy.attachment_ocr_mode == "auto" and not should_ocr:
            return "native_text_sufficient"
        if ocr_units_deferred > 0:
            return "budget_exceeded"
        return None

    @classmethod
    def _apply_pdf_fallback_if_needed(
        cls, should_ocr, full_text, tables, ocr_text_chars, file_path, heartbeat_fn,
    ):
        """Apply PDF fallback if OCR results are too sparse."""
        if not (should_ocr and len(full_text.strip()) < 50):
            return full_text, tables, ocr_text_chars
        logger.info(
            "[PPT] Empty result after all extraction (%d chars), trying PDF fallback",
            len(full_text),
        )
        fb = cls._ppt_pdf_fallback(file_path, heartbeat_fn)
        if fb is not None:
            fb_text, fb_tables, fb_chars = fb
            if len(fb_text.strip()) > len(full_text.strip()):
                full_text = fb_text
                tables = fb_tables or tables
                ocr_text_chars = max(ocr_text_chars, fb_chars)
                logger.info("[PPT] PDF fallback produced %d chars", len(full_text))
        return full_text, tables, ocr_text_chars

    @classmethod
    def _build_ppt_result(
        cls, full_text, tables, policy, should_ocr,
        ocr_units_attempted, ocr_units_extracted,
        ocr_units_deferred, native_text_chars, ocr_text_chars,
    ) -> AttachmentParseResult:
        """Build the final AttachmentParseResult for PPT parsing."""
        return AttachmentParseResult(
            extracted_text=full_text,
            extracted_tables=tables,
            confidence=0.85 if full_text.strip() else 0.0,
            ocr_mode=policy.attachment_ocr_mode,
            ocr_applied=ocr_units_extracted > 0,
            ocr_skip_reason=cls._determine_ppt_skip_reason(
                policy, should_ocr, ocr_units_deferred,
            ),
            ocr_units_attempted=ocr_units_attempted,
            ocr_units_extracted=ocr_units_extracted,
            ocr_units_deferred=ocr_units_deferred,
            native_text_chars=native_text_chars,
            ocr_text_chars=ocr_text_chars,
        )

    @classmethod
    def parse_ppt(cls, file_path: Path, heartbeat_fn=None) -> AttachmentParseResult:
        """PPT에서 슬라이드 텍스트 추출 (.pptx: python-pptx, .ppt: catppt)"""
        try:
            if str(file_path).lower().endswith(".ppt"):
                return cls._parse_legacy_ppt(file_path, heartbeat_fn=heartbeat_fn)
            from pptx import Presentation

            policy = cls.current_policy()
            prs = Presentation(file_path)
            text_parts: list[str] = []
            tables: list[dict] = []
            image_shapes: list[tuple[int, bytes]] = []

            # Stage 1: Extract text, tables, images from slides
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts, slide_tables, slide_images = (
                    cls._extract_ppt_slide_content(slide, slide_num)
                )
                tables.extend(slide_tables)
                image_shapes.extend(slide_images)
                if slide_texts:
                    text_parts.append(
                        f"[Slide {slide_num}]\n" + "\n".join(slide_texts),
                    )

            native_text = "\n\n".join(text_parts)
            native_text_chars = cls._text_chars(native_text)
            cls._emit_status(
                heartbeat_fn,
                f"native_extract ppt slides={len(prs.slides)} chars={native_text_chars}",
            )

            should_ocr = _should_ocr_ppt(policy, native_text_chars)
            ocr_preprocess, ocr_postprocess = _get_ocr_feature_flags()

            ocr_units_attempted = 0
            ocr_units_extracted = 0
            ocr_units_deferred = 0
            ocr_text_chars = 0
            extracted_slides: set[int] = set()

            # Stage 2: Slide rendering OCR
            if should_ocr and policy.slide_render_enabled and file_path:
                (
                    slide_rendered, render_texts,
                    ocr_units_attempted, ocr_units_extracted,
                    ocr_units_deferred, ocr_text_chars, extracted_slides,
                ) = cls._render_and_ocr_slides(
                    file_path, policy, heartbeat_fn,
                    ocr_preprocess, ocr_postprocess,
                )
                text_parts.extend(render_texts)
            else:
                slide_rendered = False

            # Stage 3: Shape-by-shape OCR fallback
            if should_ocr and not slide_rendered:
                (
                    shape_texts, shape_attempted, shape_extracted,
                    shape_deferred, shape_chars,
                ) = cls._shape_ocr_pass(
                    image_shapes, policy, len(prs.slides), heartbeat_fn,
                    ocr_preprocess, ocr_postprocess, extracted_slides,
                )
                text_parts.extend(shape_texts)
                ocr_units_attempted += shape_attempted
                ocr_units_extracted += shape_extracted
                ocr_units_deferred += shape_deferred
                ocr_text_chars += shape_chars

            full_text = "\n\n".join(text_parts)

            # Stage 4: PDF fallback for empty results
            full_text, tables, ocr_text_chars = cls._apply_pdf_fallback_if_needed(
                should_ocr, full_text, tables, ocr_text_chars,
                file_path, heartbeat_fn,
            )

            if ocr_units_deferred > 0:
                cls._emit_status(
                    heartbeat_fn,
                    f"ocr_skipped_budget ppt deferred={ocr_units_deferred}",
                )

            return cls._build_ppt_result(
                full_text, tables, policy, should_ocr,
                ocr_units_attempted, ocr_units_extracted,
                ocr_units_deferred, native_text_chars, ocr_text_chars,
            )

        except (RuntimeError, OSError, ValueError, TypeError, KeyError, AttributeError) as e:
            return AttachmentParseResult(
                extracted_text=f"[PPT 파싱 오류: {e}]",
                extracted_tables=[],
                confidence=0.0,
                ocr_mode=cls.current_policy().attachment_ocr_mode,
                ocr_skip_reason="parse_error",
            )

    # =================================================================
    # OCR singleton management
    # =================================================================

    _ocr_instance = None
    _ocr_type = None  # "paddle" only (amd64 Crawler Pod)
    _ocr_lock = __import__("threading").Lock()  # PaddleOCR는 thread-safe 보장 없음

    @classmethod
    def _get_ocr_instance(cls):
        """싱글톤 PaddleOCR 인스턴스 반환 (amd64 only, no fallback)."""
        if cls._ocr_instance is None:
            try:
                from src.nlp.ocr.paddle_ocr_provider import PaddleOCRProvider
                cls._ocr_instance = PaddleOCRProvider()
                cls._ocr_type = "paddle"
                logger.info("[OCR] PaddleOCR singleton created")
            except ImportError:
                logger.warning("[OCR] PaddleOCR not available (requires amd64)")
                return None
        return cls._ocr_instance

    @classmethod
    def cleanup_ocr(cls):
        """OCR 인스턴스 정리 및 메모리 해제."""
        import gc
        cls._ocr_instance = None
        cls._ocr_type = None
        # Shutdown subprocess pool if active
        if cls._ocr_process_pool is not None:
            try:
                cls._ocr_process_pool.shutdown(wait=False)
            except (RuntimeError, OSError, ValueError, TypeError, KeyError, AttributeError):
                pass
            cls._ocr_process_pool = None
        gc.collect()
        logger.info("[OCR] 메모리 정리 완료")

    # --- Subprocess-isolated OCR (SIGSEGV defense) ---
    _ocr_process_pool = None
    _ocr_pool_lock = __import__("threading").Lock()

    @classmethod
    def _ocr_extract_safe(
        cls, image_bytes: bytes, file_name: str = "", timeout: int = 1800
    ) -> tuple[str | None, float, list]:
        """Execute OCR in a forked subprocess to survive PaddleOCR SIGSEGV."""
        from concurrent.futures import ProcessPoolExecutor
        from concurrent.futures.process import BrokenProcessPool

        import multiprocessing as mp

        with cls._ocr_pool_lock:
            if cls._ocr_process_pool is None:
                ctx = mp.get_context("fork")
                cls._ocr_process_pool = ProcessPoolExecutor(
                    max_workers=1, mp_context=ctx
                )
            pool = cls._ocr_process_pool

        try:
            future = pool.submit(_ocr_worker_fn, image_bytes)
            return future.result(timeout=timeout)
        except BrokenProcessPool:
            logger.error(
                "[OCR SIGSEGV] Worker crashed on %s — restarting pool, skipping image",
                file_name,
            )
            with cls._ocr_pool_lock:
                cls._ocr_process_pool = None
            return None, 0.0, []
        except TimeoutError:
            logger.warning("[OCR Timeout] %s exceeded %ds — skipped", file_name, timeout)
            return None, 0.0, []
        except (RuntimeError, OSError, ValueError, TypeError, KeyError, AttributeError) as e:
            logger.warning("[OCR Error] %s: %s — skipped", file_name, e)
            return None, 0.0, []

    # PaddleOCR PP-OCRv5 det 모델이 극단적 종횡비에서 SIGSEGV 발생
    _OCR_MIN_DIMENSION = OCR_MIN_DIMENSION  # from _attachment_helpers
    _OCR_MAX_ASPECT_RATIO = OCR_MAX_ASPECT_RATIO  # from _attachment_helpers

    @staticmethod
    def _resize_image_if_needed(img, max_size: int = 2048):
        """큰 이미지 리사이즈 (메모리 최적화).

        Returns:
            리사이즈된 이미지 (또는 원본).
            극단적 종횡비(>8:1)나 너무 작은(<32px) 이미지는 None 반환.
        """
        width, height = img.size

        if width < AttachmentParser._OCR_MIN_DIMENSION or height < AttachmentParser._OCR_MIN_DIMENSION:
            logger.debug("[OCR] 이미지 스킵 (너무 작음): %dx%d", width, height)
            return None

        if width > max_size or height > max_size:
            img = _downscale_image(img, width, height, max_size)
            if img is None:
                return None

        return _pad_extreme_aspect_ratio(img)

    # =================================================================
    # OCR slide image — split into stages
    # =================================================================

    @classmethod
    def _ocr_slide_image(
        cls,
        png_bytes: bytes,
        slide_num: int,
        preprocess: bool = True,
        layout_analysis: bool = True,
        postprocess: bool = True,
    ) -> str | None:
        """OCR a rendered slide image with preprocessing and layout analysis."""
        from PIL import Image

        try:
            img = Image.open(io.BytesIO(png_bytes))
            if img.mode != "RGB":
                img = img.convert("RGB")

            img_original = img.copy()

            if preprocess:
                img = _preprocess_slide_image(img, slide_num)

            ocr_text = _try_slide_layout_ocr(img_original, layout_analysis, slide_num)

            if not ocr_text:
                ocr_text = cls._fallback_standard_ocr(img, slide_num)
                if not ocr_text:
                    return None

            if ocr_text and postprocess:
                ocr_text = _postprocess_slide_text(ocr_text, slide_num)

            if ocr_text:
                ocr_text = _filter_ocr_noise(ocr_text)

            return ocr_text if ocr_text and ocr_text.strip() else None

        except (RuntimeError, OSError, ValueError, TypeError, KeyError, AttributeError) as e:
            logger.error("[OCR] Slide %d OCR error: %s", slide_num, e)
            return None

    @classmethod
    def _fallback_standard_ocr(cls, img, slide_num: int) -> str | None:
        """Run standard OCR on a preprocessed slide image."""
        img_buffer = io.BytesIO()
        img.save(img_buffer, format="PNG")
        with cls._ocr_lock:
            ocr_text, ocr_conf, _ = cls._ocr_extract_safe(
                img_buffer.getvalue(), f"rendered_slide_{slide_num}"
            )
        if not ocr_text or ocr_conf <= 0.3:
            return None
        return ocr_text

    # =================================================================
    # Image parsing — split into helpers
    # =================================================================

    @classmethod
    def _parse_image_sync(
        cls,
        file_path: Path,
        content: bytes,
        use_ocr: bool = True,
    ) -> AttachmentParseResult:
        """이미지 OCR 및 메타데이터 추출 (동기 내부 구현)."""
        try:
            policy = cls.current_policy()
            from PIL import Image

            img = Image.open(io.BytesIO(content))
            width, height = img.size
            format_type = img.format or "unknown"
            metadata_text = f"[Image: {width}x{height}, {format_type}, {len(content):,} bytes]"

            if policy.attachment_ocr_mode == "off" or not use_ocr:
                return _image_result_no_ocr(metadata_text, policy, "disabled")

            if policy.ocr_max_images_per_attachment <= 0:
                cls._emit_status(None, "ocr_skipped_budget image deferred=1")
                return _image_result_no_ocr(
                    metadata_text, policy, "budget_exceeded", ocr_units_deferred=1,
                )

            if len(content) >= 10_000_000:
                return _image_result_no_ocr(
                    metadata_text, policy, "image_too_large",
                )

            return cls._perform_image_ocr(img, content, file_path, metadata_text, policy)

        except (RuntimeError, OSError, ValueError, TypeError, KeyError, AttributeError) as e:
            return AttachmentParseResult(
                extracted_text=f"[이미지 파싱 오류: {e}]",
                extracted_tables=[],
                confidence=0.0,
                ocr_mode=cls.current_policy().attachment_ocr_mode,
                ocr_skip_reason="parse_error",
            )

    @classmethod
    def _perform_image_ocr(cls, img, content, file_path, metadata_text, policy):
        """Execute OCR on an image and return the result."""
        import gc

        try:
            cls._emit_status(None, f"ocr_processing image file={file_path.name}")
            ocr = cls._get_ocr_instance()
            if ocr is None:
                return _image_result_no_ocr(metadata_text, policy, "ocr_unavailable")

            img = cls._resize_image_if_needed(img)
            if img is None:
                return _image_result_no_ocr(metadata_text, policy, "guard_rejected")

            if img.mode != "RGB":
                img = img.convert("RGB")

            img_buffer = io.BytesIO()
            img.save(img_buffer, format="PNG")
            resized_content = img_buffer.getvalue()

            with cls._ocr_lock:
                ocr_text, ocr_conf, ocr_tables = cls._ocr_extract_safe(
                    resized_content, file_path.name,
                )

            if ocr_text and ocr_conf > 0.3:
                full_text = f"{metadata_text}\n\n{ocr_text}"
                del img, img_buffer, resized_content
                gc.collect()

                return AttachmentParseResult(
                    extracted_text=full_text,
                    extracted_tables=ocr_tables,
                    confidence=ocr_conf,
                    ocr_mode=policy.attachment_ocr_mode,
                    ocr_applied=True,
                    ocr_units_attempted=1,
                    ocr_units_extracted=1,
                    ocr_text_chars=cls._text_chars(ocr_text),
                )

        except (RuntimeError, OSError, ValueError, TypeError, KeyError, AttributeError) as ocr_error:
            logger.warning("[OCR Warning] %s: %s", file_path.name, ocr_error)

        return AttachmentParseResult(
            extracted_text=metadata_text,
            extracted_tables=[],
            confidence=0.5,
            ocr_mode=policy.attachment_ocr_mode,
            ocr_skip_reason="ocr_failed",
            ocr_units_attempted=1,
        )

    @classmethod
    def parse_image(
        cls,
        file_path: Path,
        content: bytes,
        use_ocr: bool = True,
    ) -> AttachmentParseResult:
        """이미지 OCR 및 메타데이터 추출 (동기 호출용 래퍼)."""
        return cls._parse_image_sync(file_path, content, use_ocr)

    @classmethod
    async def parse_image_async(
        cls,
        file_path: Path,
        content: bytes,
        use_ocr: bool = True,
    ) -> AttachmentParseResult:
        """이미지 OCR 비동기 처리."""
        return await asyncio.to_thread(cls._parse_image_sync, file_path, content, use_ocr)

